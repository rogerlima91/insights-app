import io
import json
import os
import re
from datetime import date
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import plotly.graph_objects as go
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(page_title="Insights App", layout="wide")

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY    = "#14113b"
CYAN    = "#00b2a9"
GREEN   = "#34b233"
PURPLE  = "#6e2ca9"
MAGENTA = "#c70099"

# ── Global CSS ────────────────────────────────────────────────────────────────
# STYLE LOCK: Do not remove or modify this CSS block.
st.markdown("""
<style>
    /* ── Base font and body ─────────────────────────────────────── */
    html, body, [class*="css"] {
        font-family: "Inter", system-ui, -apple-system, "Segoe UI", sans-serif;
        font-size: 15px;
        color: #374151;
    }

    /* ── Page background ────────────────────────────────────────── */
    .stApp {
        background-color: #F3F4F6;
    }

    /* ── Sidebar — purple SaaS style ────────────────────────────── */
    [data-testid="stSidebar"] {
        background-color: #7C3AED !important;
        border-right: none !important;
    }
    section[data-testid="stSidebar"] > div {
        background-color: #7C3AED !important;
    }
    section[data-testid="stSidebar"] > div > div {
        background-color: #7C3AED !important;
    }
    section[data-testid="stSidebar"] * {
        background-color: transparent !important;
    }
    section[data-testid="stSidebar"] > div,
    section[data-testid="stSidebar"] > div > div,
    section[data-testid="stSidebar"] > div > div > div {
        background-color: #7C3AED !important;
    }
    /* All text inside sidebar: white */
    [data-testid="stSidebar"] *,
    [data-testid="stSidebar"] p,
    [data-testid="stSidebar"] span,
    [data-testid="stSidebar"] label,
    [data-testid="stSidebar"] small {
        color: #FFFFFF !important;
    }
    /* Sidebar headings — white, no left border */
    [data-testid="stSidebar"] h1,
    [data-testid="stSidebar"] h2,
    [data-testid="stSidebar"] h3,
    [data-testid="stSidebar"] h4 {
        color: #FFFFFF !important;
        border-left: none !important;
        padding-left: 0 !important;
        margin-top: 1rem !important;
    }
    /* Force inner sidebar containers to match purple — prevents white gap below nav */
    [data-testid="stSidebarContent"],
    [data-testid="stSidebarUserContent"] {
        background-color: #7C3AED !important;
    }
    /* Sidebar nav container — transparent with subtle border, no white bar */
    [data-testid="stSidebarNav"] {
        background: transparent !important;
        border: 1px solid rgba(255,255,255,0.3) !important;
        border-radius: 8px !important;
    }
    /* Sidebar nav links */
    [data-testid="stSidebarNav"] a span {
        color: rgba(255,255,255,0.80) !important;
        font-weight: 500;
        text-transform: capitalize;
    }
    [data-testid="stSidebarNav"] a:hover span {
        color: #FFFFFF !important;
    }
    /* Active page — white pill with purple text */
    [data-testid="stSidebarNavLink"][aria-selected="true"],
    [data-testid="stSidebarNav"] a[aria-current="page"] {
        background: #FFFFFF !important;
        border-radius: 20px !important;
    }
    [data-testid="stSidebarNavLink"][aria-selected="true"] span,
    [data-testid="stSidebarNav"] a[aria-current="page"] span {
        color: #7C3AED !important;
        font-weight: 700 !important;
    }
    /* Sidebar selectbox input */
    [data-testid="stSidebar"] [data-testid="stSelectbox"] > div > div {
        background: rgba(255,255,255,0.15) !important;
        border-color: rgba(255,255,255,0.35) !important;
    }
    /* Sidebar divider */
    [data-testid="stSidebar"] hr {
        border-color: rgba(255,255,255,0.25) !important;
    }

    /* ── Main area headings ──────────────────────────────────────── */
    h1, h2, h3, h4, h5, h6 {
        font-weight: 700 !important;
        color: #111827 !important;
    }
    h2, h3 {
        margin-top: 2rem !important;
        padding-top: 0.25rem !important;
        border-bottom: none !important;
        padding-bottom: 0 !important;
        border-left: none !important;
        padding-left: 0 !important;
    }

    /* ── KPI metric cards ───────────────────────────────────────── */
    [data-testid="metric-container"] {
        background: #FFFFFF;
        border: none;
        border-radius: 16px;
        padding: 20px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        border-top: 4px solid #7C3AED;
    }
    [data-testid="metric-container"] label {
        font-size: 12px;
        font-weight: 600;
        color: #6B7280;
        text-transform: uppercase;
        letter-spacing: 0.05em;
    }
    [data-testid="metric-container"] [data-testid="stMetricValue"] {
        font-size: 26px;
        font-weight: 700;
        color: #111827;
    }

    /* ── Chart cards — wrap Plotly chart output in white card ───── */
    .element-container:has([data-testid="stPlotlyChart"]) {
        background: #FFFFFF;
        border-radius: 12px;
        padding: 16px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
    }

    /* ── Primary buttons — purple ───────────────────────────────── */
    .stButton > button[kind="primary"],
    [data-testid="baseButton-primary"] {
        background-color: #7C3AED !important;
        color: #FFFFFF !important;
        border: none !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        font-size: 15px !important;
        padding: 0.5rem 1.25rem !important;
    }
    .stButton > button[kind="primary"]:hover,
    [data-testid="baseButton-primary"]:hover {
        background-color: #6D28D9 !important;
    }

    /* ── File upload box ─────────────────────────────────────────── */
    [data-testid="stFileUploader"] {
        border: 2px dashed #7C3AED !important;
        border-radius: 12px;
        padding: 10px;
        background: #FFFFFF;
    }

    /* ── DSP source badges ───────────────────────────────────────── */
    .source-badge {
        display: inline-block;
        padding: 2px 10px;
        border-radius: 20px;
        font-size: 11px;
        font-weight: 700;
        margin-left: 6px;
    }
    .badge-dv360   { background: #EFF6FF; color: #2563EB; }
    .badge-ttd     { background: #F5F3FF; color: #7C3AED; }
    .badge-generic { background: #F0FDF4; color: #16A34A; }
</style>
""", unsafe_allow_html=True)
# STYLE LOCK

# ── Column name normalisation map ─────────────────────────────────────────────
# Maps the many different column names DSPs use → our standard internal names.
# Add more aliases here as you encounter new DSP export formats.
COLUMN_MAP = {
    # Campaign / Brand (high-level grouping)
    "campaign":                    "campaign",
    "campaign name":               "campaign",
    "campaign_name":               "campaign",
    "insertion order":             "campaign",
    "advertiser":                  "campaign",
    "brand name":                  "campaign",
    "brand":                       "campaign",

    # Line item (more granular — used for the top-10 line items chart)
    "line item":                   "line_item",
    "line item name":              "line_item",
    "ad group":                    "line_item",
    "ad group name":               "line_item",
    "creative":                    "line_item",

    # Device type (for pie chart)
    "device type":                 "device_type",
    "device":                      "device_type",
    "device_type":                 "device_type",
    "device category":             "device_type",

    # Environment / inventory type (for pie chart)
    "environment":                 "environment",
    "environment type":            "environment",
    "inventory type":              "environment",
    "supply type":                 "environment",
    "site type":                   "environment",

    # Date
    "date":                        "date",
    "day":                         "date",
    "week":                        "date",
    "month":                       "date",

    # Impressions
    "impressions":                 "impressions",
    "impression":                  "impressions",
    "served impressions":          "impressions",
    "total impressions":           "impressions",

    # Clicks
    "clicks":                      "clicks",
    "click":                       "clicks",
    "total clicks":                "clicks",
    "link clicks":                 "clicks",

    # Spend / cost
    "spend":                       "spend_usd",
    "spend (usd)":                 "spend_usd",
    "total spend":                 "spend_usd",
    "total spend (usd)":           "spend_usd",
    "media cost":                  "spend_usd",
    "media cost (usd)":            "spend_usd",
    "revenue (usd)":               "spend_usd",
    "cost":                        "spend_usd",
    "billed spend":                "spend_usd",

    # Conversions
    "conversions":                 "conversions",
    "total conversions":           "conversions",
    "post-click conversions":      "conversions",

    # CTR (sometimes pre-calculated in export)
    "ctr":                         "ctr_raw",
    "click-through rate":          "ctr_raw",

    # CPM (sometimes pre-calculated)
    "cpm":                         "cpm_raw",
    "avg. cpm":                    "cpm_raw",
    "average cpm":                 "cpm_raw",
}

# ── DSP source detector ───────────────────────────────────────────────────────
# Tries to identify which DSP the CSV came from based on column names.
def detect_source(columns):
    cols_lower = [c.lower() for c in columns]
    if any("revenue (usd)" in c or "insertion order" in c for c in cols_lower):
        return "DV360"
    if any("billed spend" in c or "partner" in c for c in cols_lower):
        return "TTD"
    return "Generic"

# ── Brand memory helpers ──────────────────────────────────────────────────────
# brand_memory.json stores per-brand context that shapes AI commentary.
BRAND_MEMORY_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                 "brand_memory.json")

def load_brand_memory():
    """Load brand memory from JSON file. Returns empty dict if file doesn't exist."""
    if os.path.exists(BRAND_MEMORY_PATH):
        with open(BRAND_MEMORY_PATH, "r") as f:
            return json.load(f)
    return {}

def save_brand_memory(memory):
    """Save brand memory dict back to JSON file."""
    with open(BRAND_MEMORY_PATH, "w") as f:
        json.dump(memory, f, indent=2)

# ── CSV loader and normaliser ─────────────────────────────────────────────────
def load_and_normalise(uploaded_file):
    """
    Reads a CSV file, normalises column names to our standard format,
    and adds a 'source_file' and 'dsp_source' column for traceability.
    """
    df = pd.read_csv(uploaded_file)

    # Detect DSP before renaming columns
    dsp = detect_source(df.columns.tolist())

    # Normalise column names: lowercase + strip whitespace, then remap
    df.columns = [c.strip().lower() for c in df.columns]
    df = df.rename(columns={k: v for k, v in COLUMN_MAP.items() if k in df.columns})

    # Multiple source columns can map to the same standard name (e.g. both
    # "CTR" and "Click-through rate" → "ctr_raw"). Keep the first occurrence.
    df = df.loc[:, ~df.columns.duplicated(keep="first")]

    # Convert numeric columns — DSP exports often include commas or $ signs
    for col in ["impressions", "clicks", "spend_usd", "conversions"]:
        if col in df.columns:
            df[col] = (
                df[col].astype(str)
                .str.replace(r"[$,]", "", regex=True)   # strip $ and commas
                .str.strip()
            )
            df[col] = pd.to_numeric(df[col], errors="coerce")

    # Calculate CTR and CPM from raw numbers (more reliable than DSP-provided values)
    if "impressions" in df.columns and "clicks" in df.columns:
        df["ctr"] = df["clicks"] / df["impressions"]
    if "impressions" in df.columns and "spend_usd" in df.columns:
        df["cpm"] = df["spend_usd"] / df["impressions"] * 1000

    # Tag each row with where it came from
    df["source_file"] = uploaded_file.name
    df["dsp_source"]  = dsp

    return df, dsp

# ── PowerPoint export — dark premium template ────────────────────────────────
# Built entirely from a blank Presentation (no template file required).
#
# Colour scheme:
#   Background  #0D1B2A  dark navy
#   Primary     #FFFFFF  white text
#   Secondary   #A8B2BC  light grey labels / footer
#   Accent      #00A8E8  electric blue
#   Box fill    #1A2D40  slightly lighter navy for cards and table rows
#
# Font: Calibri throughout.  Slide size: 13.33" × 7.5" widescreen.
#
# Slide structure:
#   1.      Executive Summary — 3 KPI cards + Revenue chart + AI best/worst insight
#   2 + 3.  Per brand (repeated): Performance Breakdown + Recommendations
#   Last.   Budget Shift Recommendations (table)

_PPT_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy background
_PPT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)    # primary text
_PPT_GREY  = RGBColor(0xA8, 0xB2, 0xBC)   # secondary / label text
_PPT_BLUE  = RGBColor(0x00, 0xA8, 0xE8)   # electric blue accent
_PPT_DIM   = RGBColor(0x1A, 0x2D, 0x40)   # lighter navy for boxes / table rows

_PW = 13.33   # slide width in inches
_PH = 7.5     # slide height in inches

def _new_blank_slide(prs):
    """Add a blank slide and fill it with the dark navy background colour."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill  = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _PPT_BG
    return slide


def _tb(slide, x, y, w, h, text, size, bold=False, color=None,
        align=PP_ALIGN.LEFT):
    """Add a single-run Calibri text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.name      = "Calibri"
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color or _PPT_WHITE
    return txb

def _box(slide, x, y, w, h, fill_color, border_color=None, border_pt=1.0):
    """Add a filled rectangle with an optional coloured border."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width     = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def _kpi_box(slide, x, y, w, h, label, value):
    """KPI card: dark box with electric blue border, small label + large bold value."""
    _box(slide, x, y, w, h, _PPT_DIM, border_color=_PPT_BLUE, border_pt=1.5)
    _tb(slide, x + 0.1, y + 0.1,  w - 0.2, 0.28,
        label.upper(), 8, color=_PPT_GREY, align=PP_ALIGN.CENTER)
    _tb(slide, x + 0.05, y + 0.4, w - 0.1, 0.55,
        value, 22, bold=True, align=PP_ALIGN.CENTER)


def _slide_footer(slide, slide_num):
    """'Insights App' logo text bottom-left and slide number bottom-right."""
    _tb(slide, 0.3,      _PH - 0.34, 2.5, 0.28, "Insights App", 8, color=_PPT_GREY)
    _tb(slide, _PW - 1.6, _PH - 0.34, 1.3, 0.28, str(slide_num),
        8, color=_PPT_GREY, align=PP_ALIGN.RIGHT)


def _add_bullets(slide, x, y, w, h, bullets, size=11):
    """
    Render a list of bullet points into a text box.
    Each bullet: electric blue dot (●) followed by white Calibri text.
    """
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True

    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)

        r_dot            = p.add_run()
        r_dot.text       = "●  "
        r_dot.font.name  = "Calibri"
        r_dot.font.size  = Pt(size - 1)
        r_dot.font.color.rgb = _PPT_BLUE

        r_text            = p.add_run()
        r_text.text       = text
        r_text.font.name  = "Calibri"
        r_text.font.size  = Pt(size)
        r_text.font.color.rgb = _PPT_WHITE


def _make_revenue_chart_pptx(camp_summary):
    """
    Build a Revenue by Brand bar chart with a dark background and electric blue bars.
    Returns a BytesIO PNG buffer, or None if spend data is unavailable.
    """
    if camp_summary.empty or "spend_usd" not in camp_summary.columns:
        return None

    data = (camp_summary[["campaign", "spend_usd"]]
            .dropna()
            .sort_values("spend_usd", ascending=False))

    fig, ax = plt.subplots(figsize=(7.5, 3.6))
    fig.patch.set_facecolor("#0D1B2A")
    ax.set_facecolor("#0D1B2A")

    bars = ax.bar(data["campaign"], data["spend_usd"],
                  color="#00A8E8", edgecolor="none", width=0.5)
    ax.yaxis.set_major_formatter(
        mticker.FuncFormatter(
            lambda v, _: f"${v/1e6:.1f}M" if v >= 1e6 else f"${v/1e3:.0f}K"
        )
    )
    ax.bar_label(bars,
                 labels=[f"${v/1e6:.2f}M" if v >= 1e6 else f"${v/1e3:.0f}K"
                         for v in data["spend_usd"]],
                 padding=4, fontsize=9, color="white", fontweight="bold")
    ax.set_ylim(0, data["spend_usd"].max() * 1.32)
    ax.tick_params(colors="white", labelsize=9)
    ax.set_ylabel("Revenue (USD)", fontsize=9, color="#A8B2BC")
    for sp in ax.spines.values():
        sp.set_color("#2A3D4E")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", linewidth=0.4, alpha=0.4, color="#2A3D4E")
    ax.set_axisbelow(True)
    plt.xticks(rotation=20, ha="right", fontsize=9, color="white")
    plt.tight_layout(pad=0.4)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="#0D1B2A")
    buf.seek(0)
    plt.close(fig)
    return buf


def _build_data_text(df_all, camp_summary):
    """
    Produce a compact text summary of all campaign data for use in AI prompts.
    Includes brand totals and a line item breakdown when available.
    """
    lines = ["=== Brand Totals ==="]
    for _, r in camp_summary.iterrows():
        parts = [f"Brand: {r['campaign']}"]
        if "impressions" in r and pd.notna(r.get("impressions")):
            parts.append(f"Impressions: {int(r['impressions']):,}")
        if "clicks" in r and pd.notna(r.get("clicks")):
            parts.append(f"Clicks: {int(r['clicks']):,}")
        if "spend_usd" in r and pd.notna(r.get("spend_usd")):
            parts.append(f"Revenue: ${r['spend_usd']:,.0f}")
        if "ctr" in r and pd.notna(r.get("ctr")):
            parts.append(f"CTR: {r['ctr']:.2%}")
        if "cpm" in r and pd.notna(r.get("cpm")):
            parts.append(f"CPM: ${r['cpm']:,.2f}")
        lines.append(" | ".join(parts))

    # Line item breakdown gives the AI enough detail to cite best/worst performers
    if "line_item" in df_all.columns and "campaign" in df_all.columns:
        lines.append("\n=== By Brand and Line Item ===")
        agg = {c: (c, "sum") for c in ["impressions", "clicks", "spend_usd"]
               if c in df_all.columns}
        grp = df_all.groupby(["campaign", "line_item"]).agg(**agg).reset_index()
        if "impressions" in grp.columns and "spend_usd" in grp.columns:
            grp["cpm"] = grp["spend_usd"] / grp["impressions"].clip(lower=1) * 1000
        if "impressions" in grp.columns and "clicks" in grp.columns:
            grp["ctr"] = grp["clicks"] / grp["impressions"].clip(lower=1)
        if "clicks" in grp.columns and "spend_usd" in grp.columns:
            grp["cpc"] = grp["spend_usd"] / grp["clicks"].clip(lower=1)
        for _, r in grp.iterrows():
            parts = [f"Brand: {r['campaign']} | Line Item: {r['line_item']}"]
            if "impressions" in r and pd.notna(r.get("impressions")):
                parts.append(f"Impressions: {int(r['impressions']):,}")
            if "spend_usd" in r and pd.notna(r.get("spend_usd")):
                parts.append(f"Revenue: ${r['spend_usd']:,.0f}")
            if "cpm" in r and pd.notna(r.get("cpm")):
                parts.append(f"CPM: ${r['cpm']:,.2f}")
            if "ctr" in r and pd.notna(r.get("ctr")):
                parts.append(f"CTR: {r['ctr']:.2%}")
            if "cpc" in r and pd.notna(r.get("cpc")):
                parts.append(f"CPC: ${r['cpc']:,.2f}")
            lines.append(" | ".join(parts))

    return "\n".join(lines)


def _ai_pptx(api_key, prompt, max_tokens=400):
    """Synchronous Anthropic API call for slide content. Returns the response text."""
    client = anthropic.Anthropic(api_key=api_key)
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=max_tokens,
        system=(
            "You are a senior programmatic advertising analyst. "
            "Write clear, concise, data-driven commentary for ad tech professionals. "
            "Be specific — always reference actual numbers from the data. "
            "Follow formatting instructions exactly — no deviations."
        ),
        messages=[{"role": "user", "content": prompt}],
    )
    return msg.content[0].text.strip()


def _parse_io_sections(text):
    """
    Parse AI breakdown text into [(section_title, [bullet_strings])] tuples.
    Looks for DISPLAY:, VIDEO:, YOUTUBE: section headers.
    """
    sections      = []
    current_title = None
    current_items = []

    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        upper = s.upper()
        if upper.startswith("DISPLAY"):
            if current_title:
                sections.append((current_title, current_items))
            current_title, current_items = "Display", []
        elif upper.startswith("YOUTUBE"):
            if current_title:
                sections.append((current_title, current_items))
            current_title, current_items = "YouTube", []
        elif upper.startswith("VIDEO"):
            if current_title:
                sections.append((current_title, current_items))
            current_title, current_items = "Video", []
        elif s and (s[0] in "•-–—●" or (len(s) > 1 and s[0].isdigit() and s[1] in ".)")):
            bullet = s.lstrip("•-–—●0123456789.) ").strip()
            if bullet:
                current_items.append(bullet)

    if current_title:
        sections.append((current_title, current_items))

    return sections


def _add_budget_table(slide, budget_text, x, y, w, h):
    """
    Parse the AI budget text and render a styled table onto the slide.
    Expected format: pipe-separated — Brand | Best IO | Recommendation
    """
    rows = []
    for line in budget_text.split("\n"):
        parts = [p.strip() for p in line.split("|")]
        if len(parts) >= 3 and parts[0]:
            rows.append(parts[:3])

    if not rows:
        _tb(slide, x, y + 0.2, w, 0.4, "No budget data available.", 11, color=_PPT_GREY)
        return

    headers   = ["Brand", "Best IO", "Recommendation"]
    gap       = 0.06
    col_ws    = [w * 0.22 - gap, w * 0.18 - gap, w * 0.57 - gap]
    row_h     = 0.5

    # Header row
    cx = x
    for i, hdr in enumerate(headers):
        _box(slide, cx, y, col_ws[i], row_h * 0.75, _PPT_BLUE)
        _tb(slide, cx + 0.08, y + 0.1, col_ws[i] - 0.16, 0.34,
            hdr, 10, bold=True)
        cx += col_ws[i] + gap

    # Data rows
    for r_i, row in enumerate(rows):
        ry = y + row_h * 0.8 + r_i * row_h
        if ry + row_h > y + h:
            break
        row_fill = _PPT_DIM if r_i % 2 == 0 else _PPT_BG
        cx = x
        for c_i, cell in enumerate(row):
            _box(slide, cx, ry, col_ws[c_i], row_h - 0.05, row_fill)
            # Highlight positive recommendations in electric blue
            txt_color = _PPT_BLUE if any(
                kw in cell.lower() for kw in ["increase", "scale", "boost", "shift to"]
            ) else _PPT_WHITE
            _tb(slide, cx + 0.08, ry + 0.08, col_ws[c_i] - 0.16, row_h - 0.12,
                cell, 9, color=txt_color)
            cx += col_ws[c_i] + gap


def build_pptx_report(api_key, camp_summary, df_all):
    """
    Build and return the full PowerPoint report as a BytesIO buffer.

    Slides:
      1.    Executive Summary — 3 KPI cards, Revenue chart, AI best/worst insight
      2+3.  Per brand (repeated): Performance Breakdown + Recommendations
      Last. Budget Shift Recommendations
    """
    prs              = Presentation()
    prs.slide_width  = Inches(_PW)
    prs.slide_height = Inches(_PH)

    slide_num = 1
    data_text = _build_data_text(df_all, camp_summary)
    campaigns = (camp_summary["campaign"].tolist()
                 if "campaign" in camp_summary.columns else [])

    # ── Slide 1: Executive Summary ────────────────────────────────────────────
    s1 = _new_blank_slide(prs)

    # Title + electric blue accent line
    _tb(s1, 0.4, 0.17, 12.5, 0.75, "Campaign Performance Summary", 28, bold=True)
    _box(s1, 0.4, 0.95, 12.5, 0.04, _PPT_BLUE)

    # Three KPI cards
    total_impr = (camp_summary["impressions"].sum()
                  if "impressions" in camp_summary.columns else 0)
    total_rev  = (camp_summary["spend_usd"].sum()
                  if "spend_usd" in camp_summary.columns else 0)
    kpi_data = [
        ("Total Impressions", f"{total_impr:,.0f}"),
        ("Total Revenue",     f"${total_rev:,.0f}"),
        ("Number of Brands",  str(len(campaigns))),
    ]
    kpi_w, kpi_gap = 3.8, 0.35
    for i, (lbl, val) in enumerate(kpi_data):
        _kpi_box(s1, 0.4 + i * (kpi_w + kpi_gap), 1.08, kpi_w, 1.05, lbl, val)

    # Revenue by Brand chart (left panel)
    chart_buf = _make_revenue_chart_pptx(camp_summary)
    if chart_buf:
        s1.shapes.add_picture(chart_buf,
                              Inches(0.4), Inches(2.32), Inches(7.8), Inches(3.95))

    # AI best/worst insight (right panel)
    exec_text = _ai_pptx(
        api_key,
        f"""Campaign performance data:
{data_text}

Write two short insight paragraphs for an executive summary slide:
1. Start "Best performer: [Brand] —" then 1-2 sentences on WHY, citing key metrics (CPM, CTR, or CPV).
2. Start "Worst performer: [Brand] —" then 1-2 sentences on WHY, citing key metrics.

Plain paragraphs only — no headings, no bullets, no extra text.""",
        max_tokens=200,
    )
    _tb(s1, 8.5, 2.32, 4.5, 0.28, "KEY INSIGHTS", 8, bold=True, color=_PPT_BLUE)
    txb = s1.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(4.45), Inches(3.6))
    tf  = txb.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text           = exec_text
    run.font.name      = "Calibri"
    run.font.size      = Pt(10.5)
    run.font.color.rgb = _PPT_GREY

    _slide_footer(s1, slide_num)
    slide_num += 1

    # ── Slides 2+3 per brand ──────────────────────────────────────────────────
    for brand in campaigns:

        # ── Performance Breakdown ─────────────────────────────────────────────
        s2 = _new_blank_slide(prs)
        _tb(s2, 0.4, 0.17, 12.5, 0.65,
            f"{brand}  —  Performance Breakdown", 20, bold=True)
        _box(s2, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)

        breakdown = _ai_pptx(
            api_key,
            f"""Campaign performance data:
{data_text}

Write a performance breakdown for brand: "{brand}"

Use EXACTLY this format — only include sections where data exists for this brand:

DISPLAY:
• [CPM, CPC, CTR focus — cite best/worst line item, max 15 words]
• [second point]
• [third point — max 3 bullets]

VIDEO:
• [CPV, VTR focus — cite best/worst line item, max 15 words]
• [second point]
• [third point — max 3 bullets]

YOUTUBE:
• [CPV, VTR focus — cite best/worst creative, max 15 words]
• [second point]
• [third point — max 3 bullets]

Only include DISPLAY / VIDEO / YOUTUBE headings that have real data for this brand.""",
            max_tokens=350,
        )

        sections = _parse_io_sections(breakdown)
        if sections:
            n     = len(sections)
            col_w = (12.5 - (n - 1) * 0.3) / n
            cx    = 0.4
            for sec_title, sec_bullets in sections:
                _box(s2, cx, 1.05, col_w, 0.38, _PPT_DIM)
                _tb(s2, cx + 0.12, 1.1, col_w - 0.24, 0.3,
                    sec_title.upper(), 11, bold=True, color=_PPT_BLUE)
                if sec_bullets:
                    _add_bullets(s2, cx + 0.08, 1.58,
                                 col_w - 0.16, 5.5, sec_bullets[:3], size=10)
                cx += col_w + 0.3

        _slide_footer(s2, slide_num)
        slide_num += 1

        # ── Recommendations ───────────────────────────────────────────────────
        s3 = _new_blank_slide(prs)
        _tb(s3, 0.4, 0.17, 12.5, 0.65,
            f"{brand}  —  Recommendations", 20, bold=True)
        _box(s3, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)

        recs_raw = _ai_pptx(
            api_key,
            f"""Campaign performance data:
{data_text}

Write exactly 5 optimisation recommendations for brand: "{brand}"

Rules:
- Each must start with an action verb: Increase / Reduce / Pause / Test / Shift / Reallocate
- Reference specific line items, creatives, or metrics from the data
- Max 18 words per line
- Return only the 5 lines — no bullet symbols, no numbering, no extra text""",
            max_tokens=250,
        )
        bullets = [ln.strip().lstrip("•-–—●0123456789.) ").strip()
                   for ln in recs_raw.split("\n")
                   if ln.strip() and not ln.strip().startswith("#")]
        bullets = [b for b in bullets if b][:5]
        _add_bullets(s3, 0.6, 1.12, 12.1, 5.85, bullets, size=13)

        _slide_footer(s3, slide_num)
        slide_num += 1

    # ── Final slide: Budget Shift Recommendations ─────────────────────────────
    s_last = _new_blank_slide(prs)
    _tb(s_last, 0.4, 0.17, 12.5, 0.65,
        "Budget Shift Recommendations", 20, bold=True)
    _box(s_last, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)

    budget_raw = _ai_pptx(
        api_key,
        f"""Campaign performance data:
{data_text}

For each brand, provide one budget reallocation recommendation.
Format EXACTLY as pipe-separated lines — one line per brand, no header row:
Brand Name | Best Performing IO | One-sentence recommendation starting with an action verb

Example:
Nike Summer 2024 | Display | Increase Display budget by 20% — lowest CPM at $3.20
Coke Q3 | YouTube | Shift 15% from Display to YouTube — VTR of 68% outperforms

Only include brands from the data. Cite actual numbers. No extra text.""",
        max_tokens=350,
    )
    _add_budget_table(s_last, budget_raw, x=0.4, y=1.08, w=12.5, h=5.9)

    _slide_footer(s_last, slide_num)

    # Serialise to buffer and return
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ── Sidebar ───────────────────────────────────────────────────────────────────
st.sidebar.markdown(
    "<div style='padding:8px 0 14px 0;'>"
    "<span style='font-size:22px;font-weight:700;color:#FFFFFF;"
    "font-family:Inter,system-ui,sans-serif;letter-spacing:-0.01em;'>"
    "Insights App</span></div>"
    "<hr style='border:none;border-top:1px solid rgba(255,255,255,0.25);margin:0 0 14px 0;'>",
    unsafe_allow_html=True,
)
st.sidebar.subheader("Loaded Files")

# ── Main header ───────────────────────────────────────────────────────────────
st.title("Ad Tech Insights & Reporting")
st.markdown(
    "<p style='color:#6b7280;font-size:14px;margin-top:-12px;'>"
    "Upload CSV exports from DV360, TTD or any DSP to generate insights and reports."
    "</p>",
    unsafe_allow_html=True,
)

# ── File uploader ─────────────────────────────────────────────────────────────
st.subheader("Upload DSP Export Files")

uploaded_files = st.file_uploader(
    "Drag and drop CSV files here, or click to browse. Multiple files accepted.",
    type=["csv"],
    accept_multiple_files=True,
)

# ── Process files ─────────────────────────────────────────────────────────────
if not uploaded_files:
    # Show a friendly prompt when nothing is uploaded yet
    st.markdown("""
    <div style='text-align:center;padding:48px 0;color:#9ca3af;'>
        <div style='font-size:48px;margin-bottom:12px;'>📂</div>
        <div style='font-size:16px;font-weight:600;'>No files uploaded yet</div>
        <div style='font-size:13px;margin-top:6px;'>Supports DV360 and TTD CSV exports</div>
    </div>
    """, unsafe_allow_html=True)
    st.sidebar.markdown("_No files loaded yet._")

else:
    # Load and stack all uploaded files into one DataFrame
    frames = []
    file_info = []   # used to populate the sidebar

    with st.spinner("Reading and normalising files…"):
        for f in uploaded_files:
            df_file, dsp = load_and_normalise(f)
            frames.append(df_file)
            file_info.append({"name": f.name, "rows": len(df_file), "dsp": dsp})

    # Combine all files into a single DataFrame
    df_all = pd.concat(frames, ignore_index=True)

    # ── Sidebar file list ─────────────────────────────────────────────────────
    badge_class = {"DV360": "badge-dv360", "TTD": "badge-ttd", "Generic": "badge-generic"}
    sidebar_html = ""
    for fi in file_info:
        bc = badge_class.get(fi["dsp"], "badge-generic")
        sidebar_html += (
            f"<div style='margin-bottom:10px;font-size:13px;'>"
            f"<b>{fi['name']}</b>"
            f"<span class='source-badge {bc}'>{fi['dsp']}</span>"
            f"<br><span style='color:#6b7280;'>{fi['rows']:,} rows</span>"
            f"</div>"
        )
    st.sidebar.markdown(sidebar_html, unsafe_allow_html=True)

    # ── Summary metrics ───────────────────────────────────────────────────────
    st.subheader("Summary Metrics")

    total_impressions = df_all["impressions"].sum()  if "impressions" in df_all.columns else 0
    total_clicks      = df_all["clicks"].sum()       if "clicks"      in df_all.columns else 0
    total_spend       = df_all["spend_usd"].sum()    if "spend_usd"   in df_all.columns else 0
    avg_ctr           = df_all["ctr"].mean()         if "ctr"         in df_all.columns else None
    avg_cpm           = df_all["cpm"].mean()         if "cpm"         in df_all.columns else None

    col1, col2, col3, col4, col5 = st.columns(5)
    col1.metric("Total Impressions", f"{total_impressions:,.0f}")
    col2.metric("Total Clicks",      f"{total_clicks:,.0f}")
    col3.metric("Total Spend",       f"${total_spend:,.0f}")
    col4.metric("Avg CTR",           f"{avg_ctr:.2%}"    if avg_ctr is not None else "N/A")
    col5.metric("Avg CPM",           f"${avg_cpm:,.2f}"  if avg_cpm is not None else "N/A")

    # Build campaign-level summary now so it's available for both insights and PPTX export
    if "campaign" in df_all.columns:
        _agg = {c: (c, "sum") for c in ["impressions", "clicks", "spend_usd"] if c in df_all.columns}
        camp_summary = df_all.groupby("campaign").agg(**_agg).reset_index()
        if "impressions" in camp_summary.columns and "clicks" in camp_summary.columns:
            camp_summary["ctr"] = camp_summary["clicks"] / camp_summary["impressions"]
        if "impressions" in camp_summary.columns and "spend_usd" in camp_summary.columns:
            camp_summary["cpm"] = camp_summary["spend_usd"] / camp_summary["impressions"] * 1000
    else:
        camp_summary = pd.DataFrame()

    # ── Charts ────────────────────────────────────────────────────────────────
    st.subheader("Performance Charts")

    # Colour palette — one colour per brand, consistent across all charts
    PALETTE = ["#2563EB", "#60A5FA", "#93C5FD", "#1D4ED8", "#3B82F6",
               "#7C3AED", "#A78BFA", "#10B981", "#F59E0B", "#EF4444"]

    def no_data_msg(msg):
        st.markdown(
            f"<div style='text-align:center;padding:40px;color:#9ca3af;"
            f"border:1px dashed #e5e7eb;border-radius:8px;font-size:13px;'>{msg}</div>",
            unsafe_allow_html=True,
        )

    def trunc(s, n=32):
        """Truncate a string to n characters for axis labels."""
        return s[:n] + "…" if len(str(s)) > n else str(s)

    def apply_chart_style(fig, xaxis_title="", yaxis_title="", horizontal=False):
        """
        Apply consistent white-card styling to every Plotly chart.
        horizontal=True swaps which axis gets the grid.
        """
        fig.update_layout(
            plot_bgcolor="#FFFFFF",
            paper_bgcolor="#FFFFFF",
            font=dict(family="Inter, system-ui, sans-serif", size=13, color="#374151"),
            margin=dict(t=36, b=56, l=70, r=24),
            showlegend=False,
            hoverlabel=dict(
                bgcolor="#FFFFFF",
                font_size=13,
                font_family="Inter, system-ui, sans-serif",
                bordercolor="#E5E7EB",
            ),
        )
        # X-axis styling
        fig.update_xaxes(
            showgrid=horizontal,          # grid on x only for horizontal bar charts
            gridcolor="#F3F4F6",
            gridwidth=1,
            linecolor="#E5E7EB",
            tickfont=dict(size=11, color="#374151"),
            title_text=xaxis_title,
            title_font=dict(size=12, color="#374151"),
        )
        # Y-axis styling
        fig.update_yaxes(
            showgrid=not horizontal,      # grid on y only for vertical bar charts
            gridcolor="#F3F4F6",
            gridwidth=1,
            linecolor="#E5E7EB",
            tickfont=dict(size=11, color="#374151"),
            title_text=yaxis_title,
            title_font=dict(size=12, color="#374151"),
        )
        return fig

    # ── Row 1: Spend by Brand | CPM by Brand ──────────────────────────────────
    r1_left, r1_right = st.columns(2)

    # Chart 1 — Total Spend by Brand
    with r1_left:
        st.markdown("**Total Spend by Brand**")
        if "spend_usd" in df_all.columns and "campaign" in df_all.columns:
            c1 = (df_all.groupby("campaign")["spend_usd"]
                  .sum().reset_index()
                  .sort_values("spend_usd", ascending=False))
            colors1 = [PALETTE[i % len(PALETTE)] for i in range(len(c1))]

            # Human-readable label for each bar and tooltip
            spend_labels = [
                f"${v/1e6:.2f}M" if v >= 1e6 else f"${v/1e3:.0f}K"
                for v in c1["spend_usd"]
            ]
            fig = go.Figure(go.Bar(
                x=c1["campaign"],
                y=c1["spend_usd"],
                marker_color=colors1,
                text=spend_labels,
                textposition="outside",
                textfont=dict(size=10, color="#374151"),
                customdata=spend_labels,
                hovertemplate="<b>%{x}</b><br>Spend: %{customdata}<extra></extra>",
            ))
            fig.update_layout(
                bargap=0.45,
                yaxis_range=[0, c1["spend_usd"].max() * 1.28],
            )
            # Y-axis: SI prefix ticks ($300k, $1.2M)
            fig.update_yaxes(tickprefix="$", tickformat=".2s")
            apply_chart_style(fig, yaxis_title="Total Spend (USD)")
            st.plotly_chart(fig, use_container_width=True, config={"displaylogo": False})
        else:
            no_data_msg("No campaign or spend column found.")

    # Chart 2 — CPM by Brand (recalculated from totals for accuracy)
    with r1_right:
        st.markdown("**CPM by Brand**")
        if not camp_summary.empty and "cpm" in camp_summary.columns:
            c2 = (camp_summary[["campaign", "cpm"]]
                  .dropna()
                  .sort_values("cpm", ascending=False))
            colors2 = [PALETTE[i % len(PALETTE)] for i in range(len(c2))]

            cpm_labels = [f"${v:,.2f}" for v in c2["cpm"]]
            fig = go.Figure(go.Bar(
                x=c2["campaign"],
                y=c2["cpm"],
                marker_color=colors2,
                text=cpm_labels,
                textposition="outside",
                textfont=dict(size=10, color="#374151"),
                customdata=cpm_labels,
                hovertemplate="<b>%{x}</b><br>CPM: %{customdata}<extra></extra>",
            ))
            fig.update_layout(
                bargap=0.45,
                yaxis_range=[0, c2["cpm"].max() * 1.28],
            )
            fig.update_yaxes(tickprefix="$", tickformat=",.2f")
            apply_chart_style(fig, yaxis_title="CPM (USD)")
            st.plotly_chart(fig, use_container_width=True, config={"displaylogo": False})
        else:
            no_data_msg("Spend and impressions data needed to calculate CPM.")

    # ── Row 2: Best / Worst line items by CPM ─────────────────────────────────
    # Brand filter options — "All Brands" + each individual brand
    _li_brands = ["All Brands"] + (
        camp_summary["campaign"].tolist() if not camp_summary.empty else []
    )
    # Column to group by — prefer line_item, fall back to campaign
    li_col = "line_item" if "line_item" in df_all.columns else "campaign"

    r2_left, r2_right = st.columns(2)

    # Chart 3 — Best performing line items by CPM (lowest CPM = most cost-efficient)
    with r2_left:
        st.markdown("**Best Performing Line Items by CPM**")
        li_brand_best = st.selectbox(
            "Filter by brand", _li_brands, key="li_best_brand",
        )
        if "spend_usd" in df_all.columns and "impressions" in df_all.columns:
            df_li = df_all.copy()
            # Apply brand filter if one is selected
            if li_brand_best != "All Brands" and "campaign" in df_li.columns:
                df_li = df_li[df_li["campaign"] == li_brand_best]
            agg_li = (df_li.groupby(li_col)
                      .agg(spend_usd=("spend_usd", "sum"),
                           impressions=("impressions", "sum"))
                      .reset_index())
            # Only keep rows with enough impressions to give a meaningful CPM
            agg_li = agg_li[agg_li["impressions"] > 0].copy()
            agg_li["cpm"] = agg_li["spend_usd"] / agg_li["impressions"] * 1000
            # Best = lowest CPM; take 10, then reverse so best is at the top of the chart
            top_best = (agg_li.sort_values("cpm", ascending=True)
                        .head(10)
                        .sort_values("cpm", ascending=False))
            labels_best = [trunc(s) for s in top_best[li_col]]

            best_labels = [f"${v:,.2f}" for v in top_best["cpm"]]
            fig = go.Figure(go.Bar(
                x=top_best["cpm"],
                y=labels_best,
                orientation="h",
                marker_color="#2563EB",
                text=best_labels,
                textposition="outside",
                textfont=dict(size=10, color="#374151"),
                customdata=best_labels,
                hovertemplate="<b>%{y}</b><br>CPM: %{customdata}<extra></extra>",
            ))
            fig.update_layout(
                bargap=0.3,
                xaxis_range=[0, top_best["cpm"].max() * 1.35],
            )
            fig.update_xaxes(tickprefix="$", tickformat=",.2f")
            apply_chart_style(fig, xaxis_title="CPM (USD)", horizontal=True)
            st.plotly_chart(fig, use_container_width=True, config={"displaylogo": False})
        else:
            no_data_msg("Spend and impressions data needed to calculate CPM.")

    # Chart 4 — Worst performing line items by CPM (highest CPM = least efficient)
    with r2_right:
        st.markdown("**Worst Performing Line Items by CPM**")
        li_brand_worst = st.selectbox(
            "Filter by brand", _li_brands, key="li_worst_brand",
        )
        if "spend_usd" in df_all.columns and "impressions" in df_all.columns:
            df_li2 = df_all.copy()
            # Apply brand filter if one is selected
            if li_brand_worst != "All Brands" and "campaign" in df_li2.columns:
                df_li2 = df_li2[df_li2["campaign"] == li_brand_worst]
            agg_li2 = (df_li2.groupby(li_col)
                       .agg(spend_usd=("spend_usd", "sum"),
                            impressions=("impressions", "sum"))
                       .reset_index())
            agg_li2 = agg_li2[agg_li2["impressions"] > 0].copy()
            agg_li2["cpm"] = agg_li2["spend_usd"] / agg_li2["impressions"] * 1000
            # Worst = highest CPM; keep descending so worst sits at the top of the chart
            top_worst = agg_li2.sort_values("cpm", ascending=False).head(10)
            labels_worst = [trunc(s) for s in top_worst[li_col]]

            worst_labels = [f"${v:,.2f}" for v in top_worst["cpm"]]
            fig = go.Figure(go.Bar(
                x=top_worst["cpm"],
                y=labels_worst,
                orientation="h",
                marker_color="#60A5FA",
                text=worst_labels,
                textposition="outside",
                textfont=dict(size=10, color="#374151"),
                customdata=worst_labels,
                hovertemplate="<b>%{y}</b><br>CPM: %{customdata}<extra></extra>",
            ))
            fig.update_layout(
                bargap=0.3,
                xaxis_range=[0, top_worst["cpm"].max() * 1.35],
            )
            fig.update_xaxes(tickprefix="$", tickformat=",.2f")
            apply_chart_style(fig, xaxis_title="CPM (USD)", horizontal=True)
            st.plotly_chart(fig, use_container_width=True, config={"displaylogo": False})
        else:
            no_data_msg("Spend and impressions data needed to calculate CPM.")

    # ── AI Insights ───────────────────────────────────────────────────────────
    st.subheader("AI Brand Insights")

    # Get the Anthropic API key — check Streamlit secrets first, then env var
    api_key = (
        st.secrets.get("ANTHROPIC_API_KEY")
        if "ANTHROPIC_API_KEY" in st.secrets
        else os.environ.get("ANTHROPIC_API_KEY")
    )

    if not api_key:
        st.warning(
            "No Anthropic API key found. Add `ANTHROPIC_API_KEY` to your "
            "Streamlit secrets or environment variables to enable AI insights."
        )
    elif "campaign" not in df_all.columns:
        st.info("A 'Campaign Name' or 'Brand Name' column is required to generate per-brand insights.")
    else:
        # Build an aggregated summary per campaign (recalculate from totals, not averages)
        has_spend  = "spend_usd"   in df_all.columns
        has_clicks = "clicks"      in df_all.columns
        has_imps   = "impressions" in df_all.columns

        agg_cols = {c: (c, "sum") for c in ["impressions", "clicks", "spend_usd"] if c in df_all.columns}
        camp_summary = df_all.groupby("campaign").agg(**agg_cols).reset_index()

        if has_imps and has_clicks:
            camp_summary["ctr"] = camp_summary["clicks"] / camp_summary["impressions"]
        if has_imps and has_spend:
            camp_summary["cpm"] = camp_summary["spend_usd"] / camp_summary["impressions"] * 1000

        # Format the data as a text table for the prompt.
        # Includes campaign-level totals AND a campaign × environment breakdown
        # when an environment column is present (gives Claude the Web/App/YouTube split).
        def format_summary_table(df):
            lines = ["=== Brand Totals ==="]
            for _, r in camp_summary.iterrows():
                parts = [f"Brand: {r['campaign']}"]
                if "impressions" in r and pd.notna(r["impressions"]):
                    parts.append(f"Impressions: {int(r['impressions']):,}")
                if "clicks" in r and pd.notna(r["clicks"]):
                    parts.append(f"Clicks: {int(r['clicks']):,}")
                if "spend_usd" in r and pd.notna(r["spend_usd"]):
                    parts.append(f"Spend: ${r['spend_usd']:,.0f}")
                if "ctr" in r and pd.notna(r["ctr"]):
                    parts.append(f"CTR: {r['ctr']:.2%}")
                if "cpm" in r and pd.notna(r["cpm"]):
                    parts.append(f"CPM: ${r['cpm']:,.2f}")
                lines.append(" | ".join(parts))

            # Helper: aggregate a dimension column and append formatted rows to lines
            def _add_breakdown(label_col, section_title, label_key):
                if label_col not in df.columns:
                    return
                lines.append(f"\n=== {section_title} ===")
                agg = {c: (c, "sum") for c in ["impressions", "clicks", "spend_usd"]
                       if c in df.columns}
                grp = df.groupby(["campaign", label_col]).agg(**agg).reset_index()
                if "impressions" in grp.columns and "clicks" in grp.columns:
                    grp["ctr"] = grp["clicks"] / grp["impressions"]
                if "impressions" in grp.columns and "spend_usd" in grp.columns:
                    grp["cpm"] = grp["spend_usd"] / grp["impressions"] * 1000
                if "clicks" in grp.columns and "spend_usd" in grp.columns:
                    grp["cpc"] = grp["spend_usd"] / grp["clicks"]
                for _, r in grp.iterrows():
                    parts = [f"Brand: {r['campaign']} | {label_key}: {r[label_col]}"]
                    if "impressions" in r and pd.notna(r["impressions"]):
                        parts.append(f"Impressions: {int(r['impressions']):,}")
                    if "clicks" in r and pd.notna(r["clicks"]):
                        parts.append(f"Clicks: {int(r['clicks']):,}")
                    if "spend_usd" in r and pd.notna(r["spend_usd"]):
                        parts.append(f"Spend: ${r['spend_usd']:,.0f}")
                    if "ctr" in r and pd.notna(r["ctr"]):
                        parts.append(f"CTR: {r['ctr']:.2%}")
                    if "cpm" in r and pd.notna(r["cpm"]):
                        parts.append(f"CPM: ${r['cpm']:,.2f}")
                    if "cpc" in r and pd.notna(r["cpc"]):
                        parts.append(f"CPC: ${r['cpc']:,.2f}")
                    lines.append(" | ".join(parts))

            # Breakdown by Environment (Web, App, YouTube, etc.)
            _add_breakdown("environment", "Breakdown by Brand and Environment", "Environment")

            # Breakdown by Insertion Order / Line Item name
            _add_breakdown("line_item", "Breakdown by Brand and Line Item / Creative", "Line Item")

            # Breakdown by Device Type
            _add_breakdown("device_type", "Breakdown by Brand and Device Type", "Device Type")

            return "\n".join(lines)

        all_campaigns_text = format_summary_table(df_all)
        campaign_list = camp_summary["campaign"].tolist()
        # Store brand names in session_state so other pages (e.g. Brand Memory) can read them
        st.session_state["campaign_list"] = campaign_list

        # System prompt — analyst persona, structured output required
        SYSTEM_PROMPT = (
            "You are a senior programmatic advertising analyst with deep expertise in "
            "DSP campaign performance (DV360, TTD). You write clear, concise, data-driven "
            "commentary for ad tech professionals. Be specific — reference actual numbers. "
            "Follow the section headings and structure exactly as specified in each prompt. "
            "Do not add extra sections or deviate from the requested format."
        )

        def stream_insight(campaign_name: str, all_data: str, brand_context: str = ""):
            """
            Stream per-brand analysis using the standard Display / Video / YouTube
            insertion-order structure. Brand memory is injected as an override block
            at the end of the prompt so it takes priority over all default instructions.
            Yields text chunks from the Claude API.
            """
            client_ai = anthropic.Anthropic(api_key=api_key)

            # Default structure: Campaign Overview + three insertion-order sections.
            # The AI skips any section whose insertion order is absent from the data.
            default_structure = (
                f"**{campaign_name} - Campaign Overview**\n"
                f"Summarise overall campaign performance. Focus on total Revenue (Spend) "
                f"and total Impressions for this brand. Keep to 2-3 sentences.\n\n"
                f"**Display**\n"
                f"Summarise performance for the Display insertion order, focusing on "
                f"average CPM. Then identify:\n"
                f"- Best performing Line Item by CPM (state exact CPM, CPC, CTR)\n"
                f"- Worst performing Line Item by CPM (state exact CPM, CPC, CTR)\n"
                f"- Best performing Creative by CPM (state exact CPM, CPC, CTR)\n"
                f"- Worst performing Creative by CPM (state exact CPM, CPC, CTR)\n"
                f"If no Display insertion order data exists for this brand, omit this "
                f"section entirely — do not mention it at all.\n\n"
                f"**Video**\n"
                f"Summarise performance for the Video insertion order, focusing on CPV "
                f"and VTR. Then identify:\n"
                f"- Best performing Line Item by CPV (state exact CPV, VTR)\n"
                f"- Worst performing Line Item by CPV (state exact CPV, VTR)\n"
                f"- Best performing Creative by CPV (state exact CPV, VTR)\n"
                f"- Worst performing Creative by CPV (state exact CPV, VTR)\n"
                f"If no Video insertion order data exists for this brand, omit this "
                f"section entirely — do not mention it at all.\n\n"
                f"**YouTube**\n"
                f"Summarise performance for the YouTube insertion order, focusing on CPV "
                f"and VTR. Then identify:\n"
                f"- Best performing Line Item by CPV (state exact CPV, VTR)\n"
                f"- Worst performing Line Item by CPV (state exact CPV, VTR)\n"
                f"- Best performing Creative by CPV (state exact CPV, VTR)\n"
                f"- Worst performing Creative by CPV (state exact CPV, VTR)\n"
                f"If no YouTube insertion order data exists for this brand, omit this "
                f"section entirely — do not mention it at all."
            )

            # Brand memory override — appended after the default instructions so it
            # takes explicit priority. Only included when brand context exists.
            if brand_context:
                override_section = (
                    f"\n\nBRAND MEMORY OVERRIDE — These instructions take priority over "
                    f"all default instructions above. Where there is any conflict, always "
                    f"follow these brand-specific instructions instead:\n\n"
                    f"{brand_context}"
                )
            else:
                override_section = ""

            prompt = (
                f"Here is the full performance data for all campaigns in this report:\n\n"
                f"{all_data}\n\n"
                f"Write the analysis for the '{campaign_name}' brand using exactly "
                f"this structure and these headings (use ** for bold):\n\n"
                f"{default_structure}"
                f"{override_section}\n\n"
                f"Use only the data provided — do not invent numbers."
            )

            with client_ai.messages.stream(
                model="claude-sonnet-4-6",
                max_tokens=1200,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text in stream.text_stream:
                    yield text

        def stream_overall(all_data: str):
            """
            Stream the cross-campaign summary sections:
            Summary, Best & Worst Performers, Optimisation Recommendations.
            Yields text chunks from the Claude API.
            """
            client_ai = anthropic.Anthropic(api_key=api_key)

            prompt = (
                f"Here is the full performance data for all campaigns in this report:\n\n"
                f"{all_data}\n\n"
                f"Write the following three sections using exactly these headings "
                f"(use ** for bold):\n\n"
                f"**Summary**\n"
                f"Overall commentary across all campaigns and environments. 3-4 sentences.\n\n"
                f"**Best & Worst Performers**\n"
                f"- Best performing device type and why\n"
                f"- Worst performing device type and why\n"
                f"- Best performing creative and why\n"
                f"- Worst performing creative and why\n\n"
                f"**Optimisation Recommendations**\n"
                f"3-5 actionable recommendations based on the data above. Use bullet points.\n\n"
                f"Use only the data provided. Be specific with numbers where possible."
            )

            with client_ai.messages.stream(
                model="claude-sonnet-4-6",
                max_tokens=600,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text in stream.text_stream:
                    yield text

        # STYLE LOCK: Do not change insights display formatting
        def _insight_html(text):
            """
            Convert the subset of markdown used in insights to explicitly-styled HTML.
            Forces ALL text to plain black (#111111) so Streamlit's theme accent color
            (teal/cyan) is never applied to bold headers or any other element.

            Handles:
              **text**  → <strong> (black, no color override from theme)
              - item    → <ul><li> list
              blank line → paragraph break
            """
            lines  = text.split("\n")
            parts  = []
            in_ul  = False
            p_style = "margin:4px 0;color:#111111;font-size:14px;line-height:1.7;"
            li_style = "margin:2px 0;color:#111111;font-size:14px;"

            def apply_bold(s):
                # Replace **…** with <strong> that explicitly inherits black
                return re.sub(r'\*\*(.+?)\*\*',
                              r'<strong style="color:#111111;">\1</strong>', s)

            for line in lines:
                stripped = line.strip()
                if not stripped:
                    if in_ul:
                        parts.append("</ul>")
                        in_ul = False
                    continue
                if stripped.startswith("- "):
                    if not in_ul:
                        parts.append("<ul style='margin:6px 0 6px 18px;padding:0;'>")
                        in_ul = True
                    parts.append(f"<li style='{li_style}'>{apply_bold(stripped[2:])}</li>")
                else:
                    if in_ul:
                        parts.append("</ul>")
                        in_ul = False
                    parts.append(f"<p style='{p_style}'>{apply_bold(stripped)}</p>")

            if in_ul:
                parts.append("</ul>")

            return (
                "<div style='font-family:-apple-system,Segoe UI,Roboto,sans-serif;"
                "color:#111111;'>"
                + "\n".join(parts)
                + "</div>"
            )

        # Generate Insights button — analysis only runs when clicked
        if st.button("✨ Generate Insights", type="primary"):
            st.session_state["insights_triggered"] = True

        if st.session_state.get("insights_triggered"):
            # Load brand memory once so it's available for every campaign
            _insights_bm = load_brand_memory()

            for campaign in campaign_list:
                # Look up brand context using partial matching (case-insensitive).
                # e.g. brand key "Nike" matches campaign "Nike Summer 2024".
                _campaign_brand_context = ""
                for _bm_key, _bm_val in _insights_bm.items():
                    if _bm_key.lower() in campaign.lower():
                        _campaign_brand_context = _bm_val.get("rationale", "")
                        break

                # One expander per campaign, collapsed by default
                with st.expander(f"**{campaign}**", expanded=False):
                    # Show indicator when brand memory is applied
                    if _campaign_brand_context:
                        st.markdown(
                            "<p style='color:#2563EB;font-size:13px;font-weight:600;"
                            "margin:0 0 8px 0;'>✓ Brand context applied</p>",
                            unsafe_allow_html=True,
                        )
                    placeholder = st.empty()
                    accumulated = ""
                    for chunk in stream_insight(campaign, all_campaigns_text,
                                                _campaign_brand_context):
                        accumulated += chunk
                        placeholder.markdown(_insight_html(accumulated),
                                             unsafe_allow_html=True)
                    # Save completed text for PPTX export
                    st.session_state.setdefault("insights_text", {})[campaign] = accumulated

            # Overall summary — runs once after all per-campaign sections
            with st.expander("**📊 Summary & Recommendations**", expanded=False):
                placeholder = st.empty()
                accumulated = ""
                for chunk in stream_overall(all_campaigns_text):
                    accumulated += chunk
                    placeholder.markdown(_insight_html(accumulated),
                                         unsafe_allow_html=True)
                st.session_state["insights_overall"] = accumulated

    # ── Data preview ──────────────────────────────────────────────────────────
    st.subheader("Data Preview")

    # Show total rows and which files contributed
    total_rows  = len(df_all)
    file_count  = len(uploaded_files)
    st.markdown(
        f"<p style='color:#6b7280;font-size:13px;margin-top:-8px;'>"
        f"{total_rows:,} rows loaded from {file_count} file{'s' if file_count > 1 else ''}.</p>",
        unsafe_allow_html=True,
    )

    # Format the preview table — keep raw numbers readable
    preview_df = df_all.copy()
    if "impressions" in preview_df.columns:
        preview_df["impressions"] = preview_df["impressions"].apply(
            lambda x: f"{x:,.0f}" if pd.notna(x) else ""
        )
    if "clicks" in preview_df.columns:
        preview_df["clicks"] = preview_df["clicks"].apply(
            lambda x: f"{x:,.0f}" if pd.notna(x) else ""
        )
    if "spend_usd" in preview_df.columns:
        preview_df["spend_usd"] = preview_df["spend_usd"].apply(
            lambda x: f"${x:,.2f}" if pd.notna(x) else ""
        )
    if "ctr" in preview_df.columns:
        preview_df["ctr"] = preview_df["ctr"].apply(
            lambda x: f"{x:.2%}" if pd.notna(x) else ""
        )
    if "cpm" in preview_df.columns:
        preview_df["cpm"] = preview_df["cpm"].apply(
            lambda x: f"${x:,.2f}" if pd.notna(x) else ""
        )

    st.dataframe(preview_df, use_container_width=True)

    # ── Sidebar: Export ───────────────────────────────────────────────────────
    st.sidebar.markdown("---")
    st.sidebar.subheader("Export")

    # Load the API key for report generation (same source as insights section)
    _rpt_key = (
        st.secrets.get("ANTHROPIC_API_KEY")
        if "ANTHROPIC_API_KEY" in st.secrets
        else os.environ.get("ANTHROPIC_API_KEY")
    )

    # "Generate Report" builds the PPTX (makes AI calls per brand — takes ~10–30s)
    if st.sidebar.button("📊 Generate Report", type="primary", key="gen_report_btn"):
        if not _rpt_key:
            st.sidebar.warning(
                "Add ANTHROPIC_API_KEY to Streamlit secrets to generate AI-powered slides."
            )
        elif camp_summary.empty:
            st.sidebar.warning("Upload a CSV file first.")
        else:
            with st.spinner(
                "Building report — generating AI slides for each brand… "
                "This takes about 10–30 seconds."
            ):
                try:
                    pptx_buf = build_pptx_report(_rpt_key, camp_summary, df_all)
                    st.session_state["pptx_report"] = pptx_buf
                    st.sidebar.success("Report ready — click Download below.")
                except Exception as e:
                    st.error(f"Report generation failed: {e}")

    # Show download button once the report has been generated
    if st.session_state.get("pptx_report"):
        st.sidebar.download_button(
            label     = "📥 Download Report (.pptx)",
            data      = st.session_state["pptx_report"],
            file_name = f"{date.today():%Y-%m-%d}_campaign_report.pptx",
            mime      = "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
