import io
import json
import os
import random
from datetime import date, datetime, timedelta

import anthropic
import pandas as pd
import plotly.graph_objects as go
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# STYLE LOCK: Do not remove or modify this CSS block.
st.markdown("""
<style>
    html, body, [class*="css"] {
        font-family: "Inter", system-ui, -apple-system, "Segoe UI", sans-serif;
        font-size: 15px;
        color: #374151;
    }
    .stApp { background-color: #F3F4F6; }
    h1, h2, h3, h4, h5, h6 { font-weight: 700 !important; color: #111827 !important; }
    h2, h3 {
        margin-top: 2rem !important;
        padding-top: 0.25rem !important;
        border-bottom: none !important;
        padding-bottom: 0 !important;
        border-left: none !important;
        padding-left: 0 !important;
    }
    [data-testid="metric-container"] {
        background: #FFFFFF;
        border: none;
        border-radius: 16px;
        padding: 20px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        border-top: 4px solid #7C3AED;
    }
    [data-testid="metric-container"] label {
        font-size: 12px;
        font-weight: 600;
        color: #6B7280;
        text-transform: uppercase;
        letter-spacing: 0.05em;
    }
    [data-testid="metric-container"] [data-testid="stMetricValue"] {
        font-size: 26px;
        font-weight: 700;
        color: #111827;
    }
    .element-container:has([data-testid="stPlotlyChart"]) {
        background: #FFFFFF;
        border-radius: 12px;
        padding: 16px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
    }
    .stButton > button[kind="primary"],
    [data-testid="baseButton-primary"] {
        background-color: #7C3AED !important;
        color: #FFFFFF !important;
        border: none !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        font-size: 15px !important;
        padding: 0.5rem 1.25rem !important;
    }
    .stButton > button[kind="primary"]:hover,
    [data-testid="baseButton-primary"]:hover { background-color: #6D28D9 !important; }
    .stButton > button[kind="secondary"],
    [data-testid="baseButton-secondary"] { border-radius: 8px !important; font-weight: 600 !important; }
    [data-testid="stDataFrame"] {
        background: #FFFFFF;
        border-radius: 12px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.06);
    }
</style>
""", unsafe_allow_html=True)
# STYLE LOCK

# ── Channel configuration ─────────────────────────────────────────────────────
# Each entry defines a channel/sub-channel with realistic daily spend and
# performance ranges for mock data generation.
CHANNELS = [
    {"channel": "Google Search",  "display": "Google Search Brand",
     "campaign": "Optus_SEM_Brand_May2026",       "ad_group": "Brand Core Terms",
     "daily_spend_range": (800, 1400),   "ctr_range": (4.5, 7.2),
     "cpa_range": (18, 28),  "cpo_range": (22, 35),
     "conv_rate_range": (3.2, 5.8),  "impressions_range": (12000, 22000)},
    {"channel": "Google Search",  "display": "Google Search NonBrand",
     "campaign": "Optus_SEM_NonBrand_May2026",    "ad_group": "Mobile Plans",
     "daily_spend_range": (1200, 2200),  "ctr_range": (2.1, 3.8),
     "cpa_range": (35, 55),  "cpo_range": (40, 65),
     "conv_rate_range": (1.8, 3.2),  "impressions_range": (35000, 65000)},
    {"channel": "Programmatic",   "display": "Programmatic Display",
     "campaign": "Optus_DV360_Prospecting_May2026","ad_group": "Display Broad",
     "daily_spend_range": (600, 1100),   "ctr_range": (0.15, 0.45),
     "cpa_range": (45, 75),  "cpo_range": (55, 85),
     "conv_rate_range": (0.8, 1.8),  "impressions_range": (80000, 150000)},
    {"channel": "Programmatic",   "display": "Programmatic Retargeting",
     "campaign": "Optus_DV360_Retargeting_May2026","ad_group": "Site Visitors",
     "daily_spend_range": (400, 800),    "ctr_range": (0.45, 0.95),
     "cpa_range": (28, 45),  "cpo_range": (32, 52),
     "conv_rate_range": (1.8, 3.5),  "impressions_range": (40000, 80000)},
    {"channel": "Meta",           "display": "Meta Prospecting",
     "campaign": "Optus_Meta_Prospecting_May2026", "ad_group": "Lookalike Audience",
     "daily_spend_range": (700, 1300),   "ctr_range": (1.2, 2.8),
     "cpa_range": (38, 62),  "cpo_range": (45, 72),
     "conv_rate_range": (1.2, 2.4),  "impressions_range": (45000, 90000)},
    {"channel": "Meta",           "display": "Meta Retargeting",
     "campaign": "Optus_Meta_Retargeting_May2026", "ad_group": "Website Custom Audience",
     "daily_spend_range": (350, 650),    "ctr_range": (1.8, 3.5),
     "cpa_range": (25, 42),  "cpo_range": (30, 48),
     "conv_rate_range": (2.2, 4.0),  "impressions_range": (18000, 35000)},
    {"channel": "Affiliates",     "display": "Affiliates",
     "campaign": "Optus_Affiliates_May2026",       "ad_group": "Comparison Sites",
     "daily_spend_range": (500, 900),    "ctr_range": (8.5, 14.2),
     "cpa_range": (15, 25),  "cpo_range": (18, 28),
     "conv_rate_range": (6.5, 11.0),  "impressions_range": (3000, 7000)},
]

# Default CPA targets per display name (used for alert thresholds)
DEFAULT_CPA_TARGETS = {
    "Google Search Brand":      25,
    "Google Search NonBrand":   45,
    "Programmatic Display":     60,
    "Programmatic Retargeting": 38,
    "Meta Prospecting":         50,
    "Meta Retargeting":         35,
    "Affiliates":               20,
}

# Blue palette for charts (7 entries, one per channel)
CHART_COLOURS = ["#1E40AF","#2563EB","#3B82F6","#60A5FA","#93C5FD","#BFDBFE","#7C3AED"]

# Column name candidates for mapping uploaded file headers to standard names
UPLOAD_COLUMN_CANDIDATES = {
    "Channel":     ["channel", "platform", "source", "network"],
    "Campaign":    ["campaign", "campaign name"],
    "Ad Group":    ["ad group", "ad set", "line item"],
    "Date":        ["date"],
    "Spend":       ["spend", "cost", "amount spent"],
    "Impressions": ["impressions"],
    "Clicks":      ["clicks"],
    "Conversions": ["conversions", "orders", "leads"],
    "CPA":         ["cpa", "cost per acquisition", "cost per lead"],
    "CPO":         ["cpo", "cost per order"],
    "CTR":         ["ctr", "click through rate"],
    "Conv Rate":   ["conv. rate", "conversion rate"],
}


# ── Mock data generation ──────────────────────────────────────────────────────

@st.cache_data
def generate_mock_data():
    """Generate 30 days of realistic mock data for all channels.
    Uses a fixed random seed so the numbers are consistent across reruns."""
    rng = random.Random(42)
    rows = []

    # Loop through 30 days in May 2026
    for day_offset in range(30):
        current_date = date(2026, 5, 1) + timedelta(days=day_offset)

        # Generate a row for each channel config
        for ch in CHANNELS:
            impressions = rng.randint(*ch["impressions_range"])
            ctr = rng.uniform(*ch["ctr_range"])
            clicks = max(round(impressions * ctr / 100), 1)
            conv_rate = rng.uniform(*ch["conv_rate_range"])
            conversions = max(round(clicks * conv_rate / 100), 0)
            spend = round(rng.uniform(*ch["daily_spend_range"]), 2)
            # CPA = spend divided by conversions (avoid divide-by-zero)
            cpa = round(spend / max(conversions, 1), 2)
            # CPO is slightly higher than CPA to model order-level cost
            cpo = round(cpa * rng.uniform(1.05, 1.25), 2)

            rows.append({
                "Date":         pd.Timestamp(current_date),
                "Channel":      ch["channel"],
                "Display Name": ch["display"],
                "Campaign":     ch["campaign"],
                "Ad Group":     ch["ad_group"],
                "Spend":        spend,
                "Impressions":  impressions,
                "Clicks":       clicks,
                "Conversions":  conversions,
                "CPA":          cpa,
                "CPO":          cpo,
                "CTR":          round(ctr, 2),
                "Conv Rate":    round(conv_rate, 2),
            })

    return pd.DataFrame(rows)


# ── Upload file parser ────────────────────────────────────────────────────────

def parse_upload_files(uploaded_files):
    """Parse one or more uploaded CSV or XLSX files.
    Maps their column names to the standard internal names and returns
    a combined DataFrame. Returns None if nothing could be parsed."""
    parsed_dfs = []

    for uploaded_file in uploaded_files:
        try:
            # Try Excel first, fall back to CSV
            if uploaded_file.name.endswith(".xlsx"):
                df = pd.read_excel(uploaded_file)
            else:
                df = pd.read_csv(uploaded_file)

            # Build a lowercase lookup so we can match columns case-insensitively
            lower_col_map = {col.lower().strip(): col for col in df.columns}

            # Map the file's actual columns to our standard names
            rename_dict = {}
            for standard_name, candidates in UPLOAD_COLUMN_CANDIDATES.items():
                for candidate in candidates:
                    if candidate in lower_col_map:
                        rename_dict[lower_col_map[candidate]] = standard_name
                        break

            df = df.rename(columns=rename_dict)

            # Parse the Date column to datetime if it exists
            if "Date" in df.columns:
                df["Date"] = pd.to_datetime(df["Date"], errors="coerce")

            parsed_dfs.append(df)

        except Exception as e:
            st.error(f"Could not parse {uploaded_file.name}: {e}")

    if not parsed_dfs:
        return None

    # Combine all files into one DataFrame
    combined = pd.concat(parsed_dfs, ignore_index=True)

    # If no Display Name column, derive it from the Channel column
    if "Display Name" not in combined.columns:
        if "Channel" in combined.columns:
            combined["Display Name"] = combined["Channel"]
        else:
            combined["Display Name"] = "Unknown"

    return combined


# ── API key helper ────────────────────────────────────────────────────────────

def get_api_key():
    """Read the Anthropic API key from Streamlit secrets or environment variable."""
    try:
        return st.secrets["ANTHROPIC_API_KEY"]
    except Exception:
        return os.environ.get("ANTHROPIC_API_KEY", "")


# ── PPTX builder ─────────────────────────────────────────────────────────────

def build_pptx(summary_dict, channel_df, ai_text=""):
    """Build a dark navy PowerPoint report with KPI summary, channel table,
    and optional AI insights. Returns a BytesIO buffer ready for download."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colour palette
    BG    = RGBColor(0x0D, 0x1B, 0x2A)  # Dark navy background
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREY  = RGBColor(0xA8, 0xB2, 0xBC)
    BLUE  = RGBColor(0x00, 0xA8, 0xE8)

    blank_layout = prs.slide_layouts[6]  # Completely blank layout

    def set_bg(slide):
        """Fill slide background with dark navy."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

    def add_textbox(slide, text, left, top, width, height,
                    font_size=14, bold=False, colour=WHITE, align=PP_ALIGN.LEFT):
        """Helper: add a styled text box to a slide."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = colour
        return txBox

    def add_footer(slide, slide_num):
        """Add 'Pacebird' footer label and slide number."""
        add_textbox(slide, "Pacebird", 0.3, 7.1, 3, 0.35,
                    font_size=9, colour=GREY)
        add_textbox(slide, str(slide_num), 12.7, 7.1, 0.5, 0.35,
                    font_size=9, colour=GREY, align=PP_ALIGN.RIGHT)

    # ── Slide 1: Title + KPI summary ─────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)

    add_textbox(slide1, "Cross-Channel Performance Dashboard",
                0.5, 0.4, 12, 0.7, font_size=28, bold=True, colour=WHITE)
    add_textbox(slide1, "Optus — May 2026",
                0.5, 1.1, 6, 0.4, font_size=16, colour=BLUE)

    # KPI values: 2 rows of metrics
    kpi_items = list(summary_dict.items())
    col_positions = [0.5, 3.0, 5.5, 8.0, 10.5]

    # Row 1: first 3 KPIs
    for i, (label, value) in enumerate(kpi_items[:3]):
        x = col_positions[i]
        add_textbox(slide1, label.upper(), x, 2.0, 2.3, 0.35,
                    font_size=9, colour=GREY)
        add_textbox(slide1, value, x, 2.35, 2.3, 0.55,
                    font_size=22, bold=True, colour=WHITE)

    # Row 2: remaining KPIs
    for i, (label, value) in enumerate(kpi_items[3:]):
        x = col_positions[i]
        add_textbox(slide1, label.upper(), x, 3.2, 2.3, 0.35,
                    font_size=9, colour=GREY)
        add_textbox(slide1, value, x, 3.55, 2.3, 0.55,
                    font_size=22, bold=True, colour=WHITE)

    add_footer(slide1, 1)

    # ── Slide 2: Channel Performance table ───────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)

    add_textbox(slide2, "Channel Performance",
                0.5, 0.3, 12, 0.55, font_size=22, bold=True, colour=WHITE)

    # Write one line per channel row
    y_pos = 1.1
    for _, row in channel_df.iterrows():
        channel_name = row.get("Channel", "Unknown")
        spend = row.get("Spend", 0)
        conversions = row.get("Conversions", 0)
        cpa = row.get("CPA", 0)
        line = (f"{channel_name}  |  Spend: A${spend:,.0f}  |  "
                f"Conversions: {conversions:,.0f}  |  CPA: A${cpa:.2f}")
        add_textbox(slide2, line, 0.5, y_pos, 12, 0.35,
                    font_size=11, colour=WHITE)
        y_pos += 0.42

    add_footer(slide2, 2)

    # ── Slide 3: AI insights (only if text provided) ──────────────────────────
    if ai_text:
        slide3 = prs.slides.add_slide(blank_layout)
        set_bg(slide3)

        add_textbox(slide3, "AI Performance Insights",
                    0.5, 0.3, 12, 0.55, font_size=22, bold=True, colour=WHITE)
        # Truncate to 1400 chars to fit comfortably on the slide
        add_textbox(slide3, ai_text[:1400], 0.5, 1.1, 12.3, 5.8,
                    font_size=11, colour=WHITE)

        add_footer(slide3, 3)

    # Save to an in-memory buffer and return it
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═════════════════════════════════════════════════════════════════════════════
# PAGE LAYOUT
# ═════════════════════════════════════════════════════════════════════════════

# Read data mode — set by wrapper pages (api=default, upload=file-upload mode)
_telco_data_mode = st.session_state.get('data_mode', 'api')
# Don't pop it — the mode needs to persist for the full page render

st.title("Cross-Channel Performance Dashboard")
st.markdown(
    "<p style='color:#6B7280; margin-top:-12px;'>Unified biddable media performance across "
    "SEM, programmatic, social and affiliate channels for Optus.</p>",
    unsafe_allow_html=True,
)

# Change 8: Show API banner in API mode, or show file uploader in upload mode
if _telco_data_mode == 'api':
    # API Data mode — show live data banner, then render mock data below
    st.markdown("""
    <div style="background:#EFF6FF;border-left:4px solid #2563EB;border-radius:8px;
    padding:10px 16px;margin-bottom:12px;">
        <strong>📡 Live API Data Mode</strong> — showing cross-channel data for Optus.
    </div>
    """, unsafe_allow_html=True)

if _telco_data_mode == 'upload':
    # File Upload mode — show uploader and stop before mock data renders
    uploaded_telco = st.file_uploader(
        "Upload cross-channel data",
        type=["csv", "xlsx", "xls", "tsv"],
        help="Upload a CSV/Excel file with channel performance data",
        key="telco_upload"
    )
    if uploaded_telco is None:
        st.markdown("""
        <div style="text-align:center;padding:60px 20px;background:#FFFFFF;
        border-radius:12px;box-shadow:0 4px 12px rgba(0,0,0,0.08);margin:20px 0;">
            <div style="font-size:48px;margin-bottom:16px;">📂</div>
            <div style="font-size:20px;font-weight:700;color:#111827;margin-bottom:8px;">
                No data loaded yet
            </div>
            <div style="font-size:15px;color:#6B7280;margin-bottom:4px;">
                Drag and drop cross-channel data to get started
            </div>
            <div style="font-size:13px;color:#9CA3AF;">
                Accepted formats: CSV, XLSX, XLS, TSV
            </div>
        </div>
        """, unsafe_allow_html=True)
        st.stop()
    else:
        # Process uploaded file
        try:
            import pandas as pd
            if uploaded_telco.name.endswith(".csv"):
                _telco_df = pd.read_csv(uploaded_telco)
            else:
                _telco_df = pd.read_excel(uploaded_telco)
            st.success(f"✅ Loaded {len(_telco_df):,} rows from {uploaded_telco.name}")
            # Show a preview table
            st.dataframe(_telco_df.head(20), use_container_width=True)
            st.info("Full cross-channel analysis will be available in a future update.")
        except Exception as _e:
            st.error(f"Could not read file: {_e}")
        st.stop()

# ── Upload expander — shown only in API mode for optional data override ────────
# (When in upload mode, we already handled everything above and called st.stop())
with st.expander("📁 Upload Channel Reports"):
    uploaded_files = st.file_uploader(
        "Upload channel export files (CSV or XLSX)",
        type=["csv", "xlsx"],
        accept_multiple_files=True,
        key="cc_upload",
    )
    if uploaded_files:
        parsed = parse_upload_files(uploaded_files)
        if parsed is not None:
            st.session_state["telco_channel_data"] = parsed
            n_ch   = parsed["Channel"].nunique()
            n_days = parsed["Date"].nunique() if "Date" in parsed.columns else 0
            st.success(
                f"✅ {len(uploaded_files)} file(s) loaded — {n_ch} channels, {n_days} days"
            )
        else:
            st.error("Could not parse the uploaded file(s). Check column names and format.")

# ── Data source: uploaded or mock ─────────────────────────────────────────────
if st.session_state.get("telco_channel_data") is not None:
    active_df = st.session_state["telco_channel_data"]
    n_ch   = active_df["Channel"].nunique()
    n_days = active_df["Date"].nunique() if "Date" in active_df.columns else 0
    st.success(f"✅ Using uploaded data — {n_ch} channels, {n_days} days")
else:
    active_df = generate_mock_data()
    st.warning("⚠️ Using mock data — upload reports to use live data")

# Determine which column to use for granular display names
dn_col = "Display Name" if "Display Name" in active_df.columns else "Channel"

# ── SECTION 1 — GLOBAL FILTERS ────────────────────────────────────────────────
st.subheader("Filters")

col_f1, col_f2, col_f3 = st.columns(3)

with col_f1:
    date_range = st.date_input(
        "Date range",
        value=[active_df["Date"].min(), active_df["Date"].max()],
        key="cc_dates",
    )

with col_f2:
    all_channels = sorted(active_df["Channel"].unique().tolist())
    selected_channels = st.multiselect(
        "Channel",
        options=all_channels,
        default=all_channels,
        key="cc_ch_filter",
    )

with col_f3:
    all_campaigns = sorted(active_df["Campaign"].unique().tolist())
    selected_campaigns = st.multiselect(
        "Campaign",
        options=all_campaigns,
        default=all_campaigns,
        key="cc_camp_filter",
    )

# Apply filters to the DataFrame
# Convert date_range inputs to Timestamps for comparison with the Date column
if isinstance(date_range, (list, tuple)) and len(date_range) == 2:
    start_date = pd.Timestamp(date_range[0])
    end_date   = pd.Timestamp(date_range[1])
else:
    start_date = active_df["Date"].min()
    end_date   = active_df["Date"].max()

filtered_df = active_df[
    (active_df["Date"] >= start_date) &
    (active_df["Date"] <= end_date) &
    (active_df["Channel"].isin(selected_channels)) &
    (active_df["Campaign"].isin(selected_campaigns))
].copy()

if filtered_df.empty:
    st.info("No data matches the selected filters. Adjust the date range or channel selection.")
    st.stop()

st.divider()

# ── SECTION 2 — SUMMARY METRIC CARDS ─────────────────────────────────────────
st.subheader("Summary Metrics")

# Compute summary totals — these are also reused in Section 8 (export)
total_spend    = filtered_df["Spend"].sum()
total_conv     = filtered_df["Conversions"].sum()
blended_cpa    = total_spend / max(total_conv, 1)
# Approximate orders as 82% of conversions for CPO calculation
blended_cpo    = total_spend / max(round(total_conv * 0.82), 1)
blended_ctr    = filtered_df["Clicks"].sum() / max(filtered_df["Impressions"].sum(), 1) * 100

col_m1, col_m2, col_m3, col_m4, col_m5 = st.columns(5)

with col_m1:
    st.metric("Total Spend", f"A${total_spend:,.0f}")
with col_m2:
    st.metric("Total Conversions", f"{int(total_conv):,}")
with col_m3:
    st.metric("Blended CPA", f"A${blended_cpa:.2f}")
with col_m4:
    st.metric("Blended CPO", f"A${blended_cpo:.2f}")
with col_m5:
    st.metric("Blended CTR", f"{blended_ctr:.2f}%")

st.divider()

# ── SECTION 3 — CHANNEL PERFORMANCE TABLE ────────────────────────────────────
st.subheader("Channel Performance")

# Aggregate filtered data by Channel
channel_agg = (
    filtered_df.groupby("Channel")
    .agg(
        Spend=("Spend", "sum"),
        Impressions=("Impressions", "sum"),
        Clicks=("Clicks", "sum"),
        Conversions=("Conversions", "sum"),
    )
    .reset_index()
)

# Compute derived metrics — clip(1) avoids divide-by-zero
channel_agg["CTR%"]         = channel_agg["Clicks"] / channel_agg["Impressions"].clip(1) * 100
channel_agg["CPA"]          = channel_agg["Spend"] / channel_agg["Conversions"].clip(1)
channel_agg["CPO"]          = channel_agg["CPA"] * 1.18
channel_agg["Conv Rate%"]   = channel_agg["Conversions"] / channel_agg["Clicks"].clip(1) * 100
channel_agg["% of Spend"]   = channel_agg["Spend"] / channel_agg["Spend"].sum() * 100

# Sort by highest spend first
channel_agg = channel_agg.sort_values("Spend", ascending=False).reset_index(drop=True)

# Colour-code the CPA column: green = efficient, amber = ok, red = expensive
def colour_cpa(val):
    if val < 30:
        return "color: #10B981; font-weight: bold"
    elif val <= 50:
        return "color: #F59E0B; font-weight: bold"
    else:
        return "color: #EF4444; font-weight: bold"

# Format the display table
display_agg = channel_agg.copy()
styled = (
    display_agg.style
    .map(colour_cpa, subset=["CPA"])
    .format({
        "Spend":       "A${:,.0f}",
        "Impressions": "{:,.0f}",
        "Clicks":      "{:,.0f}",
        "Conversions": "{:,.0f}",
        "CTR%":        "{:.2f}%",
        "CPA":         "A${:.2f}",
        "CPO":         "A${:.2f}",
        "Conv Rate%":  "{:.2f}%",
        "% of Spend":  "{:.1f}%",
    })
)

st.dataframe(styled, use_container_width=True, hide_index=True)

st.divider()

# ── SECTION 4 — SPEND BY CHANNEL ─────────────────────────────────────────────
st.subheader("Spend by Channel")

# Sort ascending so the largest bar appears at the top of the horizontal chart
spend_sorted = channel_agg.sort_values("Spend", ascending=True)

# Assign colours by channel index position (cycling through palette)
bar_colours = [
    CHART_COLOURS[i % len(CHART_COLOURS)]
    for i in range(len(spend_sorted))
]

spend_fig = go.Figure(go.Bar(
    x=spend_sorted["Spend"],
    y=spend_sorted["Channel"],
    orientation="h",
    marker_color=bar_colours,
    text=[f"A${v:,.0f}" for v in spend_sorted["Spend"]],
    textposition="outside",
))

spend_fig.update_layout(
    height=280,
    margin=dict(l=20, r=80, t=20, b=20),
    paper_bgcolor="rgba(0,0,0,0)",
    plot_bgcolor="rgba(0,0,0,0)",
    xaxis=dict(
        showgrid=True,
        gridcolor="#F3F4F6",
        tickprefix="A$",
        zeroline=False,
    ),
    yaxis=dict(showgrid=False),
)

st.plotly_chart(spend_fig, use_container_width=True)

st.divider()

# ── SECTION 5 — CPA/CPO OPTIMISATION TRACKER ─────────────────────────────────
st.subheader("CPA/CPO Optimisation Tracker")

# Initialise CPA targets in session state if not already set
if "telco_cpa_targets" not in st.session_state:
    st.session_state["telco_cpa_targets"] = DEFAULT_CPA_TARGETS.copy()

# Get the list of all granular display names for the trend chart
all_display_names = sorted(active_df[dn_col].unique().tolist())

col_left, col_right = st.columns([1, 1])

with col_left:
    metric_choice = st.radio(
        "Metric",
        ["CPA", "CPO", "CTR", "Conv Rate"],
        horizontal=True,
        key="cc_metric",
    )
    trend_channels_sel = st.multiselect(
        "Channels",
        options=all_display_names,
        default=all_display_names[:4],
        key="cc_trend_ch",
    )

with col_right:
    # Only show the CPA target editor when CPA metric is selected
    if metric_choice == "CPA":
        with st.expander("🎯 Edit CPA Targets"):
            new_targets = {}
            for dn in all_display_names:
                current_target = st.session_state["telco_cpa_targets"].get(dn, 40)
                new_val = st.number_input(
                    dn,
                    value=float(current_target),
                    step=1.0,
                    key=f"tgt_{dn}",
                )
                new_targets[dn] = new_val
            # Persist any edits back to session state
            st.session_state["telco_cpa_targets"].update(new_targets)

# Build daily trend data: group by date + display name
daily_trend = (
    filtered_df.groupby(["Date", dn_col])
    .agg(
        Spend=("Spend", "sum"),
        Conversions=("Conversions", "sum"),
        Clicks=("Clicks", "sum"),
        Impressions=("Impressions", "sum"),
        CTR=("CTR", "mean"),
        Conv_Rate=("Conv Rate", "mean"),
    )
    .reset_index()
)

# Recalculate metrics from raw numbers for accuracy
daily_trend["CPA"]       = daily_trend["Spend"] / daily_trend["Conversions"].clip(1)
daily_trend["CPO"]       = daily_trend["CPA"] * 1.18
daily_trend["CTR_calc"]  = daily_trend["Clicks"] / daily_trend["Impressions"].clip(1) * 100
daily_trend["Conv Rate"] = daily_trend["Conv_Rate"]

# Map UI metric choice to the DataFrame column name
metric_col_map = {
    "CPA":       "CPA",
    "CPO":       "CPO",
    "CTR":       "CTR_calc",
    "Conv Rate": "Conv Rate",
}
selected_metric_col = metric_col_map[metric_choice]

# Build the trend chart
trend_fig = go.Figure()
cpa_targets = st.session_state["telco_cpa_targets"]

for i, ch in enumerate(trend_channels_sel):
    color = CHART_COLOURS[i % len(CHART_COLOURS)]
    ch_data = daily_trend[daily_trend[dn_col] == ch].sort_values("Date")

    if ch_data.empty:
        continue

    # Main trend line with markers
    trend_fig.add_trace(go.Scatter(
        x=ch_data["Date"],
        y=ch_data[selected_metric_col],
        mode="lines+markers",
        name=ch,
        line=dict(color=color, width=2),
        marker=dict(size=5, color=color),
    ))

    # For CPA metric: add a dashed target line and flag days over target
    if metric_choice == "CPA" and ch in cpa_targets:
        target_val = cpa_targets[ch]

        # Dashed reference line at the target value
        trend_fig.add_trace(go.Scatter(
            x=ch_data["Date"],
            y=[target_val] * len(ch_data),
            mode="lines",
            name=f"{ch} target",
            line=dict(color=color, width=1, dash="dash"),
            showlegend=True,
        ))

        # Red circle markers on days where CPA exceeds the target
        over_target = ch_data[ch_data["CPA"] > target_val]
        if not over_target.empty:
            trend_fig.add_trace(go.Scatter(
                x=over_target["Date"],
                y=over_target["CPA"],
                mode="markers",
                marker=dict(color="#EF4444", size=9, symbol="circle"),
                showlegend=False,
                hovertemplate="⚠️ CPA: A$%{y:.2f}<extra></extra>",
            ))

trend_fig.update_layout(
    height=420,
    xaxis_title="Date",
    yaxis_title=metric_choice,
    paper_bgcolor="rgba(0,0,0,0)",
    plot_bgcolor="rgba(0,0,0,0)",
    xaxis=dict(showgrid=True, gridcolor="#F3F4F6"),
    yaxis=dict(showgrid=True, gridcolor="#F3F4F6"),
    legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="left", x=0),
    margin=dict(l=20, r=20, t=20, b=80),
)

st.plotly_chart(trend_fig, use_container_width=True)

st.divider()

# ── SECTION 6 — PERFORMANCE ALERTS ───────────────────────────────────────────
st.subheader("Performance Alerts")

# Use the full active dataset (not filtered by date) but respect channel/campaign filters
alert_df = active_df[
    (active_df["Channel"].isin(selected_channels)) &
    (active_df["Campaign"].isin(selected_campaigns))
].copy()

alerts = []

for dn in all_display_names:
    dn_data = alert_df[alert_df[dn_col] == dn].copy()
    if dn_data.empty:
        continue

    # Aggregate to daily level for this display channel
    daily_dn = (
        dn_data.groupby("Date")
        .agg(
            Spend=("Spend", "sum"),
            Conversions=("Conversions", "sum"),
            Clicks=("Clicks", "sum"),
            Impressions=("Impressions", "sum"),
        )
        .reset_index()
        .sort_values("Date")
    )

    # Recompute daily CPA and CTR from aggregated numbers
    daily_dn["CPA"] = daily_dn["Spend"] / daily_dn["Conversions"].clip(1)
    daily_dn["CTR"] = daily_dn["Clicks"] / daily_dn["Impressions"].clip(1) * 100

    target_cpa = cpa_targets.get(dn, 40)

    # Alert 1 — Latest day CPA is more than 15% above target
    if len(daily_dn) > 0:
        latest_cpa = daily_dn.iloc[-1]["CPA"]
        if latest_cpa > target_cpa * 1.15:
            alerts.append(
                f"**{dn}** — CPA A${latest_cpa:.2f} is "
                f"{((latest_cpa / target_cpa) - 1) * 100:.0f}% above target "
                f"(A${target_cpa:.0f})"
            )

    # Alert 2 — CTR declining for 3 consecutive days
    if len(daily_dn) >= 3:
        last3_ctr = daily_dn["CTR"].iloc[-3:].tolist()
        if last3_ctr[0] > last3_ctr[1] > last3_ctr[2]:
            alerts.append(
                f"**{dn}** — CTR declining for 3 consecutive days "
                f"({last3_ctr[0]:.2f}% → {last3_ctr[1]:.2f}% → {last3_ctr[2]:.2f}%)"
            )

    # Alert 3 — Recent 7-day spend pacing 10%+ above the daily average
    avg_daily_spend    = daily_dn["Spend"].mean()
    recent_avg_spend   = daily_dn["Spend"].iloc[-7:].mean() if len(daily_dn) >= 7 else avg_daily_spend
    if recent_avg_spend > avg_daily_spend * 1.10:
        pct_above = ((recent_avg_spend / avg_daily_spend) - 1) * 100
        alerts.append(
            f"**{dn}** — 7-day avg spend A${recent_avg_spend:,.0f}/day is "
            f"{pct_above:.0f}% above the campaign average "
            f"(A${avg_daily_spend:,.0f}/day)"
        )

# Render alerts or a clear-pass message
if alerts:
    alert_html = "<div style='background:#FFFBEB; border:1px solid #FCD34D; border-radius:10px; padding:16px;'>"
    for alert in alerts:
        alert_html += f"<p style='margin:4px 0;'>⚠️ {alert}</p>"
    alert_html += "</div>"
    st.markdown(alert_html, unsafe_allow_html=True)
else:
    st.success("✅ No active alerts — all channels within normal thresholds.")

st.divider()

# ── SECTION 7 — AI PERFORMANCE INSIGHTS ──────────────────────────────────────
st.subheader("AI Performance Insights")

api_key = get_api_key()

if not api_key:
    st.warning("No Anthropic API key found. Add ANTHROPIC_API_KEY to .streamlit/secrets.toml or as an environment variable.")
else:
    if st.button("Generate Performance Insights", type="primary", key="cc_ai_btn"):
        # Build a readable summary of channel performance to pass to the prompt
        channel_summary_lines = []
        for _, row in channel_agg.iterrows():
            channel_summary_lines.append(
                f"- {row['Channel']}: Spend A${row['Spend']:,.0f}, "
                f"Conversions {row['Conversions']:,.0f}, "
                f"CPA A${row['CPA']:.2f}, "
                f"CTR {row['CTR%']:.2f}%"
            )
        channel_summary_str = "\n".join(channel_summary_lines)

        # Build a targets summary string
        targets_str = "\n".join(
            [f"- {ch}: A${tgt:.0f}" for ch, tgt in cpa_targets.items()]
        )

        # Identify channels currently exceeding their CPA target
        flagged = []
        for _, row in channel_agg.iterrows():
            # Match channel_agg Channel to display name targets (best-effort)
            for dn, tgt in cpa_targets.items():
                if row["Channel"].lower() in dn.lower() or dn.lower() in row["Channel"].lower():
                    if row["CPA"] > tgt:
                        flagged.append(f"{dn} (CPA A${row['CPA']:.2f} vs target A${tgt:.0f})")
                    break
        flagged_str = ", ".join(flagged) if flagged else "None"

        prompt = f"""You are a senior programmatic media analyst reviewing cross-channel paid media performance for Optus Australia.

CHANNEL PERFORMANCE SUMMARY:
{channel_summary_str}

CPA TARGETS BY CHANNEL:
{targets_str}

CHANNELS EXCEEDING CPA TARGET:
{flagged_str}

OVERALL METRICS:
- Total Spend: A${total_spend:,.0f}
- Total Conversions: {int(total_conv):,}
- Blended CPA: A${blended_cpa:.2f}
- Blended CPO: A${blended_cpo:.2f}
- Blended CTR: {blended_ctr:.2f}%

Write a concise performance insights report (4–6 paragraphs) covering:
1. Overall cross-channel performance assessment
2. Top performing channels and why
3. Underperforming channels with specific concerns
4. Budget reallocation recommendations with channel-specific reasoning
5. Actionable next steps for the media team

Use specific metrics and A$ values throughout. Be direct and professional."""

        with st.spinner("Generating insights via Claude..."):
            client = anthropic.Anthropic(api_key=api_key)
            message = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=800,
                messages=[{"role": "user", "content": prompt}],
            )
            st.session_state["cc_ai_insights"] = message.content[0].text

    # Display saved insights in a white card
    if st.session_state.get("cc_ai_insights"):
        insights_html = (
            "<div style='background:#FFFFFF; border-radius:12px; padding:20px; "
            "box-shadow:0 4px 12px rgba(0,0,0,0.08);'>"
            + st.session_state["cc_ai_insights"].replace("\n", "<br>")
            + "</div>"
        )
        st.markdown(insights_html, unsafe_allow_html=True)

st.divider()

# ── SECTION 8 — EXPORT ────────────────────────────────────────────────────────
st.subheader("Export")

if st.button("Download Performance Report (.pptx)", key="cc_pptx_btn"):
    # Build summary dict reusing totals computed in Section 2
    summary = {
        "Total Spend":       f"A${total_spend:,.0f}",
        "Total Conversions": f"{int(total_conv):,}",
        "Blended CPA":       f"A${blended_cpa:.2f}",
        "Blended CPO":       f"A${blended_cpo:.2f}",
        "Blended CTR":       f"{blended_ctr:.2f}%",
    }
    pptx_buf = build_pptx(
        summary,
        channel_agg,
        st.session_state.get("cc_ai_insights", ""),
    )
    st.session_state["cc_pptx"] = pptx_buf

if st.session_state.get("cc_pptx"):
    st.download_button(
        label="Save .pptx",
        data=st.session_state["cc_pptx"],
        file_name=f"{date.today()}_optus_cross_channel.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

print("Cross-Channel Performance Dashboard loaded.")

# ── Action Log ────────────────────────────────────────────────────────────────
st.markdown("---")

def _load_action_log_telco():
    log_path = "action_log.json"
    if os.path.exists(log_path):
        try:
            with open(log_path, "r") as f:
                return json.load(f)
        except Exception:
            return []
    return []

def _save_action_log_telco(log_data):
    with open("action_log.json", "w") as f:
        json.dump(log_data, f, indent=2)

with st.expander("📋 Action Log", expanded=False):
    _action_log = _load_action_log_telco()

    # Filter to just this page's actions
    _page_name = "Telco Cross-Channel"
    _page_log = [e for e in _action_log if e.get("page", "") == _page_name]

    if not _page_log:
        st.info("No actions logged yet. Push actions will appear here.")
    else:
        st.markdown(f"**{len(_page_log)} action(s) logged for this page**")
        for _i, _entry in enumerate(reversed(_page_log)):
            _global_idx = _action_log.index(_entry)
            _log_c1, _log_c2, _log_c3, _log_c4 = st.columns([2, 2, 2, 3])
            with _log_c1:
                st.markdown(f"**{_entry.get('timestamp', '')[:16]}**")
            with _log_c2:
                st.markdown(_entry.get('action', ''))
            with _log_c3:
                st.markdown(f"✅ {_entry.get('simulated_response', '')}")
            with _log_c4:
                _current_note = _entry.get("user_note", "")
                _new_note = st.text_input(
                    "Add note",
                    value=_current_note,
                    key=f"telco_log_note_{_global_idx}",
                    placeholder="What actually happened...",
                    label_visibility="collapsed"
                )
                if _new_note != _current_note:
                    _full_log = _load_action_log_telco()
                    _full_log[_global_idx]["user_note"] = _new_note
                    _save_action_log_telco(_full_log)
                    st.rerun()
            st.markdown("---")
