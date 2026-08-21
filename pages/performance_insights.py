import io
import json
import os
import re
import random
from datetime import date
import sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import streamlit as st
import pandas as pd
from utils.design_system import (
    metric_card, apply_plotly_style, filter_widget, PLOTLY_CONFIG,
    PRIMARY, SECONDARY, SUCCESS, WARNING, DANGER, WHITE, TEXT_PRI, TEXT_SEC,
)
from utils.mock_data import generate_api_mock_data
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY    = "#14113b"
CYAN    = "#00b2a9"
GREEN   = "#34b233"
PURPLE  = "#6e2ca9"
MAGENTA = "#c70099"

# STYLE LOCK: Pacebird design system — primary #F5A623 orange, secondary #1B2A4A navy, font Poppins.
# Shared CSS is applied globally by app.py via utils/design_system.get_css().
# STYLE LOCK

# ── Source mode detection ──────────────────────────────────────────────────────
# _pi_source_mode is set by ongoing_performance.py before delegating here.
# 'ongoing' = API Data section, 'one-off' = Upload Report section.
_pi_source_mode = st.session_state.get('_pi_source_mode', 'one-off')
st.session_state.pop('_pi_source_mode', None)  # reset after reading


# ── Column name normalisation maps ────────────────────────────────────────────
# COLUMN_MAP: internal dimension name → list of possible source column names.
# For each internal name, the first matching source column found in the data is
# used; already-claimed columns are skipped so priority is top-to-bottom within
# each internal name. Mixed-case here — matching is done after lowercasing.
COLUMN_MAP = {
    "advertiser":      ["Advertiser", "Partner", "Brand"],
    "campaign":        ["Campaign", "Campaign Name", "Brand Name", "Order"],
    "insertion_order": ["Insertion Order", "Campaign", "Campaign Name", "Order"],
    "line_item":       ["Line Item", "Line Item Name", "Ad Group", "Ad Group Name"],
    "creative":        ["Creative", "Ad"],
}

# METRIC_MAP: lowercased source column name → internal metric/dimension name.
# Covers numeric, date, device, and environment columns across all DSPs.
METRIC_MAP = {
    # Device type — covers DV360, TTD, and generic DSP column names
    "device type":                 "device_type",
    "device":                      "device_type",
    "device_type":                 "device_type",
    "device category":             "device_type",
    "device make":                 "device_type",
    # Environment / inventory — covers DV360, TTD, Amazon, and generic names
    "environment":                 "environment",
    "environment type":            "environment",
    "inventory type":              "environment",
    "supply type":                 "environment",
    "site type":                   "environment",
    "app/url":                     "environment",
    "platform":                    "environment",
    "device environment":          "environment",
    # Buy type — covers DV360, TTD, Amazon, and generic DSP deal-type columns
    "buy type":                    "buy_type",
    "deal type":                   "buy_type",
    "inventory source":            "buy_type",
    "auction type":                "buy_type",
    # Date
    "date":                        "date",
    "day":                         "date",
    "week":                        "date",
    "month":                       "date",
    # Impressions
    "impressions":                 "impressions",
    "impression":                  "impressions",
    "served impressions":          "impressions",
    "total impressions":           "impressions",
    # Clicks
    "clicks":                      "clicks",
    "click":                       "clicks",
    "total clicks":                "clicks",
    "link clicks":                 "clicks",
    # Spend / cost — USD and AUD variants all normalise to spend_usd
    "spend":                       "spend_usd",
    "spend (usd)":                 "spend_usd",
    "spend (aud)":                 "spend_usd",
    "total spend":                 "spend_usd",
    "total spend (usd)":           "spend_usd",
    "total spend (aud)":           "spend_usd",
    "media cost":                  "spend_usd",
    "media cost (usd)":            "spend_usd",
    "revenue (usd)":               "spend_usd",
    "revenue (aud)":               "spend_usd",
    "cost":                        "spend_usd",
    "billed spend":                "spend_usd",
    # Conversions
    "conversions":                 "conversions",
    "total conversions":           "conversions",
    "post-click conversions":      "conversions",
    # CTR (pre-calculated in export)
    "ctr":                         "ctr_raw",
    "click-through rate":          "ctr_raw",
    # CPM (pre-calculated in export) — currency-suffixed variants also accepted
    "cpm":                         "cpm_raw",
    "cpm (usd)":                   "cpm_raw",
    "cpm (aud)":                   "cpm_raw",
    "avg. cpm":                    "cpm_raw",
    "average cpm":                 "cpm_raw",
}

# ── DSP source detector ───────────────────────────────────────────────────────
# Identifies which DSP the CSV came from by checking for DSP-specific column names.
# More specific signals are checked first to avoid false positives.
def detect_source(columns):
    cols_lower = [c.strip().lower() for c in columns]
    # Amazon: "Detail Page Views (DPV)" is unique to Amazon DSP reports
    if any("detail page views" in c for c in cols_lower):
        return "Amazon"
    # DV360 YouTube: "Ad Format" column only appears in YouTube reports
    if any("ad format" in c for c in cols_lower):
        return "DV360 YouTube"
    # TTD: "Supply Vendor" is a TTD-specific column
    if any("supply vendor" in c for c in cols_lower):
        return "TTD"
    # DV360 Display: "Insertion Order" is a DV360-specific hierarchy level
    if any("insertion order" in c for c in cols_lower):
        return "DV360"
    # TTD fallback: uses "Ad Group" without an Insertion Order level
    if any("ad group" in c for c in cols_lower):
        return "TTD"
    # Amazon fallback: standalone "order" column (exact match)
    if any(c == "order" for c in cols_lower):
        return "Amazon"
    # DV360 / TTD fallback based on spend column naming
    if any("revenue (usd)" in c for c in cols_lower):
        return "DV360"
    if any("billed spend" in c for c in cols_lower):
        return "TTD"
    return "Generic"

# ── Brand memory helpers ──────────────────────────────────────────────────────
# brand_memory.json stores per-brand context that shapes AI commentary.
BRAND_MEMORY_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                 "brand_memory.json")

def load_brand_memory():
    """Load brand memory from JSON file. Returns empty dict if file doesn't exist."""
    if os.path.exists(BRAND_MEMORY_PATH):
        with open(BRAND_MEMORY_PATH, "r") as f:
            return json.load(f)
    return {}

def save_brand_memory(memory):
    """Save brand memory dict back to JSON file."""
    with open(BRAND_MEMORY_PATH, "w") as f:
        json.dump(memory, f, indent=2)


# ── File loader and normaliser ────────────────────────────────────────────────
def load_and_normalise(uploaded_file):
    """
    Reads a CSV, TSV, or Excel file, normalises column names to our standard
    format, and adds a 'source_file' and 'dsp_source' column for traceability.
    """
    name = uploaded_file.name.lower()
    if name.endswith(".xlsx") or name.endswith(".xls"):
        df = pd.read_excel(uploaded_file)
    elif name.endswith(".tsv"):
        df = pd.read_csv(uploaded_file, sep="\t")
    else:
        df = pd.read_csv(uploaded_file)

    # Detect DSP before renaming columns
    dsp = detect_source(df.columns.tolist())

    # Normalise column names: lowercase + strip whitespace
    df.columns = [c.strip().lower() for c in df.columns]

    # Apply dimension column mapping (COLUMN_MAP).
    # For each internal name, pick the first source column that exists AND hasn't
    # already been claimed by a higher-priority internal name.
    _used_src = set()
    _dim_rename = {}
    for _internal, _src_list in COLUMN_MAP.items():
        for _src in _src_list:
            _src_lower = _src.lower()
            if _src_lower in df.columns and _src_lower not in _used_src:
                _dim_rename[_src_lower] = _internal
                _used_src.add(_src_lower)
                break

    # Apply metric/date column mapping (METRIC_MAP) for all remaining columns
    _metric_rename = {k: v for k, v in METRIC_MAP.items()
                      if k in df.columns and k not in _used_src}

    # Merge both rename dicts and apply in one pass
    df = df.rename(columns={**_dim_rename, **_metric_rename})

    # Multiple source columns can map to the same standard name — keep first occurrence
    df = df.loc[:, ~df.columns.duplicated(keep="first")]

    # Convert numeric columns — DSP exports often include commas or $ signs
    for col in ["impressions", "clicks", "spend_usd", "conversions"]:
        if col in df.columns:
            df[col] = (
                df[col].astype(str)
                .str.replace(r"[$,]", "", regex=True)   # strip $ and commas
                .str.strip()
            )
            df[col] = pd.to_numeric(df[col], errors="coerce")

    # Calculate CTR and CPM from raw numbers (more reliable than DSP-provided values)
    if "impressions" in df.columns and "clicks" in df.columns:
        df["ctr"] = df["clicks"] / df["impressions"]
    if "impressions" in df.columns and "spend_usd" in df.columns:
        df["cpm"] = df["spend_usd"] / df["impressions"] * 1000

    # Tag each row with where it came from
    df["source_file"] = uploaded_file.name
    df["dsp_source"]  = dsp

    return df, dsp


def detect_grouping_column(df):
    """
    Find the best column to group AI insights by.
    Checks in priority order and returns (column_name, display_label).
    Falls back to the first non-numeric, non-date column if none match.
    """
    # Priority list: (normalised column name after COLUMN_MAP, display label)
    priority = [
        ("advertiser",      "Advertiser"),
        ("campaign",        "Campaign"),
        ("insertion_order", "Insertion Order"),
        ("line_item",       "Line Item"),
        ("creative",        "Creative"),
    ]
    for col, label in priority:
        if col in df.columns:
            return col, label

    # Fall back to the first column that looks like a text/category dimension
    _skip = {"source_file", "dsp_source", "date", "ctr", "cpm",
             "advertiser", "campaign", "insertion_order", "line_item", "creative"}
    for col in df.columns:
        if col in _skip:
            continue
        if df[col].dtype == object:
            return col, col.replace("_", " ").title()

    return None, None



# ── PowerPoint export — dark premium template ────────────────────────────────
# Built entirely from a blank Presentation (no template file required).
#
# Colour scheme:
#   Background  #0D1B2A  dark navy
#   Primary     #FFFFFF  white text
#   Secondary   #A8B2BC  light grey labels / footer
#   Accent      #00A8E8  electric blue
#   Box fill    #1A2D40  slightly lighter navy for cards and table rows
#
# Font: Calibri throughout.  Slide size: 13.33" × 7.5" widescreen.
#
# Slide structure:
#   1.      Executive Summary — 3 KPI cards + Revenue chart + AI best/worst insight
#   2 + 3.  Per brand (repeated): Performance Breakdown + Recommendations
#   Last.   Budget Shift Recommendations (table)

_PPT_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy background
_PPT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)    # primary text
_PPT_GREY  = RGBColor(0xA8, 0xB2, 0xBC)   # secondary / label text
_PPT_BLUE  = RGBColor(0x00, 0xA8, 0xE8)   # electric blue accent
_PPT_DIM   = RGBColor(0x1A, 0x2D, 0x40)   # lighter navy for boxes / table rows

_PW = 13.33   # slide width in inches
_PH = 7.5     # slide height in inches

def _new_blank_slide(prs):
    """Add a blank slide and fill it with the dark navy background colour."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill  = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _PPT_BG
    return slide


def _tb(slide, x, y, w, h, text, size, bold=False, color=None,
        align=PP_ALIGN.LEFT):
    """Add a single-run Calibri text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.name      = "Calibri"
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color or _PPT_WHITE
    return txb

def _box(slide, x, y, w, h, fill_color, border_color=None, border_pt=1.0):
    """Add a filled rectangle with an optional coloured border."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width     = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def _kpi_box(slide, x, y, w, h, label, value):
    """KPI card: dark box with electric blue border, small label + large bold value."""
    _box(slide, x, y, w, h, _PPT_DIM, border_color=_PPT_BLUE, border_pt=1.5)
    _tb(slide, x + 0.1, y + 0.1,  w - 0.2, 0.28,
        label.upper(), 8, color=_PPT_GREY, align=PP_ALIGN.CENTER)
    _tb(slide, x + 0.05, y + 0.4, w - 0.1, 0.55,
        value, 22, bold=True, align=PP_ALIGN.CENTER)


def _slide_footer(slide, slide_num):
    """'Insights App' logo text bottom-left and slide number bottom-right."""
    _tb(slide, 0.3,      _PH - 0.34, 2.5, 0.28, "Insights App", 8, color=_PPT_GREY)
    _tb(slide, _PW - 1.6, _PH - 0.34, 1.3, 0.28, str(slide_num),
        8, color=_PPT_GREY, align=PP_ALIGN.RIGHT)


def _add_bullets(slide, x, y, w, h, bullets, size=11):
    """
    Render a list of bullet points into a text box.
    Each bullet: electric blue dot (●) followed by white Calibri text.
    """
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True

    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)

        r_dot            = p.add_run()
        r_dot.text       = "●  "
        r_dot.font.name  = "Calibri"
        r_dot.font.size  = Pt(size - 1)
        r_dot.font.color.rgb = _PPT_BLUE

        r_text            = p.add_run()
        r_text.text       = text
        r_text.font.name  = "Calibri"
        r_text.font.size  = Pt(size)
        r_text.font.color.rgb = _PPT_WHITE


def _make_revenue_chart_pptx(camp_summary):
    """
    Build a Revenue by Brand bar chart with a dark background and electric blue bars.
    Returns a BytesIO PNG buffer, or None if spend data is unavailable.
    """
    if camp_summary.empty or "spend_usd" not in camp_summary.columns:
        return None

    # Find whichever grouping column was used to build camp_summary
    _c = next((c for c in ("advertiser", "campaign", "insertion_order", "line_item")
               if c in camp_summary.columns), None)
    if _c is None:
        return None

    data = (camp_summary[[_c, "spend_usd"]]
            .dropna()
            .sort_values("spend_usd", ascending=False))

    fig, ax = plt.subplots(figsize=(7.5, 3.6))
    fig.patch.set_facecolor("#0D1B2A")
    ax.set_facecolor("#0D1B2A")

    bars = ax.bar(data[_c], data["spend_usd"],
                  color="#00A8E8", edgecolor="none", width=0.5)
    ax.yaxis.set_major_formatter(
        mticker.FuncFormatter(
            lambda v, _: f"A${v/1e6:.1f}M" if v >= 1e6 else f"A${v/1e3:.0f}K"
        )
    )
    ax.bar_label(bars,
                 labels=[f"A${v/1e6:.2f}M" if v >= 1e6 else f"A${v/1e3:.0f}K"
                         for v in data["spend_usd"]],
                 padding=4, fontsize=9, color="white", fontweight="bold")
    ax.set_ylim(0, data["spend_usd"].max() * 1.32)
    ax.tick_params(colors="white", labelsize=9)
    ax.set_ylabel("Revenue (USD)", fontsize=9, color="#A8B2BC")
    for sp in ax.spines.values():
        sp.set_color("#2A3D4E")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", linewidth=0.4, alpha=0.4, color="#2A3D4E")
    ax.set_axisbelow(True)
    plt.xticks(rotation=20, ha="right", fontsize=9, color="white")
    plt.tight_layout(pad=0.4)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="#0D1B2A")
    buf.seek(0)
    plt.close(fig)
    return buf


def _build_data_text(df_all, camp_summary):
    """
    Produce a compact text summary of all campaign data for use in AI prompts.
    Includes brand totals and a line item breakdown when available.
    """
    # Find whichever top-level grouping column was used to build camp_summary
    _c = next((c for c in ("advertiser", "campaign", "insertion_order", "line_item")
               if c in camp_summary.columns), None)

    lines = ["=== Brand Totals ==="]
    for _, r in camp_summary.iterrows():
        brand_val = r[_c] if _c else "Unknown"
        parts = [f"Brand: {brand_val}"]
        if "impressions" in r and pd.notna(r.get("impressions")):
            parts.append(f"Impressions: {int(r['impressions']):,}")
        if "clicks" in r and pd.notna(r.get("clicks")):
            parts.append(f"Clicks: {int(r['clicks']):,}")
        if "spend_usd" in r and pd.notna(r.get("spend_usd")):
            parts.append(f"Revenue: ${r['spend_usd']:,.0f}")
        if "ctr" in r and pd.notna(r.get("ctr")):
            parts.append(f"CTR: {r['ctr']:.2%}")
        if "cpm" in r and pd.notna(r.get("cpm")):
            parts.append(f"CPM: ${r['cpm']:,.2f}")
        lines.append(" | ".join(parts))

    # Line item breakdown gives the AI enough detail to cite best/worst performers
    _c_all = next((c for c in ("advertiser", "campaign", "insertion_order")
                   if c in df_all.columns), None)
    if "line_item" in df_all.columns and _c_all:
        lines.append("\n=== By Brand and Line Item ===")
        agg = {c: (c, "sum") for c in ["impressions", "clicks", "spend_usd"]
               if c in df_all.columns}
        grp = df_all.groupby([_c_all, "line_item"]).agg(**agg).reset_index()
        if "impressions" in grp.columns and "spend_usd" in grp.columns:
            grp["cpm"] = grp["spend_usd"] / grp["impressions"].clip(lower=1) * 1000
        if "impressions" in grp.columns and "clicks" in grp.columns:
            grp["ctr"] = grp["clicks"] / grp["impressions"].clip(lower=1)
        if "clicks" in grp.columns and "spend_usd" in grp.columns:
            grp["cpc"] = grp["spend_usd"] / grp["clicks"].clip(lower=1)
        for _, r in grp.iterrows():
            parts = [f"Brand: {r[_c_all]} | Line Item: {r['line_item']}"]
            if "impressions" in r and pd.notna(r.get("impressions")):
                parts.append(f"Impressions: {int(r['impressions']):,}")
            if "spend_usd" in r and pd.notna(r.get("spend_usd")):
                parts.append(f"Revenue: ${r['spend_usd']:,.0f}")
            if "cpm" in r and pd.notna(r.get("cpm")):
                parts.append(f"CPM: ${r['cpm']:,.2f}")
            if "ctr" in r and pd.notna(r.get("ctr")):
                parts.append(f"CTR: {r['ctr']:.2%}")
            if "cpc" in r and pd.notna(r.get("cpc")):
                parts.append(f"CPC: ${r['cpc']:,.2f}")
            lines.append(" | ".join(parts))

    return "\n".join(lines)


def _ai_pptx(api_key, prompt, max_tokens=400):
    """Synchronous Anthropic API call for slide content. Returns the response text."""
    client = anthropic.Anthropic(api_key=api_key)
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=max_tokens,
        system=(
            "You are a senior programmatic advertising analyst. "
            "Write clear, concise, data-driven commentary for ad tech professionals. "
            "Be specific — always reference actual numbers from the data. "
            "Follow formatting instructions exactly — no deviations."
        ),
        messages=[{"role": "user", "content": prompt}],
    )
    return msg.content[0].text.strip()


def _parse_io_sections(text):
    """
    Parse AI breakdown text into [(section_title, [bullet_strings])] tuples.
    Looks for DISPLAY:, VIDEO:, YOUTUBE: section headers.
    """
    sections      = []
    current_title = None
    current_items = []

    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        upper = s.upper()
        if upper.startswith("DISPLAY"):
            if current_title:
                sections.append((current_title, current_items))
            current_title, current_items = "Display", []
        elif upper.startswith("YOUTUBE"):
            if current_title:
                sections.append((current_title, current_items))
            current_title, current_items = "YouTube", []
        elif upper.startswith("VIDEO"):
            if current_title:
                sections.append((current_title, current_items))
            current_title, current_items = "Video", []
        elif s and (s[0] in "•-–—●" or (len(s) > 1 and s[0].isdigit() and s[1] in ".)")):
            bullet = s.lstrip("•-–—●0123456789.) ").strip()
            if bullet:
                current_items.append(bullet)

    if current_title:
        sections.append((current_title, current_items))

    return sections


def _add_budget_table(slide, budget_text, x, y, w, h):
    """
    Parse the AI budget text and render a styled table onto the slide.
    Expected format: pipe-separated — Brand | Best IO | Recommendation
    """
    rows = []
    for line in budget_text.split("\n"):
        parts = [p.strip() for p in line.split("|")]
        if len(parts) >= 3 and parts[0]:
            rows.append(parts[:3])

    if not rows:
        _tb(slide, x, y + 0.2, w, 0.4, "No budget data available.", 11, color=_PPT_GREY)
        return

    headers   = ["Brand", "Best IO", "Recommendation"]
    gap       = 0.06
    col_ws    = [w * 0.22 - gap, w * 0.18 - gap, w * 0.57 - gap]
    row_h     = 0.5

    # Header row
    cx = x
    for i, hdr in enumerate(headers):
        _box(slide, cx, y, col_ws[i], row_h * 0.75, _PPT_BLUE)
        _tb(slide, cx + 0.08, y + 0.1, col_ws[i] - 0.16, 0.34,
            hdr, 10, bold=True)
        cx += col_ws[i] + gap

    # Data rows
    for r_i, row in enumerate(rows):
        ry = y + row_h * 0.8 + r_i * row_h
        if ry + row_h > y + h:
            break
        row_fill = _PPT_DIM if r_i % 2 == 0 else _PPT_BG
        cx = x
        for c_i, cell in enumerate(row):
            _box(slide, cx, ry, col_ws[c_i], row_h - 0.05, row_fill)
            # Highlight positive recommendations in electric blue
            txt_color = _PPT_BLUE if any(
                kw in cell.lower() for kw in ["increase", "scale", "boost", "shift to"]
            ) else _PPT_WHITE
            _tb(slide, cx + 0.08, ry + 0.08, col_ws[c_i] - 0.16, row_h - 0.12,
                cell, 9, color=txt_color)
            cx += col_ws[c_i] + gap


def _make_daily_chart_pptx(df_all):
    """
    Build a daily spend or impressions line chart for the PPTX daily performance slide.
    Returns a BytesIO PNG buffer, or None if required data is unavailable.
    """
    if "date" not in df_all.columns:
        return None
    _metric = ("spend_usd"   if "spend_usd"   in df_all.columns else
               "impressions" if "impressions" in df_all.columns else None)
    if _metric is None:
        return None

    daily = (df_all.groupby("date")[_metric].sum().reset_index().sort_values("date"))

    fig, ax = plt.subplots(figsize=(11, 3.6))
    fig.patch.set_facecolor("#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    ax.plot(daily["date"], daily[_metric],
            color="#00A8E8", linewidth=2, marker="o", markersize=4)
    ax.fill_between(daily["date"], daily[_metric], alpha=0.15, color="#00A8E8")
    label = "Spend (USD)" if _metric == "spend_usd" else "Impressions"
    if _metric == "spend_usd":
        ax.yaxis.set_major_formatter(
            mticker.FuncFormatter(
                lambda v, _: f"A${v/1e6:.1f}M" if v >= 1e6 else f"A${v/1e3:.0f}K"
            )
        )
    ax.set_ylabel(label, fontsize=9, color="#A8B2BC")
    ax.tick_params(colors="white", labelsize=9)
    for sp in ax.spines.values():
        sp.set_color("#2A3D4E")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", linewidth=0.4, alpha=0.4, color="#2A3D4E")
    ax.set_axisbelow(True)
    plt.xticks(rotation=20, ha="right", fontsize=9, color="white")
    plt.tight_layout(pad=0.4)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="#0D1B2A")
    buf.seek(0)
    plt.close(fig)
    return buf


def build_pptx_report(api_key, camp_summary, df_all, sections=None,
                      ai_analysis_cache=None):
    """
    Build and return the full PowerPoint report as a BytesIO buffer.

    sections dict controls which slides are included:
        summary_metrics        — Slide 1 KPI cards
        performance_charts     — Slide 1 revenue chart
        daily_performance      — Daily trend chart slide
        period_comparison      — Period-over-period comparison slide
        ai_advertiser_insights — Per-brand breakdown + recommendations + budget shift slides
        ai_driven_analysis     — AI-driven analysis summary slide

    Passing sections=None (default) includes all slide types.
    ai_analysis_cache: the ai_analysis_result dict from session state (reused if present).
    """
    if sections is None:
        sections = {k: True for k in (
            "summary_metrics", "performance_charts", "daily_performance",
            "period_comparison", "ai_advertiser_insights", "ai_driven_analysis",
        )}

    prs              = Presentation()
    prs.slide_width  = Inches(_PW)
    prs.slide_height = Inches(_PH)

    slide_num = 1
    data_text = _build_data_text(df_all, camp_summary)
    # Find whichever top-level grouping column was used to build camp_summary
    _camp_col = next((c for c in ("advertiser", "campaign", "insertion_order", "line_item")
                      if c in camp_summary.columns), None)
    campaigns = camp_summary[_camp_col].tolist() if _camp_col else []

    # ── Slide 1: Executive Summary — KPIs + Revenue chart ────────────────────
    if sections.get("summary_metrics") or sections.get("performance_charts"):
        s1 = _new_blank_slide(prs)
        _tb(s1, 0.4, 0.17, 12.5, 0.75, "Campaign Performance Summary", 28, bold=True)
        _box(s1, 0.4, 0.95, 12.5, 0.04, _PPT_BLUE)

        # KPI cards — only when summary_metrics is checked
        if sections.get("summary_metrics"):
            total_impr = (camp_summary["impressions"].sum()
                          if "impressions" in camp_summary.columns else 0)
            total_rev  = (camp_summary["spend_usd"].sum()
                          if "spend_usd" in camp_summary.columns else 0)
            kpi_data = [
                ("Total Impressions", f"{total_impr:,.0f}"),
                ("Total Revenue",     f"A${total_rev:,.0f}"),
                ("Number of Brands",  str(len(campaigns))),
            ]
            kpi_w, kpi_gap = 3.8, 0.35
            for i, (lbl, val) in enumerate(kpi_data):
                _kpi_box(s1, 0.4 + i * (kpi_w + kpi_gap), 1.08, kpi_w, 1.05, lbl, val)

        # Revenue chart — only when performance_charts is checked
        if sections.get("performance_charts"):
            chart_buf = _make_revenue_chart_pptx(camp_summary)
            if chart_buf:
                s1.shapes.add_picture(chart_buf,
                                      Inches(0.4), Inches(2.32), Inches(7.8), Inches(3.95))

        # AI best/worst insight — always shown on the executive summary slide
        exec_text = _ai_pptx(
            api_key,
            f"""Campaign performance data:
{data_text}

Write two short insight paragraphs for an executive summary slide:
1. Start "Best performer: [Brand] —" then 1-2 sentences on WHY, citing key metrics (CPM, CTR, or CPV).
2. Start "Worst performer: [Brand] —" then 1-2 sentences on WHY, citing key metrics.

Plain paragraphs only — no headings, no bullets, no extra text.""",
            max_tokens=200,
        )
        _tb(s1, 8.5, 2.32, 4.5, 0.28, "KEY INSIGHTS", 8, bold=True, color=_PPT_BLUE)
        txb = s1.shapes.add_textbox(Inches(8.5), Inches(2.7), Inches(4.45), Inches(3.6))
        tf  = txb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text           = exec_text
        run.font.name      = "Calibri"
        run.font.size      = Pt(10.5)
        run.font.color.rgb = _PPT_GREY

        _slide_footer(s1, slide_num)
        slide_num += 1

    # ── Daily Performance slide ───────────────────────────────────────────────
    if sections.get("daily_performance") and "date" in df_all.columns:
        s_daily = _new_blank_slide(prs)
        _tb(s_daily, 0.4, 0.17, 12.5, 0.65, "Daily Performance", 20, bold=True)
        _box(s_daily, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)
        daily_buf = _make_daily_chart_pptx(df_all)
        if daily_buf:
            s_daily.shapes.add_picture(daily_buf,
                                       Inches(0.4), Inches(1.05),
                                       Inches(12.5), Inches(5.5))
        _slide_footer(s_daily, slide_num)
        slide_num += 1

    # ── Period Comparison slide ───────────────────────────────────────────────
    if sections.get("period_comparison") and "date" in df_all.columns:
        try:
            _dp = pd.to_datetime(df_all["date"], errors="coerce")
            _dr = (_dp.max() - _dp.min()).days
            if _dr >= 14:
                s_pc = _new_blank_slide(prs)
                _tb(s_pc, 0.4, 0.17, 12.5, 0.65,
                    "Period Comparison — Last 30 Days vs Prior 30 Days", 20, bold=True)
                _box(s_pc, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)

                today    = _dp.max()
                curr_df  = df_all[_dp >= today - pd.Timedelta(days=29)]
                prior_df = df_all[(_dp >= today - pd.Timedelta(days=59)) &
                                  (_dp <  today - pd.Timedelta(days=29))]

                def _safe_sum(d, col):
                    return d[col].sum() if col in d.columns and len(d) > 0 else 0

                y_pos = 1.1
                for m_col, m_lbl, fmt in [
                    ("spend_usd",   "Spend",       "A${:,.0f}"),
                    ("impressions", "Impressions",  "{:,.0f}"),
                    ("clicks",      "Clicks",       "{:,.0f}"),
                ]:
                    c_v = _safe_sum(curr_df, m_col)
                    p_v = _safe_sum(prior_df, m_col)
                    chg = f"{(c_v - p_v) / p_v * 100:+.1f}%" if p_v > 0 else "N/A"
                    row_txt = (f"{m_lbl}:   Current {fmt.format(c_v)}"
                               f"   |   Prior {fmt.format(p_v)}"
                               f"   |   Change {chg}")
                    _tb(s_pc, 0.5, y_pos, 12.0, 0.5, row_txt, 12, color=_PPT_WHITE)
                    y_pos += 0.7

                _slide_footer(s_pc, slide_num)
                slide_num += 1
        except Exception:
            pass  # Skip period comparison slide if date parsing fails

    # ── Per-brand slides — AI Advertiser Insights ─────────────────────────────
    if sections.get("ai_advertiser_insights"):
        for brand in campaigns:

            # ── Performance Breakdown ─────────────────────────────────────────
            s2 = _new_blank_slide(prs)
            _tb(s2, 0.4, 0.17, 12.5, 0.65,
                f"{brand}  —  Performance Breakdown", 20, bold=True)
            _box(s2, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)

            breakdown = _ai_pptx(
                api_key,
                f"""Campaign performance data:
{data_text}

Write a performance breakdown for brand: "{brand}"

Use EXACTLY this format — only include sections where data exists for this brand:

DISPLAY:
• [CPM, CPC, CTR focus — cite best/worst line item, max 15 words]
• [second point]
• [third point — max 3 bullets]

VIDEO:
• [CPV, VTR focus — cite best/worst line item, max 15 words]
• [second point]
• [third point — max 3 bullets]

YOUTUBE:
• [CPV, VTR focus — cite best/worst creative, max 15 words]
• [second point]
• [third point — max 3 bullets]

Only include DISPLAY / VIDEO / YOUTUBE headings that have real data for this brand.""",
                max_tokens=350,
            )

            parsed_sections = _parse_io_sections(breakdown)
            if parsed_sections:
                n     = len(parsed_sections)
                col_w = (12.5 - (n - 1) * 0.3) / n
                cx    = 0.4
                for sec_title, sec_bullets in parsed_sections:
                    _box(s2, cx, 1.05, col_w, 0.38, _PPT_DIM)
                    _tb(s2, cx + 0.12, 1.1, col_w - 0.24, 0.3,
                        sec_title.upper(), 11, bold=True, color=_PPT_BLUE)
                    if sec_bullets:
                        _add_bullets(s2, cx + 0.08, 1.58,
                                     col_w - 0.16, 5.5, sec_bullets[:3], size=10)
                    cx += col_w + 0.3

            _slide_footer(s2, slide_num)
            slide_num += 1

            # ── Recommendations ───────────────────────────────────────────────
            s3 = _new_blank_slide(prs)
            _tb(s3, 0.4, 0.17, 12.5, 0.65,
                f"{brand}  —  Recommendations", 20, bold=True)
            _box(s3, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)

            recs_raw = _ai_pptx(
                api_key,
                f"""Campaign performance data:
{data_text}

Write exactly 5 optimisation recommendations for brand: "{brand}"

Rules:
- Each must start with an action verb: Increase / Reduce / Pause / Test / Shift / Reallocate
- Reference specific line items, creatives, or metrics from the data
- Max 18 words per line
- Return only the 5 lines — no bullet symbols, no numbering, no extra text""",
                max_tokens=250,
            )
            bullets = [ln.strip().lstrip("•-–—●0123456789.) ").strip()
                       for ln in recs_raw.split("\n")
                       if ln.strip() and not ln.strip().startswith("#")]
            bullets = [b for b in bullets if b][:5]
            _add_bullets(s3, 0.6, 1.12, 12.1, 5.85, bullets, size=13)

            _slide_footer(s3, slide_num)
            slide_num += 1

        # ── Budget Shift Recommendations — part of AI Advertiser Insights ─────
        s_last = _new_blank_slide(prs)
        _tb(s_last, 0.4, 0.17, 12.5, 0.65,
            "Budget Shift Recommendations", 20, bold=True)
        _box(s_last, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)

        budget_raw = _ai_pptx(
            api_key,
            f"""Campaign performance data:
{data_text}

For each brand, provide one budget reallocation recommendation.
Format EXACTLY as pipe-separated lines — one line per brand, no header row:
Brand Name | Best Performing IO | One-sentence recommendation starting with an action verb

Example:
Nike Summer 2024 | Display | Increase Display budget by 20% — lowest CPM at $3.20
Coke Q3 | YouTube | Shift 15% from Display to YouTube — VTR of 68% outperforms

Only include brands from the data. Cite actual numbers. No extra text.""",
            max_tokens=350,
        )
        _add_budget_table(s_last, budget_raw, x=0.4, y=1.08, w=12.5, h=5.9)

        _slide_footer(s_last, slide_num)
        slide_num += 1

    # ── AI-Driven Analysis slide ──────────────────────────────────────────────
    if sections.get("ai_driven_analysis"):
        s_ai = _new_blank_slide(prs)
        _tb(s_ai, 0.4, 0.17, 12.5, 0.65, "AI-Driven Analysis", 20, bold=True)
        _box(s_ai, 0.4, 0.86, 12.5, 0.04, _PPT_BLUE)

        if ai_analysis_cache and isinstance(ai_analysis_cache, dict):
            # Reuse cached result from session state — no extra API call needed
            summary_txt   = ai_analysis_cache.get("summary", "")
            anomaly_txt   = ai_analysis_cache.get("top_anomaly", "")
        else:
            # No cached result — generate brief text via AI
            ai_driven_raw = _ai_pptx(
                api_key,
                f"""Campaign performance data:
{data_text}

Provide a brief AI analysis in two parts:
1. SUMMARY: One paragraph on the most important patterns in this dataset.
2. TOP FINDING: The single most important anomaly or insight.

Format exactly as:
SUMMARY: [text]
TOP FINDING: [text]""",
                max_tokens=300,
            )
            summary_txt = ""
            anomaly_txt = ""
            for _line in ai_driven_raw.split("\n"):
                if _line.startswith("SUMMARY:"):
                    summary_txt = _line.replace("SUMMARY:", "").strip()
                elif _line.startswith("TOP FINDING:"):
                    anomaly_txt = _line.replace("TOP FINDING:", "").strip()

        y_pos = 1.1
        if summary_txt:
            _tb(s_ai, 0.5, y_pos, 12.0, 0.28, "SUMMARY", 10, bold=True, color=_PPT_BLUE)
            y_pos += 0.35
            txb_s = s_ai.shapes.add_textbox(Inches(0.5), Inches(y_pos),
                                            Inches(12.0), Inches(1.5))
            tf_s        = txb_s.text_frame
            tf_s.word_wrap = True
            run_s       = tf_s.paragraphs[0].add_run()
            run_s.text  = summary_txt
            run_s.font.name      = "Calibri"
            run_s.font.size      = Pt(11)
            run_s.font.color.rgb = _PPT_WHITE
            y_pos += 1.7
        if anomaly_txt:
            _tb(s_ai, 0.5, y_pos, 12.0, 0.28, "TOP FINDING", 10, bold=True, color=_PPT_BLUE)
            y_pos += 0.35
            txb_a = s_ai.shapes.add_textbox(Inches(0.5), Inches(y_pos),
                                            Inches(12.0), Inches(1.2))
            tf_a        = txb_a.text_frame
            tf_a.word_wrap = True
            run_a       = tf_a.paragraphs[0].add_run()
            run_a.text  = anomaly_txt
            run_a.font.name      = "Calibri"
            run_a.font.size      = Pt(11)
            run_a.font.color.rgb = _PPT_WHITE

        _slide_footer(s_ai, slide_num)
        slide_num += 1

    # Serialise to buffer and return
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ── File uploader ─────────────────────────────────────────────────────────────
# Expander label and default state depend on whether data was loaded on the
# previous run. Session state "_uploader_loaded_info" is set after files are
# processed so the label and collapse state update on the next rerun.
_loaded_info = st.session_state.get("_uploader_loaded_info")
if _loaded_info:
    _exp_label    = (f"📁 Upload Data — {_loaded_info['label']} loaded "
                     f"({_loaded_info['rows']:,} rows)")
    _exp_expanded = False
else:
    _exp_label    = "📁 Upload Data"
    _exp_expanded = True

# Only show file upload in File Upload (one-off) mode
if _pi_source_mode != 'ongoing':
    with st.expander(_exp_label, expanded=_exp_expanded):
        uploaded_files = st.file_uploader(
            "Drag and drop files here, or click to browse. Accepts CSV, TSV, Excel (.xlsx / .xls). Multiple files accepted.",
            type=["csv", "tsv", "xlsx", "xls"],
            accept_multiple_files=True,
        )
else:
    uploaded_files = None

# ── Main header ───────────────────────────────────────────────────────────────
st.title("Performance & Insights")

# ── Process files ─────────────────────────────────────────────────────────────
if not uploaded_files and _pi_source_mode != 'ongoing':
    # Clear stored file info so expander resets to expanded when files are removed
    st.session_state.pop("_uploader_loaded_info", None)

    # File Upload mode empty state — stop rendering when no file is uploaded
    if _pi_source_mode == 'one-off':
        st.markdown("""
        <div style="text-align:center;padding:60px 20px;background:#FFFFFF;
        border-radius:12px;box-shadow:0 4px 12px rgba(0,0,0,0.08);margin:20px 0;">
            <div style="font-size:48px;margin-bottom:16px;">📂</div>
            <div style="font-size:20px;font-weight:700;color:#111827;margin-bottom:8px;">
                No data loaded yet
            </div>
            <div style="font-size:15px;color:#6B7280;margin-bottom:4px;">
                Drag and drop a DSP export to get started
            </div>
            <div style="font-size:13px;color:#9CA3AF;">
                Accepted formats: CSV, XLSX, XLS, TSV
            </div>
        </div>
        """, unsafe_allow_html=True)
        st.stop()

    st.stop()

else:
    # ── API (ongoing) mode: load mock data ────────────────────────────────────
    if _pi_source_mode == 'ongoing':
        df_all = generate_api_mock_data()
        # Share with portfolio overview page
        st.session_state["portfolio_df"] = df_all

    else:
        # ── File Upload mode: load and stack all uploaded files ───────────────
        frames = []
        file_info = []

        with st.spinner("Reading and normalising files…"):
            for f in (uploaded_files or []):
                df_file, dsp = load_and_normalise(f)
                frames.append(df_file)
                file_info.append({"name": f.name, "rows": len(df_file), "dsp": dsp})

        # Combine all files into a single DataFrame
        df_all = pd.concat(frames, ignore_index=True)

        # Store file info in session state so the expander label updates on next rerun
        _file_label = (uploaded_files[0].name if len(uploaded_files or []) == 1
                       else f"{len(uploaded_files or [])} files")
        st.session_state["_uploader_loaded_info"] = {
            "label": _file_label,
            "rows":  len(df_all),
        }

    # Parse date column to datetime so the date range filter works correctly
    if "date" in df_all.columns:
        df_all["date"] = pd.to_datetime(df_all["date"], errors="coerce")

    # ── Global filter bar — all filters in one row ─────────────────────────────
    # Order: ADVERTISER | DSP | CAMPAIGN | LINE ITEM | CREATIVE | METRIC | DATE RANGE
    # Cascading: advertiser selection narrows campaign options, campaign narrows
    # line item, line item narrows creative.
    _adv_col  = "advertiser" if "advertiser" in df_all.columns else (
                "campaign"   if "campaign"   in df_all.columns else None)
    _has_date = "date" in df_all.columns

    # Build metric options before the filter bar so the Metric selectbox can use them
    _METRIC_LABELS = {
        "spend_usd": "Spend", "impressions": "Impressions", "clicks": "Clicks",
        "conversions": "Conversions", "cpm": "CPM", "ctr": "CTR",
    }
    avail_metrics = {k: v for k, v in _METRIC_LABELS.items() if k in df_all.columns}
    for _col in df_all.select_dtypes(include="number").columns:
        if _col not in avail_metrics and _col not in ("source_file",):
            avail_metrics[_col] = _col.replace("_", " ").title()

    # Shared label style — identical across every filter widget
    _CL = (
        "font-size:11px;font-weight:600;text-transform:uppercase;"
        "letter-spacing:0.04em;color:#1B2A4A;margin:0 0 4px 0;"
        "font-family:'Poppins',system-ui,sans-serif;line-height:1.2;display:block;"
    )

    _fc = st.columns([1.2, 0.8, 1.5, 1.5, 1.2, 1.0, 1.8])

    # ── ADVERTISER ────────────────────────────────────────────────────────────
    with _fc[0]:
        st.markdown(f'<span style="{_CL}">Advertiser</span>', unsafe_allow_html=True)
        if _adv_col:
            _adv_opts = sorted(df_all[_adv_col].dropna().unique().tolist())
            sel_adv = st.multiselect("Advertiser", _adv_opts, default=[],
                                     placeholder="All", key="gf_adv",
                                     label_visibility="collapsed")
        else:
            sel_adv = []

    # ── DSP ───────────────────────────────────────────────────────────────────
    with _fc[1]:
        st.markdown(f'<span style="{_CL}">DSP</span>', unsafe_allow_html=True)
        _dsp_opts = sorted(df_all["dsp_source"].dropna().unique().tolist())
        sel_dsp = st.multiselect("DSP", _dsp_opts, default=[],
                                 placeholder="All", key="gf_dsp",
                                 label_visibility="collapsed")

    # Apply advertiser + DSP to get cascading base for Campaign/LI/Creative
    # Empty list means "All" — no filter applied
    _df_base = df_all.copy()
    if sel_adv and _adv_col:
        _df_base = _df_base[_df_base[_adv_col].isin(sel_adv)]
    if sel_dsp:
        _df_base = _df_base[_df_base["dsp_source"].isin(sel_dsp)]

    # ── CAMPAIGN (cascades from advertiser + DSP) ─────────────────────────────
    with _fc[2]:
        st.markdown(f'<span style="{_CL}">Campaign</span>', unsafe_allow_html=True)
        if "campaign" in _df_base.columns:
            _camp_opts = sorted(_df_base["campaign"].dropna().unique().tolist())
            sel_campaign = st.multiselect("Campaign", _camp_opts, default=[],
                                          placeholder="All", key="gf_campaign",
                                          label_visibility="collapsed")
        else:
            sel_campaign = []

    _df_camp = _df_base.copy()
    if sel_campaign and "campaign" in _df_camp.columns:
        _df_camp = _df_camp[_df_camp["campaign"].isin(sel_campaign)]

    # ── LINE ITEM (cascades from campaign) ────────────────────────────────────
    with _fc[3]:
        st.markdown(f'<span style="{_CL}">Line Item</span>', unsafe_allow_html=True)
        if "line_item" in _df_camp.columns:
            _li_opts = sorted(_df_camp["line_item"].dropna().unique().tolist())
            sel_line_item = st.multiselect("Line Item", _li_opts, default=[],
                                           placeholder="All", key="gf_line_item",
                                           label_visibility="collapsed")
        else:
            sel_line_item = []

    _df_li = _df_camp.copy()
    if sel_line_item and "line_item" in _df_li.columns:
        _df_li = _df_li[_df_li["line_item"].isin(sel_line_item)]

    # ── CREATIVE (cascades from line item) ────────────────────────────────────
    with _fc[4]:
        st.markdown(f'<span style="{_CL}">Creative</span>', unsafe_allow_html=True)
        if "creative" in _df_li.columns:
            _cr_opts = sorted(_df_li["creative"].dropna().unique().tolist())
            sel_creative = st.multiselect("Creative", _cr_opts, default=[],
                                          placeholder="All", key="gf_creative",
                                          label_visibility="collapsed")
        else:
            sel_creative = []

    # ── METRIC ────────────────────────────────────────────────────────────────
    with _fc[5]:
        st.markdown(f'<span style="{_CL}">Metric</span>', unsafe_allow_html=True)
        sel_label = st.selectbox("Metric", list(avail_metrics.values()),
                                 key="gf_metric", label_visibility="collapsed")
    sel_col = next(k for k, v in avail_metrics.items() if v == sel_label)

    # ── DATE RANGE ────────────────────────────────────────────────────────────
    with _fc[6]:
        st.markdown(f'<span style="{_CL}">Date Range</span>', unsafe_allow_html=True)
        if _has_date:
            _min_date   = df_all["date"].dropna().min().date()
            _max_date   = df_all["date"].dropna().max().date()
            _date_range = st.date_input(
                "Date Range",
                value=(_min_date, _max_date),
                min_value=_min_date,
                max_value=_max_date,
                key="gf_date_range",
                label_visibility="collapsed",
            )
        else:
            _date_range = None

    # ── Apply all filters to build df_metrics (drives every section below) ────
    df_metrics = _df_li.copy()
    if sel_creative and "creative" in df_metrics.columns:
        df_metrics = df_metrics[df_metrics["creative"].isin(sel_creative)]
    if _has_date and isinstance(_date_range, (list, tuple)) and len(_date_range) == 2:
        _start, _end = _date_range
        df_metrics = df_metrics[
            (df_metrics["date"].dt.date >= _start) &
            (df_metrics["date"].dt.date <= _end)
        ]

    # Also narrow df_all by date + DSP so sections that still reference df_all
    # (period comparison, AI insights, PPTX) see a consistently filtered dataset
    if sel_dsp:
        df_all = df_all[df_all["dsp_source"].isin(sel_dsp)].copy()
    if _has_date and isinstance(_date_range, (list, tuple)) and len(_date_range) == 2:
        _start, _end = _date_range
        df_all = df_all[
            (df_all["date"].dt.date >= _start) &
            (df_all["date"].dt.date <= _end)
        ].copy()

    # ── Summary metrics ───────────────────────────────────────────────────────
    total_impressions = df_metrics["impressions"].sum()  if "impressions" in df_metrics.columns else 0
    total_clicks      = df_metrics["clicks"].sum()       if "clicks"      in df_metrics.columns else 0
    total_spend       = df_metrics["spend_usd"].sum()    if "spend_usd"   in df_metrics.columns else 0
    # Always recalculate CTR from totals — never average a pre-calculated CTR column
    avg_ctr           = (total_clicks / total_impressions) if total_impressions > 0 else None
    avg_cpm           = df_metrics["cpm"].mean()         if "cpm"         in df_metrics.columns else None

    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        st.markdown(metric_card("Total Impressions", f"{total_impressions:,.0f}"), unsafe_allow_html=True)
    with col2:
        st.markdown(metric_card("Total Clicks", f"{total_clicks:,.0f}"), unsafe_allow_html=True)
    with col3:
        st.markdown(metric_card("Total Spend", f"A${total_spend:,.0f}"), unsafe_allow_html=True)
    with col4:
        st.markdown(metric_card("Avg CTR", f"{avg_ctr:.2%}" if avg_ctr is not None else "N/A"), unsafe_allow_html=True)
    with col5:
        st.markdown(metric_card("Avg CPM", f"A${avg_cpm:,.2f}" if avg_cpm is not None else "N/A"), unsafe_allow_html=True)

    # 24px breathing room between the metric cards row and the first chart row
    st.markdown('<div style="margin-top:24px;"></div>', unsafe_allow_html=True)

    # Build top-level summary for PPTX export — group by best available dimension
    _top_col = next((c for c in ("advertiser", "campaign", "insertion_order")
                     if c in df_metrics.columns), None)
    if _top_col:
        _agg = {c: (c, "sum") for c in ["impressions", "clicks", "spend_usd"] if c in df_metrics.columns}
        camp_summary = df_metrics.groupby(_top_col).agg(**_agg).reset_index()
        if "impressions" in camp_summary.columns and "clicks" in camp_summary.columns:
            camp_summary["ctr"] = camp_summary["clicks"] / camp_summary["impressions"]
        if "impressions" in camp_summary.columns and "spend_usd" in camp_summary.columns:
            camp_summary["cpm"] = camp_summary["spend_usd"] / camp_summary["impressions"] * 1000
    else:
        camp_summary = pd.DataFrame()

    # ── Charts ────────────────────────────────────────────────────────────────
    # Colour palette — one colour per brand, consistent across all charts
    # Design-system palette only — RAG colors (#10B981 / #F59E0B / #EF4444)
    # are reserved for status indicators and must never appear in performance charts.
    PALETTE = ["#1B2A4A", "#F5A623", "#2C4A7A", "#F7B84B", "#162238",
               "#E8951A", "#4A7AB5", "#3D5A80", "#6B85A8", "#F7C566"]

    def no_data_msg(msg):
        st.markdown(
            f"<div style='text-align:center;padding:40px;color:#9ca3af;"
            f"border:1px dashed #e5e7eb;border-radius:8px;font-size:13px;'>{msg}</div>",
            unsafe_allow_html=True,
        )

    def trunc(s, n=20):
        """Truncate a string to n characters for axis labels."""
        return s[:n] + "…" if len(str(s)) > n else str(s)

    def apply_chart_style(fig, xaxis_title="", yaxis_title="", horizontal=False, height=380):
        """
        Apply the shared Pacebird Plotly template, then bar-chart-specific overrides.
        horizontal=True swaps grid to x-axis for horizontal bar charts.
        height controls the chart height in pixels (default 380 for bar charts).
        """
        apply_plotly_style(fig, height=height)
        fig.update_layout(
            showlegend=False,
            margin=dict(l=40, r=16, t=44, b=36),  # tight margins for compact cards
        )
        if horizontal:
            # Horizontal bars: grid on x, clear y grid
            fig.update_xaxes(showgrid=True, gridcolor="rgba(27,42,74,0.08)")
            fig.update_yaxes(showgrid=False, showline=False)
        if xaxis_title:
            fig.update_xaxes(title_text=xaxis_title)
        if yaxis_title:
            fig.update_yaxes(title_text=yaxis_title)
        return fig

    # ── Chart dimension configs ────────────────────────────────────────────────
    # Title is the dimension name only (no "Spend by" prefix).
    # dim_col is the normalised internal column name to group by.
    # None = auto-detect from df_all columns.
    CHART_CONFIGS = [
        {"title": "Advertiser",       "dim_col": "advertiser",  "horizontal": False},
        {"title": "Campaign",         "dim_col": "campaign",    "horizontal": True},
        {"title": "Line Item",        "dim_col": "line_item",   "horizontal": True},
        {"title": "Creative",         "dim_col": "creative",    "horizontal": True},
        {"title": "Audience Segment", "dim_col": None,          "horizontal": False},
    ]

    # Auto-detect audience / segment column for the last chart
    for _cfg in CHART_CONFIGS:
        if _cfg["dim_col"] is None:
            for _col in df_metrics.columns:
                if any(kw in _col.lower() for kw in ("audience", "segment")):
                    _cfg["dim_col"] = _col
                    break

    def fmt_val(v, col):
        """Format a metric value for chart data labels."""
        if col in ("spend_usd", "cpm"):
            if v >= 1_000_000:
                return f"A${v/1e6:.2f}M"
            elif v >= 1_000:
                return f"A${v/1e3:.0f}K"
            return f"A${v:.2f}"
        elif col == "ctr":
            return f"{v:.2%}"
        return f"{v:,.0f}"

    def get_agg(df, dim_col, metric_col):
        """
        Aggregate df by dim_col for the chosen metric.
        CPM and CTR are recalculated from raw totals — never summed or
        averaged directly, which would give wrong results for rate metrics.
        """
        if metric_col == "cpm" and "spend_usd" in df.columns and "impressions" in df.columns:
            grp = (df.groupby(dim_col)
                     .agg(spend_usd=("spend_usd", "sum"),
                          impressions=("impressions", "sum"))
                     .reset_index())
            grp["cpm"] = grp["spend_usd"] / grp["impressions"].clip(lower=1) * 1000
            return grp[[dim_col, "cpm"]]
        elif metric_col == "ctr" and "clicks" in df.columns and "impressions" in df.columns:
            grp = (df.groupby(dim_col)
                     .agg(clicks=("clicks", "sum"),
                          impressions=("impressions", "sum"))
                     .reset_index())
            grp["ctr"] = grp["clicks"] / grp["impressions"].clip(lower=1)
            return grp[[dim_col, "ctr"]]
        else:
            return df.groupby(dim_col)[metric_col].sum().reset_index()

    # Only render charts whose dimension column was detected in the data
    visible_charts = [
        cfg for cfg in CHART_CONFIGS
        if cfg["dim_col"] and cfg["dim_col"] in df_metrics.columns
    ]

    if not visible_charts or not avail_metrics:
        no_data_msg("No chartable dimensions or metrics found in the uploaded data.")
    else:
        # df_metrics from the global filter bar drives all charts
        df_chart_base = df_metrics.copy()

        # ── 2×2 chart grid — no per-chart filter widgets ─────────────────────────
        for _i in range(0, len(visible_charts), 2):
            _row_cols = st.columns(2)
            for _j, cfg in enumerate(visible_charts[_i:_i+2]):
                with _row_cols[_j]:
                    dim_col    = cfg["dim_col"]
                    title      = cfg["title"]
                    horizontal = cfg.get("horizontal", False)

                    # Sort ascending for horizontal so top bar = highest value
                    agg_df = (get_agg(df_chart_base, dim_col, sel_col)
                              .sort_values(sel_col, ascending=horizontal)
                              .head(15))

                    if agg_df.empty or agg_df[sel_col].sum() == 0:
                        no_data_msg(f"No {sel_label} data for {title}.")
                        continue

                    bar_labels  = [fmt_val(v, sel_col) for v in agg_df[sel_col]]
                    colors      = [PALETTE[k % len(PALETTE)] for k in range(len(agg_df))]
                    full_names  = agg_df[dim_col].astype(str).tolist()
                    trunc_names = [trunc(s) for s in full_names]

                    if horizontal:
                        # Horizontal bars — category on y-axis
                        # Used for Campaign, Line Item, and Creative where names are long
                        fig = go.Figure(go.Bar(
                            y=trunc_names,
                            x=agg_df[sel_col],
                            customdata=full_names,
                            orientation="h",
                            marker_color=colors,
                            text=bar_labels,
                            textposition="outside",
                            textfont=dict(size=10, color=SECONDARY),
                            hovertemplate=(
                                f"<b>%{{customdata}}</b><br>"
                                f"{sel_label}: %{{text}}<extra></extra>"
                            ),
                        ))
                        max_v = agg_df[sel_col].max()
                        fig.update_layout(
                            bargap=0.35,
                            xaxis_range=[0, max_v * 1.35] if max_v > 0 else [0, 1],
                            title=dict(
                                text=title,
                                font=dict(size=15, color=SECONDARY,
                                          family="Poppins, system-ui, sans-serif"),
                                x=0.02, xanchor="left",
                            ),
                        )
                        if sel_col in ("spend_usd", "cpm"):
                            fig.update_xaxes(tickprefix="A$", tickformat=",.0f")
                        elif sel_col == "ctr":
                            fig.update_xaxes(tickformat=".1%")
                        else:
                            fig.update_xaxes(tickformat=",")
                        apply_chart_style(fig, horizontal=True, height=380)
                    else:
                        # Vertical bars — category on x-axis
                        # Used for Advertiser and Creative where names are shorter
                        fig = go.Figure(go.Bar(
                            x=trunc_names,
                            y=agg_df[sel_col],
                            customdata=full_names,
                            marker_color=colors,
                            text=bar_labels,
                            textposition="outside",
                            textfont=dict(size=10, color=SECONDARY),
                            hovertemplate=(
                                f"<b>%{{customdata}}</b><br>"
                                f"{sel_label}: %{{text}}<extra></extra>"
                            ),
                        ))
                        max_v = agg_df[sel_col].max()
                        fig.update_layout(
                            bargap=0.45,
                            yaxis_range=[0, max_v * 1.3] if max_v > 0 else [0, 1],
                            title=dict(
                                text=title,
                                font=dict(size=15, color=SECONDARY,
                                          family="Poppins, system-ui, sans-serif"),
                                x=0.02, xanchor="left",
                            ),
                        )
                        if sel_col in ("spend_usd", "cpm"):
                            fig.update_yaxes(tickprefix="A$", tickformat=",.0f")
                        elif sel_col == "ctr":
                            fig.update_yaxes(tickformat=".1%")
                        else:
                            fig.update_yaxes(tickformat=",")
                        apply_chart_style(fig, height=380)
                        fig.update_xaxes(tickangle=0)

                    st.plotly_chart(fig, use_container_width=True, config=PLOTLY_CONFIG,
                                    key=f"chart_{title.replace(' ', '_')}_{sel_col}")

    # ── Breakdown donut charts — DSP / Device / Environment / Exchange ─────────
    # 4-column row; each chart shows share of the selected metric by that dimension.
    # A chart is hidden when its column is absent from the data.
    _DONUT_PALETTE = ["#1B2A4A", "#F5A623", "#3D5A80", "#F7C566", "#6B85A8",
                      "#2C4A7A", "#E8951A", "#4A7AB5"]
    _DONUT_DIMS = [
        ("dsp_source",  "DSP"),
        ("device_type", "Device"),
        ("environment", "Environment"),
        ("buy_type",    "Buy Type"),
    ]
    _donut_visible = [(col, lbl) for col, lbl in _DONUT_DIMS
                      if col in df_metrics.columns and df_metrics[col].notna().any()]

    if _donut_visible and sel_col in df_metrics.columns:
        _donut_cols = st.columns(len(_donut_visible))
        for _di, (d_col, d_lbl) in enumerate(_donut_visible):
            with _donut_cols[_di]:
                _d_agg = (
                    df_metrics.groupby(d_col)[sel_col].sum()
                    .reset_index()
                    .sort_values(sel_col, ascending=False)
                )
                # Skip if all values are zero
                if _d_agg[sel_col].sum() == 0:
                    continue
                _d_fig = go.Figure(go.Pie(
                    labels=_d_agg[d_col].astype(str),
                    values=_d_agg[sel_col],
                    hole=0.5,
                    marker=dict(colors=_DONUT_PALETTE[:len(_d_agg)]),
                    textinfo="percent",
                    textfont=dict(size=11, family="Poppins, system-ui, sans-serif"),
                    hovertemplate=(
                        f"<b>%{{label}}</b><br>{sel_label}: "
                        f"%{{value:,.0f}} (%{{percent}})<extra></extra>"
                    ),
                ))
                _d_fig.update_layout(
                    title=dict(
                        text=d_lbl,
                        font=dict(size=14, color=SECONDARY,
                                  family="Poppins, system-ui, sans-serif"),
                        x=0.5, xanchor="center",
                    ),
                    showlegend=True,
                    legend=dict(
                        orientation="h",
                        yanchor="top",
                        y=-0.02,
                        xanchor="center",
                        x=0.5,
                        font=dict(size=10, family="Poppins, system-ui, sans-serif"),
                        tracegroupgap=0,
                        itemwidth=30,
                    ),
                    margin=dict(l=10, r=10, t=30, b=10),
                    height=300,
                    paper_bgcolor="rgba(0,0,0,0)",
                    plot_bgcolor="rgba(0,0,0,0)",
                )
                st.plotly_chart(_d_fig, use_container_width=True,
                                config=PLOTLY_CONFIG,
                                key=f"donut_{d_col}_{sel_col}")

    # ── Daily Performance chart ────────────────────────────────────────────────
    # Shown when the data contains a date column and at least one summable metric.
    if "date" in df_metrics.columns:
        st.subheader("Daily Performance")

        # Only include summable metrics — CPM and CTR need to be recalculated from
        # daily raw totals, not summed, so we exclude them from this chart for simplicity.
        daily_summable = {k: v for k, v in avail_metrics.items()
                          if k not in ("cpm", "ctr", "cpm_raw", "ctr_raw")}

        if not daily_summable:
            no_data_msg("No summable metrics found for the daily chart.")
        else:
            # Controls row: advertiser filter | left axis metric | right axis metric
            _d_col1, _d_col2, _d_col3 = st.columns([2, 1, 1])

            with _d_col1:
                if _adv_col:
                    daily_adv_opts = ["All"] + sorted(df_metrics[_adv_col].dropna().unique().tolist())
                    sel_daily_adv = st.selectbox(
                        "Advertiser", daily_adv_opts,
                        key="daily_adv_filter",
                        label_visibility="collapsed",
                    )
                    df_daily_src = (df_metrics if sel_daily_adv == "All"
                                    else df_metrics[df_metrics[_adv_col] == sel_daily_adv])
                else:
                    df_daily_src = df_metrics

            daily_keys   = list(daily_summable.keys())
            daily_labels = list(daily_summable.values())

            # Default left axis = Clicks, right axis = Impressions (when available)
            left_default  = daily_keys.index("clicks")      if "clicks"      in daily_keys else 0
            right_default = (daily_keys.index("impressions") if "impressions" in daily_keys
                             else (1 if len(daily_keys) > 1 else 0))

            with _d_col2:
                left_label = st.selectbox(
                    "Left axis", daily_labels,
                    index=left_default,
                    key="daily_left_metric",
                    label_visibility="collapsed",
                )
                left_col = daily_keys[daily_labels.index(left_label)]

            with _d_col3:
                right_label = st.selectbox(
                    "Right axis", daily_labels,
                    index=right_default,
                    key="daily_right_metric",
                    label_visibility="collapsed",
                )
                right_col = daily_keys[daily_labels.index(right_label)]

            # Aggregate by date — sum both selected metrics
            daily_agg_cols = list(dict.fromkeys([left_col, right_col]))  # deduplicated, ordered
            daily = (
                df_daily_src.groupby("date")[daily_agg_cols]
                .sum()
                .reset_index()
                .sort_values("date")
            )

            # Dual Y-axis layout
            fig_daily = make_subplots(specs=[[{"secondary_y": True}]])

            # Left Y-axis trace — solid blue line
            fig_daily.add_trace(
                go.Scatter(
                    x=daily["date"],
                    y=daily[left_col],
                    name=left_label,
                    mode="lines+markers",
                    line=dict(color="#1B2A4A", width=2, dash="solid"),
                    marker=dict(symbol="circle", size=6, color="#1B2A4A"),
                    hovertemplate=f"%{{x}}<br>{left_label}: <b>%{{y:,}}</b><extra></extra>",
                ),
                secondary_y=False,
            )

            # Right Y-axis trace — dashed orange line (only when a different metric is chosen)
            if right_col != left_col:
                fig_daily.add_trace(
                    go.Scatter(
                        x=daily["date"],
                        y=daily[right_col],
                        name=right_label,
                        mode="lines+markers",
                        line=dict(color="#EF9F27", width=2, dash="dash"),
                        marker=dict(symbol="circle", size=6, color="#EF9F27"),
                        hovertemplate=f"%{{x}}<br>{right_label}: <b>%{{y:,}}</b><extra></extra>",
                    ),
                    secondary_y=True,
                )

            fig_daily.update_layout(
                title=dict(text="Daily Performance", font=dict(size=15, color="#111827")),
                plot_bgcolor="#FFFFFF",
                paper_bgcolor="#FFFFFF",
                font=dict(family="Poppins, system-ui, sans-serif", color="#374151"),
                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                margin=dict(l=60, r=60, t=60, b=50),
                hovermode="x unified",
            )
            fig_daily.update_yaxes(
                title_text=left_label, secondary_y=False,
                tickformat=",", gridcolor="#F3F4F6",
                title_font=dict(color="#1B2A4A"),
                tickfont=dict(color="#1B2A4A"),
            )
            fig_daily.update_yaxes(
                title_text=right_label if right_col != left_col else "",
                secondary_y=True,
                tickformat=",",
                title_font=dict(color="#EF9F27"),
                tickfont=dict(color="#EF9F27"),
                showgrid=False,
            )
            fig_daily.update_xaxes(gridcolor="#F3F4F6")

            st.plotly_chart(fig_daily, use_container_width=True,
                            config=PLOTLY_CONFIG,
                            key="chart_daily_performance")

    # ─────────────────────────────────────────────────────────────────────────
    # PERIOD COMPARISON — Week-on-Week / Month-on-Month
    # Only shown when date column spans at least 14 days
    # ─────────────────────────────────────────────────────────────────────────
    _period_comp_available = False  # updated to True inside the block when conditions are met
    if 'date' in df_all.columns:
        try:
            df_all['_date_parsed'] = pd.to_datetime(df_all['date'], errors='coerce')
            date_range = (df_all['_date_parsed'].max() - df_all['_date_parsed'].min()).days

            if date_range >= 14:
                _period_comp_available = True
                st.markdown("---")
                st.markdown("## 📅 Period Comparison")

                period_toggle = st.radio(
                    "Comparison period:",
                    ["WoW", "MoM", "QoQ", "YoY"],
                    horizontal=True,
                    key="period_toggle"
                )

                today = df_all['_date_parsed'].max()

                if period_toggle == "WoW":
                    # Week on Week: last 7 days vs prior 7 days
                    current_start  = today - pd.Timedelta(days=6)
                    previous_start = today - pd.Timedelta(days=13)
                    previous_end   = today - pd.Timedelta(days=7)
                elif period_toggle == "MoM":
                    # Month on Month: last 30 days vs prior 30 days
                    current_start  = today - pd.Timedelta(days=29)
                    previous_start = today - pd.Timedelta(days=59)
                    previous_end   = today - pd.Timedelta(days=30)
                elif period_toggle == "QoQ":
                    # Quarter on Quarter: last 90 days vs prior 90 days
                    current_start  = today - pd.Timedelta(days=89)
                    previous_start = today - pd.Timedelta(days=179)
                    previous_end   = today - pd.Timedelta(days=90)
                else:  # YoY
                    # Year on Year: last 365 days vs prior 365 days
                    current_start  = today - pd.Timedelta(days=364)
                    previous_start = today - pd.Timedelta(days=729)
                    previous_end   = today - pd.Timedelta(days=365)

                current_df = df_all[df_all['_date_parsed'] >= current_start]
                previous_df = df_all[(df_all['_date_parsed'] >= previous_start) & (df_all['_date_parsed'] <= previous_end)]

                def safe_sum(d, col):
                    return d[col].sum() if col in d.columns and len(d) > 0 else 0

                def calc_metrics(d):
                    imps = safe_sum(d, 'impressions')
                    clicks = safe_sum(d, 'clicks')
                    spend = safe_sum(d, 'spend_usd')
                    convs = safe_sum(d, 'conversions')
                    return {
                        'Spend': spend,
                        'Impressions': imps,
                        'Clicks': clicks,
                        'CTR': (clicks / imps * 100) if imps > 0 else 0,
                        'CPA': (spend / convs) if convs > 0 else 0,
                    }

                cur = calc_metrics(current_df)
                prev = calc_metrics(previous_df)

                # Build comparison table
                rows = []
                format_map = {
                    'Spend': '${:,.0f}',
                    'Impressions': '{:,.0f}',
                    'Clicks': '{:,.0f}',
                    'CTR': '{:.2f}%',
                    'CPA': '${:.2f}',
                }
                higher_better = {'Spend': True, 'Impressions': True, 'Clicks': True, 'CTR': True, 'CPA': False}

                for metric, fmt in format_map.items():
                    c_val = cur[metric]
                    p_val = prev[metric]
                    if p_val > 0:
                        change_pct = (c_val - p_val) / p_val * 100
                        change_str = f"{change_pct:+.1f}%"
                        hib = higher_better[metric]
                        is_improvement = (change_pct > 0 and hib) or (change_pct < 0 and not hib)
                        is_decline = (change_pct < 0 and hib) or (change_pct > 0 and not hib)
                        abs_change = abs(change_pct)

                        if abs_change > 10 and is_improvement:
                            status = "✅"
                        elif abs_change > 10 and is_decline:
                            status = "⚠️"
                        else:
                            status = "➡️"
                    else:
                        change_str = "N/A"
                        status = "➡️"

                    rows.append({
                        "Metric": metric,
                        "This Period": fmt.format(c_val),
                        "Last Period": fmt.format(p_val) if p_val > 0 else "—",
                        "Change": change_str,
                        "Status": status,
                    })

                comparison_df = pd.DataFrame(rows)
                st.dataframe(comparison_df, use_container_width=True, hide_index=True)
        except Exception:
            pass  # Silently skip if date parsing fails

    # ── AI Insights ───────────────────────────────────────────────────────────
    # Detect the best grouping column — works with any DSP export structure
    group_col, group_label = detect_grouping_column(df_all)

    # Pre-compute camp_summary here so it is available for the report export dialog
    # regardless of whether the AI Insights section has been expanded/triggered.
    camp_summary = pd.DataFrame()
    if group_col:
        _cs_agg = {c: (c, "sum") for c in ["impressions", "clicks", "spend_usd"]
                   if c in df_all.columns}
        if _cs_agg:
            camp_summary = df_all.groupby(group_col).agg(**_cs_agg).reset_index()
            if "impressions" in camp_summary.columns and "clicks" in camp_summary.columns:
                camp_summary["ctr"] = camp_summary["clicks"] / camp_summary["impressions"]
            if "impressions" in camp_summary.columns and "spend_usd" in camp_summary.columns:
                camp_summary["cpm"] = camp_summary["spend_usd"] / camp_summary["impressions"] * 1000

    st.subheader(f"AI {group_label} Insights" if group_label else "AI Insights")

    # Get the Anthropic API key — check Streamlit secrets first, then env var
    api_key = (
        st.secrets.get("ANTHROPIC_API_KEY")
        if "ANTHROPIC_API_KEY" in st.secrets
        else os.environ.get("ANTHROPIC_API_KEY")
    )

    if not api_key:
        st.warning(
            "No Anthropic API key found. Add `ANTHROPIC_API_KEY` to your "
            "Streamlit secrets or environment variables to enable AI insights."
        )
    elif group_col is None:
        st.info("No suitable grouping column found in this file. AI insights require at least one text/category column.")
    else:
        # camp_summary is pre-computed above (before this if/elif/else block)
        has_spend  = "spend_usd"   in df_all.columns
        has_clicks = "clicks"      in df_all.columns
        has_imps   = "impressions" in df_all.columns

        # Format the data as a text table for the prompt.
        # Includes campaign-level totals AND a campaign × environment breakdown
        # when an environment column is present (gives Claude the Web/App/YouTube split).
        def format_summary_table(df):
            lines = [f"=== {group_label} Totals ==="]
            for _, r in camp_summary.iterrows():
                parts = [f"{group_label}: {r[group_col]}"]
                if "impressions" in r and pd.notna(r["impressions"]):
                    parts.append(f"Impressions: {int(r['impressions']):,}")
                if "clicks" in r and pd.notna(r["clicks"]):
                    parts.append(f"Clicks: {int(r['clicks']):,}")
                if "spend_usd" in r and pd.notna(r["spend_usd"]):
                    parts.append(f"Spend: ${r['spend_usd']:,.0f}")
                if "ctr" in r and pd.notna(r["ctr"]):
                    parts.append(f"CTR: {r['ctr']:.2%}")
                if "cpm" in r and pd.notna(r["cpm"]):
                    parts.append(f"CPM: ${r['cpm']:,.2f}")
                lines.append(" | ".join(parts))

            # Helper: aggregate a dimension column and append formatted rows to lines
            def _add_breakdown(label_col, section_title, label_key):
                if label_col not in df.columns:
                    return
                lines.append(f"\n=== {section_title} ===")
                agg = {c: (c, "sum") for c in ["impressions", "clicks", "spend_usd"]
                       if c in df.columns}
                grp = df.groupby([group_col, label_col]).agg(**agg).reset_index()
                if "impressions" in grp.columns and "clicks" in grp.columns:
                    grp["ctr"] = grp["clicks"] / grp["impressions"]
                if "impressions" in grp.columns and "spend_usd" in grp.columns:
                    grp["cpm"] = grp["spend_usd"] / grp["impressions"] * 1000
                if "clicks" in grp.columns and "spend_usd" in grp.columns:
                    grp["cpc"] = grp["spend_usd"] / grp["clicks"]
                for _, r in grp.iterrows():
                    parts = [f"{group_label}: {r[group_col]} | {label_key}: {r[label_col]}"]
                    if "impressions" in r and pd.notna(r["impressions"]):
                        parts.append(f"Impressions: {int(r['impressions']):,}")
                    if "clicks" in r and pd.notna(r["clicks"]):
                        parts.append(f"Clicks: {int(r['clicks']):,}")
                    if "spend_usd" in r and pd.notna(r["spend_usd"]):
                        parts.append(f"Spend: ${r['spend_usd']:,.0f}")
                    if "ctr" in r and pd.notna(r["ctr"]):
                        parts.append(f"CTR: {r['ctr']:.2%}")
                    if "cpm" in r and pd.notna(r["cpm"]):
                        parts.append(f"CPM: ${r['cpm']:,.2f}")
                    if "cpc" in r and pd.notna(r["cpc"]):
                        parts.append(f"CPC: ${r['cpc']:,.2f}")
                    lines.append(" | ".join(parts))

            # Breakdown by Environment (Web, App, YouTube, etc.)
            _add_breakdown("environment", "Breakdown by Brand and Environment", "Environment")

            # Breakdown by Insertion Order / Line Item name
            _add_breakdown("line_item", "Breakdown by Brand and Line Item / Creative", "Line Item")

            # Breakdown by Device Type
            _add_breakdown("device_type", "Breakdown by Brand and Device Type", "Device Type")

            return "\n".join(lines)

        all_campaigns_text = format_summary_table(df_all)
        campaign_list = camp_summary[group_col].tolist()
        # Store group names in session_state so other pages (e.g. Brand Settings) can read them
        st.session_state["campaign_list"] = campaign_list

        # System prompt — analyst persona, structured output required
        SYSTEM_PROMPT = (
            "You are a senior programmatic advertising analyst with deep expertise in "
            "DSP campaign performance (DV360, TTD). You write clear, concise, data-driven "
            "commentary for ad tech professionals. Be specific — reference actual numbers. "
            "Follow the section headings and structure exactly as specified in each prompt. "
            "Do not add extra sections or deviate from the requested format."
        )

        def stream_insight(campaign_name: str, all_data: str, brand_context: str = ""):
            """
            Stream per-brand analysis using the standard Display / Video / YouTube
            insertion-order structure. Brand memory is injected as an override block
            at the end of the prompt so it takes priority over all default instructions.
            Yields text chunks from the Claude API.
            """
            client_ai = anthropic.Anthropic(api_key=api_key)

            # Default structure: Campaign Overview + three insertion-order sections.
            # The AI skips any section whose insertion order is absent from the data.
            default_structure = (
                f"**{campaign_name} - Campaign Overview**\n"
                f"Summarise overall campaign performance. Focus on total Revenue (Spend) "
                f"and total Impressions for this brand. Keep to 2-3 sentences.\n\n"
                f"**Display**\n"
                f"Summarise performance for the Display insertion order, focusing on "
                f"average CPM. Then identify:\n"
                f"- Best performing Line Item by CPM (state exact CPM, CPC, CTR)\n"
                f"- Worst performing Line Item by CPM (state exact CPM, CPC, CTR)\n"
                f"- Best performing Creative by CPM (state exact CPM, CPC, CTR)\n"
                f"- Worst performing Creative by CPM (state exact CPM, CPC, CTR)\n"
                f"If no Display insertion order data exists for this brand, omit this "
                f"section entirely — do not mention it at all.\n\n"
                f"**Video**\n"
                f"Summarise performance for the Video insertion order, focusing on CPV "
                f"and VTR. Then identify:\n"
                f"- Best performing Line Item by CPV (state exact CPV, VTR)\n"
                f"- Worst performing Line Item by CPV (state exact CPV, VTR)\n"
                f"- Best performing Creative by CPV (state exact CPV, VTR)\n"
                f"- Worst performing Creative by CPV (state exact CPV, VTR)\n"
                f"If no Video insertion order data exists for this brand, omit this "
                f"section entirely — do not mention it at all.\n\n"
                f"**YouTube**\n"
                f"Summarise performance for the YouTube insertion order, focusing on CPV "
                f"and VTR. Then identify:\n"
                f"- Best performing Line Item by CPV (state exact CPV, VTR)\n"
                f"- Worst performing Line Item by CPV (state exact CPV, VTR)\n"
                f"- Best performing Creative by CPV (state exact CPV, VTR)\n"
                f"- Worst performing Creative by CPV (state exact CPV, VTR)\n"
                f"If no YouTube insertion order data exists for this brand, omit this "
                f"section entirely — do not mention it at all."
            )

            # Brand memory override — appended after the default instructions so it
            # takes explicit priority. Only included when brand context exists.
            if brand_context:
                override_section = (
                    f"\n\nBRAND MEMORY OVERRIDE — These instructions take priority over "
                    f"all default instructions above. Where there is any conflict, always "
                    f"follow these brand-specific instructions instead:\n\n"
                    f"{brand_context}"
                )
            else:
                override_section = ""

            prompt = (
                f"Here is the full performance data for all campaigns in this report:\n\n"
                f"{all_data}\n\n"
                f"Write the analysis for the '{campaign_name}' brand using exactly "
                f"this structure and these headings (use ** for bold):\n\n"
                f"{default_structure}"
                f"{override_section}\n\n"
                f"Use only the data provided — do not invent numbers."
            )

            with client_ai.messages.stream(
                model="claude-sonnet-4-6",
                max_tokens=1200,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text in stream.text_stream:
                    yield text

        def stream_overall(all_data: str):
            """
            Stream the cross-campaign summary sections:
            Summary, Best & Worst Performers, Optimisation Recommendations.
            Yields text chunks from the Claude API.
            """
            client_ai = anthropic.Anthropic(api_key=api_key)

            prompt = (
                f"Here is the full performance data for all campaigns in this report:\n\n"
                f"{all_data}\n\n"
                f"Write the following three sections using exactly these headings "
                f"(use ** for bold):\n\n"
                f"**Summary**\n"
                f"Overall commentary across all campaigns and environments. 3-4 sentences.\n\n"
                f"**Best & Worst Performers**\n"
                f"- Best performing device type and why\n"
                f"- Worst performing device type and why\n"
                f"- Best performing creative and why\n"
                f"- Worst performing creative and why\n\n"
                f"**Optimisation Recommendations**\n"
                f"3-5 actionable recommendations based on the data above. Use bullet points.\n\n"
                f"Use only the data provided. Be specific with numbers where possible."
            )

            with client_ai.messages.stream(
                model="claude-sonnet-4-6",
                max_tokens=600,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": prompt}],
            ) as stream:
                for text in stream.text_stream:
                    yield text

        # STYLE LOCK: Do not change insights display formatting
        def _insight_html(text):
            """
            Convert the subset of markdown used in insights to explicitly-styled HTML.
            Forces ALL text to plain black (#111111) so Streamlit's theme accent color
            (teal/cyan) is never applied to bold headers or any other element.

            Handles:
              **text**  → <strong> (black, no color override from theme)
              - item    → <ul><li> list
              blank line → paragraph break
            """
            lines  = text.split("\n")
            parts  = []
            in_ul  = False
            p_style = "margin:4px 0;color:#111111;font-size:14px;line-height:1.7;"
            li_style = "margin:2px 0;color:#111111;font-size:14px;"

            def apply_bold(s):
                # Replace **…** with <strong> that explicitly inherits black
                return re.sub(r'\*\*(.+?)\*\*',
                              r'<strong style="color:#111111;">\1</strong>', s)

            for line in lines:
                stripped = line.strip()
                if not stripped:
                    if in_ul:
                        parts.append("</ul>")
                        in_ul = False
                    continue
                if stripped.startswith("- "):
                    if not in_ul:
                        parts.append("<ul style='margin:6px 0 6px 18px;padding:0;'>")
                        in_ul = True
                    parts.append(f"<li style='{li_style}'>{apply_bold(stripped[2:])}</li>")
                else:
                    if in_ul:
                        parts.append("</ul>")
                        in_ul = False
                    parts.append(f"<p style='{p_style}'>{apply_bold(stripped)}</p>")

            if in_ul:
                parts.append("</ul>")

            return (
                "<div style='font-family:-apple-system,Segoe UI,Roboto,sans-serif;"
                "color:#111111;'>"
                + "\n".join(parts)
                + "</div>"
            )

        # Generate Insights button — analysis only runs when clicked
        if st.button("✨ Generate Insights", type="primary"):
            st.session_state["insights_triggered"] = True

        if st.session_state.get("insights_triggered"):
            # Load brand memory once so it's available for every campaign
            _insights_bm = load_brand_memory()

            for campaign in campaign_list:
                # Look up brand context using partial matching (case-insensitive).
                # e.g. brand key "Nike" matches campaign "Nike Summer 2024".
                _campaign_brand_context = ""
                for _bm_key, _bm_val in _insights_bm.items():
                    if _bm_key.lower() in campaign.lower():
                        _campaign_brand_context = _bm_val.get("rationale", "")
                        break

                # One expander per campaign, collapsed by default
                with st.expander(f"**{campaign}**", expanded=False):
                    # Show indicator when brand memory is applied
                    if _campaign_brand_context:
                        st.markdown(
                            "<p style='color:#F5A623;font-size:13px;font-weight:600;"
                            "margin:0 0 8px 0;'>✓ Brand context applied</p>",
                            unsafe_allow_html=True,
                        )
                    placeholder = st.empty()
                    accumulated = ""
                    for chunk in stream_insight(campaign, all_campaigns_text,
                                                _campaign_brand_context):
                        accumulated += chunk
                        placeholder.markdown(_insight_html(accumulated),
                                             unsafe_allow_html=True)
                    # Save completed text for PPTX export
                    st.session_state.setdefault("insights_text", {})[campaign] = accumulated

            # Overall summary — runs once after all per-campaign sections
            with st.expander("**📊 Summary & Recommendations**", expanded=False):
                placeholder = st.empty()
                accumulated = ""
                for chunk in stream_overall(all_campaigns_text):
                    accumulated += chunk
                    placeholder.markdown(_insight_html(accumulated),
                                         unsafe_allow_html=True)
                st.session_state["insights_overall"] = accumulated

    # ─────────────────────────────────────────────────────────────────────────
    # AI-DRIVEN DYNAMIC ANALYSIS
    # Sends column headers + 50 sample rows to Claude to recommend 3 charts
    # ─────────────────────────────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("## 🤖 AI-Driven Analysis")

    # Load brand context for the AI analysis prompt
    def get_brand_context_for_ai(df):
        try:
            with open("brand_memory.json", "r") as f:
                bm = json.load(f)
            contexts = []
            if 'campaign' in df.columns:
                for brand in df['campaign'].dropna().unique()[:3]:
                    for key, val in bm.items():
                        if key.lower() in str(brand).lower() or str(brand).lower() in key.lower():
                            rationale = val.get("rationale", "")
                            if rationale:
                                contexts.append(f"{brand}: {rationale[:200]}")
            return "\n".join(contexts)
        except Exception:
            return ""

    # Button row: Run triggers the API call; Refresh clears the cached result
    _ai_btn_cols = st.columns([2, 1])
    with _ai_btn_cols[0]:
        _run_ai_btn = st.button("🤖 Run AI Analysis", type="primary", key="run_ai_analysis")
    with _ai_btn_cols[1]:
        if st.session_state.get("ai_analysis_result") and st.button("🔄 Refresh", key="refresh_ai_analysis"):
            st.session_state.pop("ai_analysis_result", None)
            st.rerun()

    # Only call the Claude API when the Run button is clicked
    if _run_ai_btn:
        st.session_state.pop("ai_analysis_result", None)

        columns_list = list(df_all.columns)
        sample_rows  = df_all.head(50).to_string(index=False, max_cols=20)
        brand_ctx    = get_brand_context_for_ai(df_all)

        brand_ctx_injection = ""
        if brand_ctx:
            brand_ctx_injection = (
                f"\n\nBrand context to consider: {brand_ctx}\n"
                "Weight your chart recommendations toward metrics relevant to this brand's objectives."
            )

        seed_val = random.randint(1, 9999)

        prompt = f"""You are an expert data analyst reviewing a programmatic advertising dataset. Here are the columns available: {columns_list}. Here is a sample of the data:

{sample_rows}
{brand_ctx_injection}

Return a JSON response only, no other text, with this exact structure:
{{
  "insights": [
    {{
      "title": "Chart title",
      "chart_type": "bar|line|scatter|pie",
      "x_axis": "column name",
      "y_axis": "column name",
      "rationale": "Why this chart surfaces a useful insight",
      "anomaly": "Any anomaly detected in this dimension or null"
    }}
  ],
  "summary": "One paragraph summary of the most important patterns in this dataset",
  "top_anomaly": "The single most important anomaly or finding across the whole dataset"
}}
Return exactly 3 chart recommendations. Seed: {seed_val}"""

        try:
            _ai_key = (
                st.secrets.get("ANTHROPIC_API_KEY")
                if "ANTHROPIC_API_KEY" in st.secrets
                else os.environ.get("ANTHROPIC_API_KEY", "")
            )
            client_driven = anthropic.Anthropic(api_key=_ai_key)

            with st.spinner("Running AI analysis…"):
                response = client_driven.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=2000,
                    messages=[{"role": "user", "content": prompt}]
                )

            raw = response.content[0].text.strip()
            # Strip markdown code fences if Claude wraps the JSON in them
            if raw.startswith("```"):
                raw = re.sub(r'^```[a-z]*\n?', '', raw)
                raw = re.sub(r'\n?```$', '', raw)

            result = json.loads(raw)
            st.session_state["ai_analysis_result"] = result

        except Exception as e:
            st.warning(f"AI analysis unavailable: {e}")
            st.session_state["ai_analysis_result"] = None

    # Show placeholder when no analysis has been run yet
    if "ai_analysis_result" not in st.session_state:
        st.markdown("""
        <div style="text-align:center;padding:32px 20px;background:#FFFFFF;
        border-radius:12px;box-shadow:0 2px 8px rgba(0,0,0,0.06);margin:12px 0;">
            <div style="font-size:32px;margin-bottom:10px;">🤖</div>
            <div style="font-size:15px;font-weight:600;color:#111827;margin-bottom:6px;">
                AI-Driven Chart Analysis
            </div>
            <div style="font-size:13px;color:#6B7280;">
                Click <strong>🤖 Run AI Analysis</strong> above to generate
                AI-recommended chart insights for this dataset.
            </div>
        </div>
        """, unsafe_allow_html=True)

    # Display AI analysis results (only when a result is cached in session state)
    ai_result = st.session_state.get("ai_analysis_result")
    if ai_result:
        # Summary card
        summary = ai_result.get("summary", "")
        top_anomaly = ai_result.get("top_anomaly", "")

        if summary:
            st.markdown(f"""
            <div style="background:#EEF1F4;border-left:4px solid #1B2A4A;border-radius:8px;padding:16px;margin-bottom:12px;">
                <div style="font-weight:700;color:#1B2A4A;margin-bottom:6px;">📊 Dataset Summary</div>
                <div style="color:#374151;">{summary}</div>
            </div>
            """, unsafe_allow_html=True)

        if top_anomaly:
            st.markdown(f"""
            <div style="background:#FEF3C7;border-left:4px solid #F59E0B;border-radius:8px;padding:16px;margin-bottom:20px;">
                <div style="font-weight:700;color:#92400E;margin-bottom:6px;">⚠️ Top Finding</div>
                <div style="color:#374151;">{top_anomaly}</div>
            </div>
            """, unsafe_allow_html=True)

        # Render AI-recommended charts
        insights = ai_result.get("insights", [])
        if insights:
            ai_chart_cols = st.columns(min(len(insights), 3))
            for i, insight in enumerate(insights[:3]):
                with ai_chart_cols[i % 3]:
                    title = insight.get("title", f"Chart {i+1}")
                    chart_type = insight.get("chart_type", "bar")
                    x_col = insight.get("x_axis", "")
                    y_col = insight.get("y_axis", "")
                    rationale = insight.get("rationale", "")
                    anomaly = insight.get("anomaly")

                    try:
                        # Only render if both columns exist in the dataframe
                        if x_col in df_all.columns and y_col in df_all.columns:
                            plot_df = df_all[[x_col, y_col]].dropna()

                            # Aggregate text/category x-axis columns by summing y
                            if plot_df[x_col].dtype == object or pd.api.types.is_categorical_dtype(plot_df[x_col]):
                                plot_df = plot_df.groupby(x_col)[y_col].sum().reset_index()

                            if chart_type == "bar":
                                fig = px.bar(plot_df, x=x_col, y=y_col, title=title,
                                            color_discrete_sequence=["#1B2A4A"])
                            elif chart_type == "line":
                                fig = px.line(plot_df, x=x_col, y=y_col, title=title,
                                             color_discrete_sequence=["#1B2A4A"])
                            elif chart_type == "scatter":
                                fig = px.scatter(plot_df, x=x_col, y=y_col, title=title,
                                               color_discrete_sequence=["#1B2A4A"])
                            elif chart_type == "pie":
                                fig = px.pie(plot_df, names=x_col, values=y_col, title=title)
                            else:
                                fig = px.bar(plot_df, x=x_col, y=y_col, title=title,
                                            color_discrete_sequence=["#1B2A4A"])

                            fig.update_layout(
                                paper_bgcolor="rgba(0,0,0,0)",
                                plot_bgcolor="rgba(0,0,0,0)",
                                font_family="Poppins, sans-serif",
                                height=300,
                                margin=dict(l=20, r=20, t=40, b=20),
                            )
                            st.plotly_chart(fig, use_container_width=True, config=PLOTLY_CONFIG, key=f"ai_chart_{i}")

                            # Label showing this was AI-selected
                            st.markdown(f"""
                            <div style="font-size:11px;color:#6B7280;background:#F9FAFB;border-radius:4px;padding:6px 10px;margin-top:-8px;">
                                🤖 <em>AI selected · {rationale}</em>
                            </div>
                            """, unsafe_allow_html=True)

                            if anomaly:
                                st.markdown(f"""
                                <div style="font-size:11px;color:#92400E;background:#FEF3C7;border-radius:4px;padding:4px 8px;margin-top:4px;">
                                    ⚠️ {anomaly}
                                </div>
                                """, unsafe_allow_html=True)
                        else:
                            st.info(f"Column '{x_col}' or '{y_col}' not found in data")

                    except Exception as chart_err:
                        st.warning(f"Could not render chart: {chart_err}")

    # ─────────────────────────────────────────────────────────────────────────
    # NATURAL LANGUAGE QUERYING
    # User can ask plain English questions about their data
    # ─────────────────────────────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("## 💬 Ask a Question About Your Data")

    nl_question = st.text_input(
        "Ask a question about your data",
        placeholder="e.g. Which campaign had the worst CTR last week?",
        key="nl_question_input",
        label_visibility="collapsed"
    )

    if nl_question and st.button("Ask →", key="nl_submit", type="primary"):
        # Build a brief dataset summary for the prompt
        summary_parts = []
        if 'campaign' in df_all.columns:
            summary_parts.append(f"Campaigns: {', '.join(df_all['campaign'].dropna().unique()[:10])}")
        if 'impressions' in df_all.columns:
            summary_parts.append(f"Total impressions: {df_all['impressions'].sum():,.0f}")
        if 'clicks' in df_all.columns:
            summary_parts.append(f"Total clicks: {df_all['clicks'].sum():,.0f}")
        if 'spend_usd' in df_all.columns:
            summary_parts.append(f"Total spend: ${df_all['spend_usd'].sum():,.2f}")

        # Include a sample of the raw data for context
        agg_sample = df_all.head(100).to_string(index=False, max_cols=15)

        nl_prompt = f"""You are an expert programmatic advertising analyst. The user has uploaded a campaign performance dataset.

Dataset summary:
{chr(10).join(summary_parts)}

Data sample:
{agg_sample}

User question: {nl_question}

Provide a direct, specific answer in plain English. Include the specific numbers from the data. Keep your answer to 2-3 sentences maximum."""

        try:
            _nl_key = (
                st.secrets.get("ANTHROPIC_API_KEY")
                if "ANTHROPIC_API_KEY" in st.secrets
                else os.environ.get("ANTHROPIC_API_KEY", "")
            )
            client_nl = anthropic.Anthropic(api_key=_nl_key)

            with st.spinner(""):
                response = client_nl.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=500,
                    messages=[{"role": "user", "content": nl_prompt}]
                )

            answer = response.content[0].text.strip()

            # Store in query history (keep last 5)
            if "nl_query_history" not in st.session_state:
                st.session_state["nl_query_history"] = []
            st.session_state["nl_query_history"].insert(0, {
                "question": nl_question,
                "answer": answer
            })
            st.session_state["nl_query_history"] = st.session_state["nl_query_history"][:5]

            st.session_state["nl_latest_answer"] = answer

        except Exception as e:
            st.error(f"Could not get answer: {e}")

    # Show the most recent answer
    if "nl_latest_answer" in st.session_state:
        st.markdown(f"""
        <div style="background:#F0FDF4;border-left:4px solid #10B981;border-radius:8px;padding:16px;margin:12px 0;">
            <div style="font-weight:700;color:#065F46;margin-bottom:6px;">💡 Answer</div>
            <div style="color:#374151;">{st.session_state["nl_latest_answer"]}</div>
        </div>
        """, unsafe_allow_html=True)

    # Query history accordion
    if st.session_state.get("nl_query_history"):
        with st.expander("📋 Query History (last 5)", expanded=False):
            for i, item in enumerate(st.session_state["nl_query_history"]):
                st.markdown(f"**Q{i+1}:** {item['question']}")
                st.markdown(f"**A:** {item['answer']}")
                if i < len(st.session_state["nl_query_history"]) - 1:
                    st.markdown("---")


    # ── In-page Report Export ────────────────────────────────────────────────
    st.markdown("---")

    # Determine which sections have data available for the export modal
    _rpt_key       = api_key  # same Anthropic key used for AI Insights above
    _daily_avail   = "date" in df_all.columns
    _ai_ins_avail  = bool(group_col and api_key and not camp_summary.empty)

    @st.dialog("Build Your Report")
    def _show_report_dialog():
        st.caption("Select the sections to include in your PowerPoint export.")

        # Restore previous checkbox selections from session state
        _defs = st.session_state.get("_rpt_section_defaults", {
            "summary_metrics":        True,
            "performance_charts":     True,
            "daily_performance":      True,
            "period_comparison":      True,
            "ai_advertiser_insights": False,
            "ai_driven_analysis":     False,
        })

        # Select all / Deselect all buttons
        _sc = st.columns(2)
        with _sc[0]:
            if st.button("Select all", key="rpt_sel_all"):
                st.session_state["_rpt_section_defaults"] = {k: True for k in _defs}
                st.rerun()
        with _sc[1]:
            if st.button("Deselect all", key="rpt_desel_all"):
                st.session_state["_rpt_section_defaults"] = {k: False for k in _defs}
                st.rerun()

        st.markdown("")

        # One checkbox per section — disabled when data is unavailable
        s_summary = st.checkbox(
            "Summary Metrics", value=_defs.get("summary_metrics", True),
            key="rpt_cb_summary",
        )
        s_charts = st.checkbox(
            "Performance Charts", value=_defs.get("performance_charts", True),
            key="rpt_cb_charts",
        )
        s_daily = st.checkbox(
            "Daily Performance",
            value=_defs.get("daily_performance", True) and _daily_avail,
            disabled=not _daily_avail,
            key="rpt_cb_daily",
            help="Requires a date column in your data." if not _daily_avail else None,
        )
        s_period = st.checkbox(
            "Period Comparison",
            value=_defs.get("period_comparison", True) and _period_comp_available,
            disabled=not _period_comp_available,
            key="rpt_cb_period",
            help="Requires 14+ days of date data." if not _period_comp_available else None,
        )
        s_ai_ins = st.checkbox(
            "AI Advertiser Insights",
            value=_defs.get("ai_advertiser_insights", False),
            disabled=not _ai_ins_avail,
            key="rpt_cb_ai_ins",
            help="Requires a grouping column and API key." if not _ai_ins_avail else None,
        )
        s_ai_drv = st.checkbox(
            "AI-Driven Analysis",
            value=_defs.get("ai_driven_analysis", False),
            key="rpt_cb_ai_drv",
        )
        st.caption("☝ AI sections use AI credits (API calls made at export time if not already cached).")

        st.markdown("")

        # Bottom action buttons
        _bc = st.columns(2)
        with _bc[0]:
            if st.button("Generate PowerPoint", type="primary", key="rpt_gen_btn"):
                # Persist selections for next time the dialog is opened
                st.session_state["_rpt_section_defaults"] = {
                    "summary_metrics":        s_summary,
                    "performance_charts":     s_charts,
                    "daily_performance":      s_daily,
                    "period_comparison":      s_period,
                    "ai_advertiser_insights": s_ai_ins,
                    "ai_driven_analysis":     s_ai_drv,
                }
                _sections = st.session_state["_rpt_section_defaults"]
                if not any(_sections.values()):
                    st.warning("Select at least one section.")
                elif not _rpt_key:
                    st.warning("No API key found. Add ANTHROPIC_API_KEY to Streamlit secrets.")
                elif camp_summary.empty:
                    st.warning("No data loaded — upload a CSV file first.")
                else:
                    with st.spinner("Building PowerPoint… this may take 10–30 seconds."):
                        try:
                            pptx_buf = build_pptx_report(
                                _rpt_key,
                                camp_summary,
                                df_all,
                                sections=_sections,
                                ai_analysis_cache=st.session_state.get("ai_analysis_result"),
                            )
                            st.session_state["pptx_report"] = pptx_buf
                        except Exception as _e:
                            st.error(f"Report generation failed: {_e}")
                    st.rerun()
        with _bc[1]:
            if st.button("Cancel", key="rpt_cancel_btn"):
                st.rerun()

    # Centre the Generate Report button
    _, _rpt_col, _ = st.columns([2, 2, 2])
    with _rpt_col:
        if st.button("📊 Generate Report", type="primary", key="open_report_dialog_btn"):
            _show_report_dialog()

    # Show download button and success message after a report has been generated
    if st.session_state.get("pptx_report"):
        st.success("✅ Report ready — click below to download.")
        st.download_button(
            label     = "📥 Download Report (.pptx)",
            data      = st.session_state["pptx_report"],
            file_name = f"{date.today():%Y-%m-%d}_campaign_report.pptx",
            mime      = "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key       = "download_report_inline",
        )
