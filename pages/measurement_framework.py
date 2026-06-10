import io
import os
import json
from datetime import date, timedelta
import pandas as pd
import streamlit as st
import plotly.graph_objects as go
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# STYLE LOCK: Do not remove or modify this CSS block.
st.markdown("""
<style>
    html, body, [class*="css"] {
        font-family: "Inter", system-ui, -apple-system, "Segoe UI", sans-serif;
        font-size: 15px;
        color: #374151;
    }
    .stApp { background-color: #F3F4F6; }
    h1, h2, h3, h4, h5, h6 { font-weight: 700 !important; color: #111827 !important; }
    h2, h3 {
        margin-top: 2rem !important;
        padding-top: 0.25rem !important;
        border-bottom: none !important;
        padding-bottom: 0 !important;
        border-left: none !important;
        padding-left: 0 !important;
    }
    [data-testid="metric-container"] {
        background: #FFFFFF;
        border: none;
        border-radius: 16px;
        padding: 20px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        border-top: 4px solid #7C3AED;
    }
    [data-testid="metric-container"] label {
        font-size: 12px;
        font-weight: 600;
        color: #6B7280;
        text-transform: uppercase;
        letter-spacing: 0.05em;
    }
    [data-testid="metric-container"] [data-testid="stMetricValue"] {
        font-size: 26px;
        font-weight: 700;
        color: #111827;
    }
    .element-container:has([data-testid="stPlotlyChart"]) {
        background: #FFFFFF;
        border-radius: 12px;
        padding: 16px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
    }
    .stButton > button[kind="primary"],
    [data-testid="baseButton-primary"] {
        background-color: #7C3AED !important;
        color: #FFFFFF !important;
        border: none !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        font-size: 15px !important;
        padding: 0.5rem 1.25rem !important;
    }
    .stButton > button[kind="primary"]:hover,
    [data-testid="baseButton-primary"]:hover { background-color: #6D28D9 !important; }
</style>
""", unsafe_allow_html=True)
# STYLE LOCK

# ── Brand memory helpers ──────────────────────────────────────────────────────
_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BRAND_MEMORY_PATH = os.path.join(_ROOT, "brand_memory.json")


def load_brand_memory():
    if os.path.exists(BRAND_MEMORY_PATH):
        with open(BRAND_MEMORY_PATH, "r") as f:
            return json.load(f)
    return {}


# ── Mock campaign data ────────────────────────────────────────────────────────
# All values are representative of a mid-size pDOOH retail media campaign.
MOCK = {
    "advertiser":          "Sample Advertiser",
    "campaign_name":       "Brand Awareness — Q2 2026",
    "centres":             ["Melbourne Central", "Macquarie Centre", "Rouse Hill Town Centre", "Highpoint"],
    "flight_start":        date(2026, 5, 1),
    "flight_end":          date(2026, 7, 31),
    "total_budget":        42_000,
    "impressions":         1_500_000,
    "frequency":           2.0,
    "share_of_voice":      35.0,       # % of total available impressions
    "avg_cpm":             28.00,
    "fill_rate":           94.2,
    "total_spend":         42_000,
    # Delivery by centre
    "by_centre": {
        "Melbourne Central":       600_000,
        "Macquarie Centre":        350_000,
        "Rouse Hill Town Centre":  320_000,
        "Highpoint":               230_000,
    },
    # Delivery by screen type
    "by_screen": {
        "Large Format":  0.40,
        "Food Court":    0.30,
        "Corridor":      0.20,
        "Entry Panels":  0.10,
    },
    # Footfall attribution (modelled)
    "footfall_rate":       22.0,   # % of ad-exposed who entered a store
    "footfall_uplift":     8.0,    # uplift vs non-exposed control
    "footfall_categories": [
        ("Fashion",           12),
        ("Food & Beverage",    9),
        ("Electronics",        7),
    ],
}

# Internal report additional mock values
INTERNAL = {
    "monthly_target":        50_000,
    "revenue_per_screen_day": 70,   # A$ per screen per day
    "avg_clearing_cpm":       29.50,
    "total_available_imps":   1_594_915,  # impressions available in period
    # Revenue per centre (for internal report)
    "rev_by_centre": {
        "Melbourne Central":       18_000,
        "Macquarie Centre":        10_500,
        "Rouse Hill Town Centre":   9_000,
        "Highpoint":                4_500,
    },
    # Daily fill rate over campaign (simplified 7-day sample for chart)
    "daily_fill_rate": {
        "2026-05-01": 88.0, "2026-05-08": 91.5, "2026-05-15": 93.0,
        "2026-05-22": 94.8, "2026-06-01": 95.2, "2026-06-08": 94.2,
        "2026-06-15": 93.7,
    },
    # Deal type revenue split per centre
    "deal_type_rev": {
        "Melbourne Central":       {"Open": 2_000, "PMP": 8_000, "PG": 8_000},
        "Macquarie Centre":        {"Open": 500,   "PMP": 6_500, "PG": 3_500},
        "Rouse Hill Town Centre":  {"Open": 0,     "PMP": 5_000, "PG": 4_000},
        "Highpoint":               {"Open": 1_000, "PMP": 2_000, "PG": 1_500},
    },
}


# ── PPTX builder (shared dark premium template) ────────────────────────────
def build_pptx(report_type, advertiser, campaign_name, centres,
               flight_start, flight_end, impressions, avg_cpm,
               fill_rate, ai_commentary=""):
    """
    Build a dark navy PPTX report (agency or internal) using python-pptx from scratch.
    Returns a BytesIO buffer.
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    LGREY  = RGBColor(0xA8, 0xB2, 0xBC)
    BLUE   = RGBColor(0x00, 0xA8, 0xE8)

    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY

    def txt(slide, text, l, t, w, h, size=14, bold=False,
            color=WHITE, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        r  = p.add_run()
        r.text           = text
        r.font.size      = Pt(size)
        r.font.bold      = bold
        r.font.color.rgb = color
        r.font.name      = "Calibri"

    def footer(slide, n):
        txt(slide, "GPT Group — Measurement Framework", 0.2, 7.1, 7, 0.3, size=9, color=LGREY)
        txt(slide, str(n), 12.9, 7.1, 0.3, 0.3, size=9, color=LGREY, align=PP_ALIGN.RIGHT)

    centres_str  = ", ".join(centres) if centres else "Multiple Centres"
    flight_str   = (
        f"{flight_start.strftime('%d %b %Y')} – {flight_end.strftime('%d %b %Y')}"
        if flight_start and flight_end else "TBC"
    )
    report_label = "Agency Campaign Report" if report_type == "agency" else "Internal Yield Report"

    # Slide 1 — Cover
    s1 = add_slide(); set_bg(s1)
    txt(s1, "GPT GROUP",       0.5, 0.4, 12, 0.5, size=11, color=BLUE, bold=True)
    txt(s1, report_label,      0.5, 1.0, 12, 0.8, size=34, bold=True)
    txt(s1, advertiser,        0.5, 1.9, 12, 0.5, size=22, color=LGREY)
    txt(s1, campaign_name,     0.5, 2.5, 12, 0.4, size=16, color=LGREY)
    txt(s1, f"Flight: {flight_str}   ·   Centres: {centres_str}",
        0.5, 3.0, 12, 0.4, size=13, color=LGREY)
    footer(s1, 1)

    # Slide 2 — KPI Summary
    s2 = add_slide(); set_bg(s2)
    txt(s2, "Campaign KPI Summary", 0.5, 0.3, 12, 0.5, size=20, bold=True)
    kpis = [
        f"Total Impressions: {impressions:,.0f}",
        f"Estimated Reach: {impressions/2:,.0f} unique audiences",
        f"Average Frequency: 2.0x",
        f"Average CPM: A${avg_cpm:.2f}",
        f"Fill Rate: {fill_rate:.1f}%",
        f"Total Spend: A${impressions * avg_cpm / 1000:,.0f}",
    ]
    body = "\n".join(f"• {k}" for k in kpis)
    txt(s2, body, 0.5, 1.1, 12.3, 5.5, size=14, color=LGREY)
    footer(s2, 2)

    # Slide 3 — AI Commentary (if provided)
    if ai_commentary:
        s3 = add_slide(); set_bg(s3)
        txt(s3, "Campaign Commentary", 0.5, 0.3, 12, 0.5, size=20, bold=True)
        txt(s3, ai_commentary[:1000], 0.5, 1.1, 12.3, 5.5, size=12, color=LGREY)
        footer(s3, 3)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Page header ───────────────────────────────────────────────────────────────
st.title("Measurement Framework")
st.markdown(
    "<p style='color:#6b7280;font-size:14px;margin-top:-12px;'>"
    "Post-campaign reporting for agency clients and internal GPT stakeholders."
    "</p>",
    unsafe_allow_html=True,
)

# ── Report type toggle ────────────────────────────────────────────────────────
report_type = st.radio(
    "Report view",
    ["📋 Agency Report", "📊 Internal Report"],
    horizontal=True,
    label_visibility="collapsed",
)
is_agency = report_type == "📋 Agency Report"

st.markdown("---")

# ══════════════════════════════════════════════════════════════════════════════
# AGENCY REPORT VIEW
# ══════════════════════════════════════════════════════════════════════════════
if is_agency:

    # ── Section 1: Campaign Summary ───────────────────────────────────────────
    st.subheader("Campaign Summary")
    a1, a2 = st.columns(2)
    with a1:
        advertiser   = st.text_input("Advertiser",    value=MOCK["advertiser"],    key="ag_adv")
        campaign_nm  = st.text_input("Campaign name", value=MOCK["campaign_name"], key="ag_camp")
        centres_disp = st.text_input(
            "Centre(s) activated",
            value=", ".join(MOCK["centres"]),
            key="ag_centres",
        )
    with a2:
        ag_start     = st.date_input("Flight start", value=MOCK["flight_start"], key="ag_start")
        ag_end       = st.date_input("Flight end",   value=MOCK["flight_end"],   key="ag_end")
        ag_budget    = st.number_input(
            "Total budget (A$)", value=MOCK["total_budget"],
            min_value=0, step=1_000, key="ag_budget",
        )

    # ── Section 2: Campaign KPI Scorecard ─────────────────────────────────────
    st.subheader("Campaign KPI Scorecard")

    reach     = int(MOCK["impressions"] / MOCK["frequency"])
    total_avail = MOCK["impressions"] / (MOCK["share_of_voice"] / 100)

    k1, k2, k3, k4 = st.columns(4)
    k1.metric("Total Impressions",    f"{MOCK['impressions']:,.0f}")
    k2.metric("OTS",                  f"{MOCK['impressions']:,.0f}",
              help="Opportunity to See = impressions for pDOOH")
    k3.metric("Estimated Reach",      f"{reach:,.0f}",
              help="Impressions ÷ Average Frequency")
    k4.metric("Average Frequency",    f"{MOCK['frequency']:.1f}x")

    k5, k6, k7, k8 = st.columns(4)
    k5.metric("Share of Voice",       f"{MOCK['share_of_voice']:.0f}%",
              help="Impressions ÷ Total Available Impressions")
    k6.metric("Average CPM",          f"A${MOCK['avg_cpm']:.2f}")
    k7.metric("Total Spend",          f"A${MOCK['total_spend']:,.0f}")
    k8.metric("Fill Rate",            f"{MOCK['fill_rate']:.1f}%")

    # ── Section 3: Delivery by Centre ─────────────────────────────────────────
    st.subheader("Delivery by Centre")

    centres_list = list(MOCK["by_centre"].keys())
    imps_list    = list(MOCK["by_centre"].values())

    centre_fig = go.Figure(go.Bar(
        x=centres_list,
        y=imps_list,
        marker_color="#2563EB",
        text=[f"{v/1000:.0f}k" for v in imps_list],
        textposition="outside",
        textfont=dict(size=11, family="Inter, sans-serif"),
        hovertemplate="%{x}<br>Impressions: %{y:,.0f}<extra></extra>",
    ))
    centre_fig.update_layout(
        yaxis_title="Impressions Delivered",
        height=340,
        margin=dict(t=20, b=80, l=60, r=20),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(family="Inter, sans-serif"),
        xaxis=dict(gridcolor="#F3F4F6"),
        yaxis=dict(gridcolor="#F3F4F6", zeroline=False),
    )
    st.plotly_chart(centre_fig, use_container_width=True)

    # ── Section 4: Delivery by Screen Type ────────────────────────────────────
    st.subheader("Delivery by Screen Type")

    screen_labels = list(MOCK["by_screen"].keys())
    screen_vals   = [v * MOCK["impressions"] for v in MOCK["by_screen"].values()]

    screen_fig = go.Figure(go.Pie(
        labels=screen_labels,
        values=screen_vals,
        hole=0.45,
        marker=dict(colors=["#2563EB", "#7C3AED", "#10B981", "#F59E0B"]),
        textinfo="label+percent",
        textfont=dict(family="Inter, sans-serif", size=13),
        hovertemplate="%{label}: %{value:,.0f} impressions (%{percent})<extra></extra>",
    ))
    screen_fig.update_layout(
        showlegend=True,
        legend=dict(orientation="h", yanchor="bottom", y=-0.2, xanchor="center", x=0.5),
        margin=dict(t=20, b=50, l=20, r=20),
        height=360,
        paper_bgcolor="rgba(0,0,0,0)",
        font=dict(family="Inter, sans-serif"),
    )
    st.plotly_chart(screen_fig, use_container_width=True)

    # ── Section 5: Footfall Attribution ───────────────────────────────────────
    st.subheader("Footfall Attribution")

    st.markdown(
        "<div style='background:#EFF6FF;border:1px solid #BFDBFE;border-radius:8px;"
        "padding:10px 14px;font-size:12px;color:#1E40AF;margin-bottom:12px;'>"
        "📍 Footfall attribution powered by mobile location data (Quantium/Near). "
        "Based on modelled audience matching."
        "</div>",
        unsafe_allow_html=True,
    )

    fa1, fa2, fa3 = st.columns(3)
    fa1.metric(
        "Ad-Exposed Store Entry Rate",
        f"{MOCK['footfall_rate']:.0f}%",
        help="% of ad-exposed audiences who entered a GPT centre store",
    )
    fa2.metric(
        "Uplift vs Control Group",
        f"+{MOCK['footfall_uplift']:.0f}%",
        help="Uplift vs non-exposed control group",
    )
    fa3.metric("Measurement Method", "Mobile Location")

    st.markdown("**Top Converting Retailer Categories Post-Exposure:**")
    for cat, uplift in MOCK["footfall_categories"]:
        st.markdown(
            f"<div style='display:inline-block;background:#F5F3FF;color:#7C3AED;"
            f"border-radius:6px;padding:4px 12px;margin:3px 4px;"
            f"font-size:13px;font-weight:600;'>"
            f"{cat} +{uplift}%</div>",
            unsafe_allow_html=True,
        )
    st.markdown("")  # spacing

    # ── Section 6: AI Campaign Commentary ─────────────────────────────────────
    st.subheader("AI Campaign Commentary")

    api_key = (
        st.secrets.get("ANTHROPIC_API_KEY")
        if "ANTHROPIC_API_KEY" in st.secrets
        else os.environ.get("ANTHROPIC_API_KEY")
    )

    if not api_key:
        st.warning(
            "Add `ANTHROPIC_API_KEY` to Streamlit secrets to enable AI commentary."
        )
    else:
        if st.button("✨ Generate Campaign Summary", type="primary", key="ag_ai_btn"):
            centres_for_prompt = centres_disp or ", ".join(MOCK["centres"])
            flight_days = (ag_end - ag_start).days if ag_end > ag_start else 91

            prompt = (
                f"You are writing a post-campaign report for {advertiser} who ran a pDOOH "
                f"campaign across {centres_for_prompt} at GPT Group retail centres. "
                f"Campaign delivered {MOCK['impressions']:,.0f} impressions at "
                f"A${MOCK['avg_cpm']:.2f} CPM with {MOCK['fill_rate']:.1f}% fill rate "
                f"over {flight_days} days. Total spend: A${MOCK['total_spend']:,.0f}. "
                f"Estimated reach: {reach:,.0f} unique audiences at {MOCK['frequency']:.1f}x frequency. "
                f"Footfall uplift vs control: +{MOCK['footfall_uplift']:.0f}%.\n\n"
                f"Write a professional 2-paragraph campaign summary suitable for sending "
                f"to an agency. Focus on delivery highlights, audience reach, and one "
                f"optimisation recommendation for the next campaign."
            )

            # Check brand memory for extra context
            _bm = load_brand_memory()
            for key, val in _bm.items():
                if key.lower() in advertiser.lower() or advertiser.lower() in key.lower():
                    ctx = val.get("rationale", "")
                    if ctx:
                        prompt += (
                            f"\n\nBRAND MEMORY OVERRIDE — These instructions take priority. "
                            f"Where there is any conflict, follow these instead:\n\n{ctx}"
                        )
                    break

            ai_client = anthropic.Anthropic(api_key=api_key)
            with st.spinner("Generating campaign summary…"):
                msg = ai_client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=500,
                    system=(
                        "You are a senior media account manager at GPT Group writing "
                        "professional post-campaign reports for agency clients. "
                        "Be concise, positive, and data-driven."
                    ),
                    messages=[{"role": "user", "content": prompt}],
                )
            st.session_state["mf_ag_ai_text"] = msg.content[0].text.strip()

        if st.session_state.get("mf_ag_ai_text"):
            st.markdown(
                "<div style='background:#FFFFFF;border-radius:12px;padding:20px 24px;"
                "box-shadow:0 4px 12px rgba(0,0,0,0.08);margin-top:12px;font-size:14px;"
                "line-height:1.7;color:#374151;'>"
                + st.session_state["mf_ag_ai_text"].replace("\n", "<br>") +
                "</div>",
                unsafe_allow_html=True,
            )

    # ── Section 7: Export ──────────────────────────────────────────────────────
    st.subheader("Export")

    pptx_buf = build_pptx(
        report_type  = "agency",
        advertiser   = advertiser,
        campaign_name= campaign_nm,
        centres      = [c.strip() for c in centres_disp.split(",") if c.strip()],
        flight_start = ag_start,
        flight_end   = ag_end,
        impressions  = MOCK["impressions"],
        avg_cpm      = MOCK["avg_cpm"],
        fill_rate    = MOCK["fill_rate"],
        ai_commentary= st.session_state.get("mf_ag_ai_text", ""),
    )
    slug = (advertiser or "Campaign").replace(" ", "_")
    st.download_button(
        label="📥 Download Agency Report (.pptx)",
        data=pptx_buf,
        file_name=f"{slug}_Agency_Report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ══════════════════════════════════════════════════════════════════════════════
# INTERNAL REPORT VIEW
# ══════════════════════════════════════════════════════════════════════════════
else:

    # ── Section 1: Yield Summary Cards ────────────────────────────────────────
    st.subheader("Yield Summary")

    total_rev      = sum(INTERNAL["rev_by_centre"].values())
    unsold_imps    = INTERNAL["total_available_imps"] - MOCK["impressions"]
    unsold_pct     = unsold_imps / INTERNAL["total_available_imps"] * 100
    avg_floor_cpm  = 22.0   # blended floor CPM across portfolio
    unsold_value   = unsold_imps * avg_floor_cpm / 1000
    rev_vs_target  = total_rev / INTERNAL["monthly_target"] * 100

    y1, y2, y3 = st.columns(3)
    y1.metric("Total Revenue",             f"A${total_rev:,.0f}")
    y2.metric("Revenue / Screen / Day",    f"A${INTERNAL['revenue_per_screen_day']:.0f}")
    y3.metric("Average Fill Rate",         f"{MOCK['fill_rate']:.1f}%")

    y4, y5, y6 = st.columns(3)
    y4.metric("Average Clearing CPM",      f"A${INTERNAL['avg_clearing_cpm']:.2f}")
    y5.metric(
        "Unsold Impressions",
        f"{unsold_pct:.1f}%",
        delta=f"-A${unsold_value:,.0f} unrealised",
        delta_color="inverse",
        help="Estimated revenue left on the table at floor CPM",
    )
    y6.metric(
        "Revenue vs Target",
        f"{rev_vs_target:.0f}%",
        delta=f"A${total_rev - INTERNAL['monthly_target']:,.0f} vs A${INTERNAL['monthly_target']:,.0f} target",
        delta_color="normal",
    )

    # ── Section 2: Revenue by Centre ──────────────────────────────────────────
    st.subheader("Revenue by Centre")

    centres_rev = list(INTERNAL["rev_by_centre"].keys())
    rev_vals    = list(INTERNAL["rev_by_centre"].values())
    # Mock target per centre (proportional to visitor count)
    rev_targets = [18_000, 12_000, 10_000, 5_000]

    rev_fig = go.Figure()
    rev_fig.add_trace(go.Bar(
        name="Actual Revenue",
        x=centres_rev,
        y=rev_vals,
        marker_color="#2563EB",
        text=[f"A${v:,.0f}" for v in rev_vals],
        textposition="outside",
        textfont=dict(size=11, family="Inter, sans-serif"),
    ))
    rev_fig.add_trace(go.Scatter(
        name="Target",
        x=centres_rev,
        y=rev_targets,
        mode="markers",
        marker=dict(symbol="line-ew", size=18, color="#EF4444", line=dict(width=3, color="#EF4444")),
    ))
    rev_fig.update_layout(
        yaxis_title="Revenue (A$)",
        height=360,
        margin=dict(t=20, b=80, l=70, r=20),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        legend=dict(orientation="h", yanchor="bottom", y=-0.25, xanchor="center", x=0.5),
        font=dict(family="Inter, sans-serif"),
        xaxis=dict(gridcolor="#F3F4F6"),
        yaxis=dict(gridcolor="#F3F4F6", zeroline=False),
    )
    st.plotly_chart(rev_fig, use_container_width=True)

    # ── Section 3: Fill Rate Trend ─────────────────────────────────────────────
    st.subheader("Fill Rate Trend")

    dates_trend = list(INTERNAL["daily_fill_rate"].keys())
    fills_trend = list(INTERNAL["daily_fill_rate"].values())

    fill_fig = go.Figure()
    fill_fig.add_trace(go.Scatter(
        x=dates_trend,
        y=fills_trend,
        mode="lines+markers",
        name="Fill Rate %",
        line=dict(color="#2563EB", width=2),
        marker=dict(size=7),
        hovertemplate="%{x}<br>Fill Rate: %{y:.1f}%<extra></extra>",
    ))
    fill_fig.add_hline(
        y=90, line_dash="dash", line_color="#7C3AED", line_width=2,
        annotation_text="90% target", annotation_position="top right",
        annotation_font_color="#7C3AED",
    )
    fill_fig.update_layout(
        yaxis_title="Fill Rate %",
        yaxis_range=[80, 100],
        height=300,
        margin=dict(t=20, b=60, l=60, r=20),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(family="Inter, sans-serif"),
        xaxis=dict(gridcolor="#F3F4F6"),
        yaxis=dict(gridcolor="#F3F4F6"),
    )
    st.plotly_chart(fill_fig, use_container_width=True)

    # ── Section 4: Deal Type Revenue Mix per Centre (stacked bar) ─────────────
    st.subheader("Deal Type Revenue Mix by Centre")

    dt_centres = list(INTERNAL["deal_type_rev"].keys())
    open_vals  = [INTERNAL["deal_type_rev"][c].get("Open", 0) for c in dt_centres]
    pmp_vals   = [INTERNAL["deal_type_rev"][c].get("PMP",  0) for c in dt_centres]
    pg_vals    = [INTERNAL["deal_type_rev"][c].get("PG",   0) for c in dt_centres]

    dtmix_fig = go.Figure()
    dtmix_fig.add_trace(go.Bar(name="Open", x=dt_centres, y=open_vals,
                               marker_color="#6B7280"))
    dtmix_fig.add_trace(go.Bar(name="PMP",  x=dt_centres, y=pmp_vals,
                               marker_color="#2563EB"))
    dtmix_fig.add_trace(go.Bar(name="PG",   x=dt_centres, y=pg_vals,
                               marker_color="#7C3AED"))
    dtmix_fig.update_layout(
        barmode="stack",
        yaxis_title="Revenue (A$)",
        height=360,
        margin=dict(t=20, b=80, l=70, r=20),
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        legend=dict(orientation="h", yanchor="bottom", y=-0.25, xanchor="center", x=0.5),
        font=dict(family="Inter, sans-serif"),
        xaxis=dict(gridcolor="#F3F4F6"),
        yaxis=dict(gridcolor="#F3F4F6", zeroline=False),
    )
    st.plotly_chart(dtmix_fig, use_container_width=True)

    # ── Section 5: Unsold Inventory Alert ─────────────────────────────────────
    st.subheader("Unsold Inventory Alert")

    # Flag centres where fill rate is below 85%
    # For this mock, we derive per-centre fill rates from impressions vs available
    centre_fill_rates = {
        "Melbourne Central":       96.5,
        "Macquarie Centre":        94.8,
        "Rouse Hill Town Centre":  93.1,
        "Highpoint":               82.4,   # below 85 — will trigger alert
    }
    centre_available_imps = {
        "Melbourne Central":       621_244,
        "Macquarie Centre":        368_421,
        "Rouse Hill Town Centre":  343_954,
        "Highpoint":               261_296,
    }
    centre_floor_cpms = {
        "Melbourne Central":       28,
        "Macquarie Centre":        25,
        "Rouse Hill Town Centre":  24,
        "Highpoint":               22,
    }

    alert_shown = False
    for centre, fr in centre_fill_rates.items():
        if fr < 85:
            avail   = centre_available_imps.get(centre, 0)
            floor   = centre_floor_cpms.get(centre, 20)
            unsold  = avail * (1 - fr / 100)
            unreal  = unsold * floor / 1000
            st.markdown(
                f"<div style='background:#FEF2F2;border:1px solid #FECACA;border-radius:8px;"
                f"padding:14px 18px;margin-bottom:8px;font-size:14px;'>"
                f"⚠️ <strong>{centre}</strong> — {100 - fr:.1f}% unsold — "
                f"estimated <strong>A${unreal:,.0f}</strong> unrealised revenue "
                f"at floor CPM of A${floor}"
                f"</div>",
                unsafe_allow_html=True,
            )
            alert_shown = True

    if not alert_shown:
        st.markdown(
            "<div style='background:#ECFDF5;border:1px solid #6EE7B7;border-radius:8px;"
            "padding:14px 18px;font-size:14px;color:#065F46;'>"
            "✅ All centres are above the 85% fill rate threshold."
            "</div>",
            unsafe_allow_html=True,
        )

    # ── Section 6: Export ──────────────────────────────────────────────────────
    st.subheader("Export")

    int_pptx = build_pptx(
        report_type  = "internal",
        advertiser   = "GPT Group",
        campaign_name= "Internal Yield Report — Q2 2026",
        centres      = list(MOCK["by_centre"].keys()),
        flight_start = MOCK["flight_start"],
        flight_end   = MOCK["flight_end"],
        impressions  = MOCK["impressions"],
        avg_cpm      = INTERNAL["avg_clearing_cpm"],
        fill_rate    = MOCK["fill_rate"],
    )
    st.download_button(
        label="📥 Download Internal Report (.pptx)",
        data=int_pptx,
        file_name="GPT_Internal_Yield_Report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

print("Measurement Framework page loaded.")
