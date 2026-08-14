import streamlit as st
import json
import os
import sys
import io
from datetime import date
import numpy as np

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
from utils.design_system import (
    get_css, PRIMARY, SECONDARY, SUCCESS, WARNING, DANGER,
    WHITE, TEXT_SEC, TEXT_PRI, BORDER_LIGHT, BG_PAGE,
)

# ── Apply Pacebird design system CSS ─────────────────────────────────────────
# STYLE LOCK: primary #F5A623 orange, secondary #1B2A4A navy, font Poppins.
st.markdown(get_css(), unsafe_allow_html=True)

# ── Page-specific CSS ─────────────────────────────────────────────────────────
st.markdown(f"""
<style>
    /* ── Segment cards ─────────────────────────────────────────────── */
    .seg-card {{
        background: {WHITE};
        border-radius: 16px;
        border-top: 4px solid {PRIMARY};
        padding: 16px 20px 14px 20px;
        margin-bottom: 10px;
        box-shadow: 0 2px 12px rgba(0,0,0,0.06);
    }}
    .seg-card-selected {{
        border-top: 4px solid {SUCCESS};
        box-shadow: 0 2px 18px rgba(16,185,129,0.15);
    }}
    /* ── Badges ─────────────────────────────────────────────────────── */
    .badge-cat {{
        display: inline-block;
        padding: 2px 10px;
        border-radius: 20px;
        font-size: 11px;
        font-weight: 700;
        background: rgba(245,166,35,0.14);
        color: {PRIMARY};
    }}
    .badge-match {{
        display: inline-block;
        padding: 2px 10px;
        border-radius: 20px;
        font-size: 11px;
        font-weight: 700;
    }}
    .match-high {{ background: rgba(16,185,129,0.14);  color: {SUCCESS}; }}
    .match-mid  {{ background: rgba(245,158,11,0.14);  color: {WARNING}; }}
    .match-low  {{ background: rgba(107,114,128,0.10); color: {TEXT_SEC}; }}
    .badge-signal {{
        display: inline-block;
        padding: 2px 8px;
        border-radius: 12px;
        font-size: 10px;
        font-weight: 600;
        background: rgba(27,42,74,0.08);
        color: {SECONDARY};
    }}
    .badge-act {{
        display: inline-block;
        padding: 2px 8px;
        border-radius: 12px;
        font-size: 10px;
        font-weight: 600;
        background: {BG_PAGE};
        color: {TEXT_SEC};
        margin-right: 3px;
    }}
    /* ── Meta labels and values inside cards ────────────────────────── */
    .mlabel {{
        font-size: 10px;
        color: {TEXT_SEC};
        text-transform: uppercase;
        letter-spacing: 0.06em;
        font-weight: 700;
        margin-bottom: 2px;
    }}
    .mval {{
        font-size: 14px;
        font-weight: 700;
        color: {TEXT_PRI};
    }}
    /* ── Stack and deal panels ──────────────────────────────────────── */
    .stack-panel {{
        background: {WHITE};
        border-radius: 16px;
        border-top: 4px solid {SECONDARY};
        padding: 20px 24px;
        box-shadow: 0 2px 12px rgba(0,0,0,0.06);
        margin-bottom: 12px;
    }}
    .deal-panel {{
        background: {WHITE};
        border-radius: 16px;
        border-top: 4px solid {PRIMARY};
        padding: 20px 24px;
        box-shadow: 0 2px 12px rgba(0,0,0,0.06);
    }}
    /* ── Warning callouts ───────────────────────────────────────────── */
    .warn-box {{
        background: #FFFBEB;
        border: 1px solid {WARNING};
        border-left: 4px solid {WARNING};
        border-radius: 8px;
        padding: 10px 14px;
        font-size: 13px;
        color: #92400E;
        margin-top: 8px;
        line-height: 1.6;
    }}
    /* ── Sense-check block ──────────────────────────────────────────── */
    .sense-check {{
        background: {BG_PAGE};
        border-radius: 8px;
        padding: 10px 14px;
        font-size: 13px;
        color: {TEXT_PRI};
        margin-top: 14px;
    }}
    /* ── Section dividers ───────────────────────────────────────────── */
    .sec-rule {{
        border: none;
        border-top: 2px solid {BORDER_LIGHT};
        margin: 28px 0 20px 0;
    }}
    /* ── Disclaimer bar ─────────────────────────────────────────────── */
    .disclaimer-bar {{
        background: #FFF4E0;
        border: 1px solid {WARNING};
        border-radius: 8px;
        padding: 6px 14px;
        font-size: 12px;
        color: #92400E;
        margin-bottom: 16px;
    }}
</style>
""", unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════════════════════════
# DATA & MODEL HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

@st.cache_data
def load_segments():
    """Load the Are Media audience segment taxonomy from JSON."""
    seg_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "..", "data", "audience_segments.json"
    )
    with open(seg_path) as f:
        return json.load(f)["segments"]


@st.cache_resource(show_spinner="Loading semantic model (first run only)...")
def load_model():
    """Load all-MiniLM-L6-v2. Cached globally — only downloads once."""
    from sentence_transformers import SentenceTransformer
    return SentenceTransformer("all-MiniLM-L6-v2")


@st.cache_data(show_spinner=False)
def embed_descriptions(desc_tuple):
    """
    Embed all segment descriptions once and cache the result.
    Takes a tuple so st.cache_data can hash it reliably.
    """
    model = load_model()
    return model.encode(list(desc_tuple))


def cosine_sim(a, b):
    """Cosine similarity between two numpy-compatible vectors."""
    a = np.array(a, dtype=float)
    b = np.array(b, dtype=float)
    denom = np.linalg.norm(a) * np.linalg.norm(b)
    return float(np.dot(a, b) / (denom + 1e-10))


def match_badge_class(score):
    """Return CSS class for the match % badge based on score."""
    if score >= 0.50:
        return "badge-match match-high"
    if score >= 0.30:
        return "badge-match match-mid"
    return "badge-match match-low"


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX BUILDER
# ═══════════════════════════════════════════════════════════════════════════════

def _build_proposal_pptx(ctx, sections):
    """
    Build a PowerPoint proposal deck from the stack context dict.
    Dark premium template — same visual language as QBR Generator.

    ctx keys: selected_segs, brief_text, budget, flight_start, flight_end,
              flight_days, objective, est_reach, use_and, wt_cpm_min,
              wt_cpm_max, max_min_spend, wt_index, deal_type,
              deal_rationale, floor_cpm, freq_cap, categories,
              signal_types, sense_check
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    # Dark premium palette
    BG_    = rgb("#0D1B2A")
    WHITE_ = rgb("#FFFFFF")
    GREY_  = rgb("#A8B2BC")
    ORANGE = rgb("#F5A623")
    NAVY_  = rgb("#1B2A4A")
    GREEN_ = rgb("#10B981")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]   # blank layout

    def new_slide():
        """Add a new dark-background slide."""
        sl = prs.slides.add_slide(blank)
        bg = sl.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_
        return sl

    def txt(sl, text, l, t, w, h,
            size=12, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
        """Add a text box to a slide."""
        tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(text)
        r.font.size   = Pt(size)
        r.font.bold   = bold
        r.font.italic = italic
        r.font.color.rgb = color or WHITE_

    def footer(sl, n):
        """Add standard Pacebird footer with slide number."""
        txt(sl, "Pacebird — Audience Solutions", 0.3, 7.1, 5, 0.3, size=9, color=GREY_)
        txt(sl, str(n), 12.8, 7.1, 0.5, 0.3, size=9, color=GREY_, align=PP_ALIGN.RIGHT)

    # ── Unpack context ─────────────────────────────────────────────────────────
    segs    = ctx["selected_segs"]
    brief   = ctx.get("brief_text", "")
    bgt     = ctx.get("budget", 0)
    fstart  = ctx.get("flight_start")
    fend    = ctx.get("flight_end")
    fdays   = ctx.get("flight_days")
    obj     = ctx.get("objective", "")
    reach   = ctx.get("est_reach", 0)
    use_and = ctx.get("use_and", False)
    cpm_lo  = ctx.get("wt_cpm_min", 0)
    cpm_hi  = ctx.get("wt_cpm_max", 0)
    min_sp  = ctx.get("max_min_spend", 0)
    idx     = ctx.get("wt_index", 0)
    d_type  = ctx.get("deal_type", "")
    d_rat   = ctx.get("deal_rationale", "")
    floor   = ctx.get("floor_cpm", 0)
    freq    = ctx.get("freq_cap", "")
    cats    = ctx.get("categories", [])
    sigs    = ctx.get("signal_types", [])
    sense   = ctx.get("sense_check", "")

    slide_n = 1

    # ── Slide 1: Title ─────────────────────────────────────────────────────────
    sl = new_slide()
    txt(sl, "Audience Solutions Proposal", 1.0, 2.0, 11.0, 1.0, size=32, bold=True)
    txt(sl, "Are Media — First-Party Audience Taxonomy", 1.0, 3.1, 11.0, 0.6, size=18, color=ORANGE)
    txt(sl,
        f"Prepared {date.today().strftime('%d %B %Y')}  ·  "
        "Fabricated demo data — not real Are Media inventory",
        1.0, 4.1, 11.0, 0.4, size=10, color=GREY_, italic=True)
    footer(sl, slide_n); slide_n += 1

    # ── Slide 2: Brief Summary ─────────────────────────────────────────────────
    if sections.get("brief_summary") and brief.strip():
        sl = new_slide()
        txt(sl, "Client Brief", 0.5, 0.3, 12.0, 0.6, size=22, bold=True)
        # Truncate brief for readability on slide
        brief_s = (brief.strip()[:500] + "…") if len(brief.strip()) > 500 else brief.strip()
        txt(sl, brief_s, 0.5, 1.1, 12.0, 3.5, size=12)
        params = []
        if bgt and bgt > 0:
            params.append(f"Budget: A${bgt:,}")
        if fstart and fend:
            params.append(f"Flight: {fstart.strftime('%d %b %Y')} → {fend.strftime('%d %b %Y')} ({fdays} days)")
        params.append(f"Objective: {obj}")
        for i, p in enumerate(params):
            txt(sl, f"• {p}", 0.5, 4.9 + i * 0.45, 12.0, 0.4, size=12, color=ORANGE)
        footer(sl, slide_n); slide_n += 1

    # ── Slides: Recommended Segments (3 per slide) ────────────────────────────
    if sections.get("recommended_segments") and segs:
        chunk_sz = 3
        for c_start in range(0, len(segs), chunk_sz):
            chunk = segs[c_start:c_start + chunk_sz]
            sl = new_slide()
            txt(sl, "Recommended Segments", 0.5, 0.3, 12.0, 0.6, size=22, bold=True)
            if len(segs) > chunk_sz:
                txt(sl, f"({c_start+1}–{c_start+len(chunk)} of {len(segs)})",
                    0.5, 0.85, 6.0, 0.3, size=10, color=GREY_)

            col_w = 12.0 / max(len(chunk), 1)
            for i, seg in enumerate(chunk):
                x = 0.5 + i * col_w
                # Card background shape with orange top border
                card = sl.shapes.add_shape(
                    1, Inches(x), Inches(1.3), Inches(col_w - 0.15), Inches(5.5)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = NAVY_
                card.line.color.rgb = ORANGE

                # Segment content inside card
                txt(sl, seg["segment_id"],   x+0.1, 1.4,  col_w-0.25, 0.3, size=8,  color=GREY_)
                txt(sl, seg["segment_name"], x+0.1, 1.68, col_w-0.25, 0.5, size=11, bold=True)
                txt(sl, seg["category"],     x+0.1, 2.22, col_w-0.25, 0.3, size=9,  color=ORANGE)

                desc_s = seg["description"][:210] + "…" if len(seg["description"]) > 210 else seg["description"]
                txt(sl, desc_s,              x+0.1, 2.58, col_w-0.25, 1.65, size=8, color=GREY_)

                txt(sl, f"Reach: {seg['reach_monthly']:,}",
                    x+0.1, 4.28, col_w-0.25, 0.32, size=10)
                txt(sl, f"Index: {seg['index_general_pop']:.1f}×",
                    x+0.1, 4.62, col_w-0.25, 0.32, size=10)
                txt(sl, f"CPM: A${seg['indicative_cpm_min']:.0f}–A${seg['indicative_cpm_max']:.0f}",
                    x+0.1, 4.96, col_w-0.25, 0.32, size=10)
                txt(sl, seg["signal_type"],
                    x+0.1, 5.32, col_w-0.25, 0.3,  size=8, color=GREY_)

            footer(sl, slide_n); slide_n += 1

    # ── Slide: Stack Summary ───────────────────────────────────────────────────
    if sections.get("stack_summary") and segs:
        sl = new_slide()
        txt(sl, "Stack Summary", 0.5, 0.3, 12.0, 0.6, size=22, bold=True)

        metric_cards = [
            ("Combined Reach",  f"{reach:,}",
             "AND intersection" if use_and else "OR union · approx."),
            ("Blended CPM",     f"A${cpm_lo:.0f}–A${cpm_hi:.0f}",
             "reach-weighted average"),
            ("Min Spend",       f"A${min_sp:,}",
             "highest across stack"),
            ("Wtd Avg Index",   f"{idx:.1f}×",
             "vs general population"),
        ]
        for i, (lbl, val, sub) in enumerate(metric_cards):
            x = 0.5 + i * 3.1
            c = sl.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(2.85), Inches(1.8))
            c.fill.solid(); c.fill.fore_color.rgb = NAVY_; c.line.color.rgb = ORANGE
            txt(sl, lbl, x+0.1, 1.3,  2.6, 0.3, size=9,  color=GREY_)
            txt(sl, val, x+0.1, 1.62, 2.6, 0.6, size=17, bold=True)
            txt(sl, sub, x+0.1, 2.28, 2.6, 0.3, size=8,  color=GREY_, italic=True)

        txt(sl, "Segments in stack:", 0.5, 3.3, 12.0, 0.35, size=11, color=GREY_)
        for j, seg in enumerate(segs[:8]):   # cap at 8 lines on slide
            txt(sl, f"• {seg['segment_id']}  {seg['segment_name']}",
                0.5, 3.68 + j * 0.38, 12.0, 0.35, size=10)

        txt(sl,
            "Reach figures are approximate. AND = 30% of smallest segment; "
            "OR = sum minus deduplication haircut.",
            0.5, 6.8, 12.0, 0.3, size=8, color=GREY_, italic=True)
        footer(sl, slide_n); slide_n += 1

    # ── Slide: Deal Structure ─────────────────────────────────────────────────
    if sections.get("deal_structure"):
        sl = new_slide()
        txt(sl, "Deal Structure Recommendation", 0.5, 0.3, 12.0, 0.6, size=22, bold=True)
        txt(sl, d_type,  0.5, 1.1,  12.0, 0.6, size=24, bold=True, color=ORANGE)
        txt(sl, d_rat,   0.5, 1.75, 12.0, 0.5, size=12, color=GREY_)

        rows = [
            ("Floor CPM",     f"A${floor:.2f}"),
            ("Frequency Cap", freq),
            ("Objective",     obj),
            ("Categories",    ", ".join(cats)),
            ("Signal Types",  ", ".join(sigs)),
        ]
        for j, (lbl, val) in enumerate(rows):
            y = 2.5 + j * 0.52
            txt(sl, lbl + ":", 0.5, y, 3.0, 0.4, size=11, color=GREY_)
            txt(sl, val,       3.6, y, 9.2, 0.4, size=11, bold=True)

        if sense:
            txt(sl, f"Sense-check: {sense}", 0.5, 5.55, 12.0, 0.55, size=11, color=GREEN_)

        txt(sl,
            "Indicative only. Validate availability in the ad server before committing.",
            0.5, 6.8, 12.0, 0.3, size=8, color=GREY_, italic=True)
        footer(sl, slide_n); slide_n += 1

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# EXPORT DIALOG
# Defined at module level so it is always reachable when the button is pressed.
# It reads its context from st.session_state["_as_ctx"] which is written
# at the end of Section 4 every time segments are selected.
# ═══════════════════════════════════════════════════════════════════════════════

@st.dialog("Generate Proposal")
def show_proposal_dialog():
    """Section-picker dialog — same pattern as QBR Generator."""
    ctx = st.session_state.get("_as_ctx")
    if not ctx or not ctx.get("selected_segs"):
        st.error("No stack data found. Select segments and try again.")
        return

    st.caption("Select the sections to include in your PowerPoint export.")

    defaults = st.session_state.get("_as_section_defaults", {
        "brief_summary":        True,
        "recommended_segments": True,
        "stack_summary":        True,
        "deal_structure":       True,
    })

    sc = st.columns(2)
    with sc[0]:
        if st.button("Select all", key="as_sel_all"):
            st.session_state["_as_section_defaults"] = {k: True for k in defaults}
            st.rerun()
    with sc[1]:
        if st.button("Deselect all", key="as_desel_all"):
            st.session_state["_as_section_defaults"] = {k: False for k in defaults}
            st.rerun()

    st.markdown("")
    s1 = st.checkbox("Brief Summary",         value=defaults.get("brief_summary", True),        key="as_cb_brief")
    s2 = st.checkbox("Recommended Segments",  value=defaults.get("recommended_segments", True),  key="as_cb_segs")
    s3 = st.checkbox("Stack Summary",         value=defaults.get("stack_summary", True),         key="as_cb_stack")
    s4 = st.checkbox("Deal Structure",        value=defaults.get("deal_structure", True),        key="as_cb_deal")

    st.markdown("")
    bc = st.columns(2)
    with bc[0]:
        if st.button("Build Proposal", type="primary", key="as_gen_btn"):
            st.session_state["_as_section_defaults"] = {
                "brief_summary":        s1,
                "recommended_segments": s2,
                "stack_summary":        s3,
                "deal_structure":       s4,
            }
            secs = st.session_state["_as_section_defaults"]
            if not any(secs.values()):
                st.warning("Select at least one section.")
                return
            with st.spinner("Building proposal deck..."):
                try:
                    buf = _build_proposal_pptx(ctx, secs)
                    st.session_state["_as_pptx"]    = buf
                    st.session_state["_as_filename"] = (
                        f"AudienceSolutions_{date.today().isoformat()}.pptx"
                    )
                except Exception as e:
                    st.error(f"Export failed: {e}")
                    return
    with bc[1]:
        if st.button("Cancel", key="as_cancel_btn"):
            st.rerun()

    # Download shown immediately once PPTX is ready (no rerun needed — stays in dialog)
    if "_as_pptx" in st.session_state:
        st.success("Proposal ready!")
        st.download_button(
            "📥 Download Proposal (.pptx)",
            data=st.session_state["_as_pptx"],
            file_name=st.session_state.get("_as_filename", "AudienceSolutions_Proposal.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="as_dl_btn",
        )


# ═══════════════════════════════════════════════════════════════════════════════
# PAGE CONTENT
# ═══════════════════════════════════════════════════════════════════════════════

# ── Header ─────────────────────────────────────────────────────────────────────
st.markdown("# 🎯 Audience Solutions")
st.markdown(
    '<div class="disclaimer-bar">'
    "⚠️ Demo using fabricated segment data. Not real Are Media inventory."
    "</div>",
    unsafe_allow_html=True,
)
st.markdown("Find, stack, and structure first-party audience segments from the Are Media taxonomy.")

segments   = load_segments()
desc_tuple = tuple(s["description"] for s in segments)


# ── SECTION 1: BRIEF INPUT ───────────────────────────────────────────────────
st.markdown('<hr class="sec-rule">', unsafe_allow_html=True)
st.markdown("### 1 — Client Brief")

brief_text = st.text_area(
    "Paste the client brief or campaign details",
    height=140,
    placeholder=(
        "Example: We're launching a prestige skincare range in Australia targeting women 35–55 "
        "who actively purchase premium skincare. Campaign runs 6 weeks with a A$45,000 budget. "
        "Primary goal is consideration — we want women already in the research phase with a high "
        "propensity to spend above $80 per product."
    ),
    key="as_brief",
)

# Compact structured inputs in a single row
ci1, ci2, ci3, ci4 = st.columns([2, 2, 2, 3])
with ci1:
    budget = st.number_input(
        "Budget (A$)", min_value=0, step=1000, value=0, key="as_budget",
        help="Leave at 0 to skip the budget sense-check",
    )
with ci2:
    flight_start = st.date_input("Flight start", value=None, key="as_start")
with ci3:
    flight_end   = st.date_input("Flight end",   value=None, key="as_end")
with ci4:
    objective = st.selectbox(
        "Primary objective",
        ["Awareness", "Consideration", "Conversion"],
        key="as_objective",
    )

find_btn = st.button("🔍 Find matching segments", type="primary", key="as_find")

if find_btn:
    if not brief_text.strip():
        st.warning("Please paste a brief before searching.")
    else:
        with st.spinner("Embedding brief and scoring segments..."):
            try:
                model     = load_model()
                seg_embs  = embed_descriptions(desc_tuple)
                query_vec = model.encode(brief_text.strip())
                scored = sorted(
                    [
                        (cosine_sim(query_vec, seg_embs[i]), segments[i])
                        for i in range(len(segments))
                    ],
                    key=lambda x: x[0],
                    reverse=True,
                )
                st.session_state["as_results"] = scored
                # Clear prior segment selections when a new search runs
                for seg in segments:
                    st.session_state.pop(f"as_sel_{seg['segment_id']}", None)
            except Exception as e:
                st.error(f"Semantic search failed: {e}")


# ── SECTION 2: SEGMENT MATCHES ────────────────────────────────────────────────
if "as_results" in st.session_state:
    st.markdown('<hr class="sec-rule">', unsafe_allow_html=True)
    st.markdown("### 2 — Segment Matches")

    threshold = st.slider(
        "Relevance threshold — hide segments below this match score",
        min_value=0.10, max_value=0.80, value=0.20, step=0.01,
        format="%.2f", key="as_threshold",
    )

    all_results = st.session_state["as_results"]
    shown       = [(score, seg) for score, seg in all_results if score >= threshold][:8]

    if not shown:
        st.info("No segments meet the threshold. Try lowering the slider.")
    else:
        st.caption(
            f"Top {len(shown)} segments at or above {threshold:.0%} match score. "
            "Check a segment to add it to your stack."
        )

        for score, seg in shown:
            sid    = seg["segment_id"]
            is_sel = st.session_state.get(f"as_sel_{sid}", False)

            card_cls  = "seg-card seg-card-selected" if is_sel else "seg-card"
            badge_cls = match_badge_class(score)
            act_html  = "".join(
                f'<span class="badge-act">{a}</span>' for a in seg["activation"]
            )

            col_card, col_cb = st.columns([11, 1])
            with col_card:
                st.markdown(f"""
                <div class="{card_cls}">
                    <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;margin-bottom:6px;">
                        <span style="font-weight:700;font-size:16px;color:{SECONDARY};">{seg['segment_name']}</span>
                        <span class="badge-cat">{seg['category']}</span>
                        <span class="{badge_cls}">{score:.0%} match</span>
                        <span class="badge-signal">{seg['signal_type']}</span>
                    </div>
                    <div style="color:{TEXT_SEC};font-size:13px;line-height:1.55;">{seg['description']}</div>
                    <div style="display:flex;flex-wrap:wrap;gap:20px;margin-top:12px;align-items:flex-start;">
                        <div>
                            <div class="mlabel">Monthly Reach</div>
                            <div class="mval">{seg['reach_monthly']:,}</div>
                        </div>
                        <div>
                            <div class="mlabel">Index vs Pop</div>
                            <div class="mval">{seg['index_general_pop']:.1f}×</div>
                        </div>
                        <div>
                            <div class="mlabel">CPM Range</div>
                            <div class="mval">A${seg['indicative_cpm_min']:.0f}–A${seg['indicative_cpm_max']:.0f}</div>
                        </div>
                        <div>
                            <div class="mlabel">Min Spend</div>
                            <div class="mval">A${seg['min_spend']:,}</div>
                        </div>
                        <div>
                            <div class="mlabel">Activation</div>
                            <div style="margin-top:3px;">{act_html}</div>
                        </div>
                    </div>
                </div>
                """, unsafe_allow_html=True)

            with col_cb:
                # Vertical offset to visually centre the checkbox against the card
                st.markdown("<div style='margin-top:22px;'>", unsafe_allow_html=True)
                st.checkbox(
                    "Select",
                    value=is_sel,
                    key=f"as_sel_{sid}",
                    label_visibility="collapsed",
                )
                st.markdown("</div>", unsafe_allow_html=True)

    # Collect all checked segments (scan across all results, not just shown ones)
    selected_segs = [
        seg for _, seg in st.session_state["as_results"]
        if st.session_state.get(f"as_sel_{seg['segment_id']}", False)
    ]


    # ── SECTION 3: STACK BUILDER ──────────────────────────────────────────────
    st.markdown('<hr class="sec-rule">', unsafe_allow_html=True)
    st.markdown("### 3 — Stack Builder")

    if not selected_segs:
        st.info("Check segments above to add them to your stack.")
    else:
        # Base numeric vectors from selected segments
        reaches   = [s["reach_monthly"]     for s in selected_segs]
        cpm_mins  = [s["indicative_cpm_min"] for s in selected_segs]
        cpm_maxs  = [s["indicative_cpm_max"] for s in selected_segs]
        min_spds  = [s["min_spend"]          for s in selected_segs]
        indexes   = [s["index_general_pop"]  for s in selected_segs]
        n_segs    = len(selected_segs)
        tot_reach = sum(reaches)

        # AND / OR toggle
        reach_mode = st.radio(
            "Reach estimation logic",
            [
                "OR — Union (sum with deduplication haircut)",
                "AND — Intersection (conservative overlap estimate)",
            ],
            horizontal=True,
            key="as_reach_mode",
        )
        use_and = reach_mode.startswith("AND")

        if use_and:
            # Conservative AND: 30% of the smallest segment's monthly reach
            est_reach  = int(min(reaches) * 0.30)
            reach_note = "AND — conservative 30% of smallest segment (approximate)"
        else:
            # OR with dedup haircut that grows with stack depth: 10% per extra segment, capped at 40%
            haircut    = min(0.40, 0.10 * (n_segs - 1))
            est_reach  = int(tot_reach * (1 - haircut))
            reach_note = f"OR — {haircut:.0%} dedup haircut applied (approximate)"

        # Reach-weighted CPM and index
        wt_cpm_min    = sum(r * c for r, c in zip(reaches, cpm_mins)) / tot_reach
        wt_cpm_max    = sum(r * c for r, c in zip(reaches, cpm_maxs)) / tot_reach
        wt_index      = sum(r * i for r, i in zip(reaches, indexes))  / tot_reach
        max_min_spend = max(min_spds)

        # Stack summary panel
        st.markdown(f"""
        <div class="stack-panel">
            <div style="font-weight:700;font-size:16px;color:{SECONDARY};margin-bottom:18px;">
                Stack Summary — {n_segs} segment{'s' if n_segs > 1 else ''} selected
            </div>
            <div style="display:grid;grid-template-columns:repeat(4,1fr);gap:20px;">
                <div>
                    <div class="mlabel">Est. Combined Reach</div>
                    <div style="font-size:26px;font-weight:700;color:{TEXT_PRI};">{est_reach:,}</div>
                    <div style="font-size:11px;color:{TEXT_SEC};margin-top:2px;">{reach_note}</div>
                </div>
                <div>
                    <div class="mlabel">Blended CPM</div>
                    <div style="font-size:26px;font-weight:700;color:{TEXT_PRI};">A${wt_cpm_min:.0f}–A${wt_cpm_max:.0f}</div>
                    <div style="font-size:11px;color:{TEXT_SEC};margin-top:2px;">reach-weighted average</div>
                </div>
                <div>
                    <div class="mlabel">Min Spend Required</div>
                    <div style="font-size:26px;font-weight:700;color:{TEXT_PRI};">A${max_min_spend:,}</div>
                    <div style="font-size:11px;color:{TEXT_SEC};margin-top:2px;">highest across stack</div>
                </div>
                <div>
                    <div class="mlabel">Wtd Avg Index</div>
                    <div style="font-size:26px;font-weight:700;color:{TEXT_PRI};">{wt_index:.1f}×</div>
                    <div style="font-size:11px;color:{TEXT_SEC};margin-top:2px;">vs general population</div>
                </div>
            </div>
            <div style="margin-top:14px;font-size:12px;color:{TEXT_SEC};">
                <strong>Stack:</strong> {' · '.join(s['segment_id'] for s in selected_segs)}
            </div>
        </div>
        """, unsafe_allow_html=True)

        # ── Automatic warnings ────────────────────────────────────────────────

        # Derive flight_days here so it's available for both warning checks and deal logic
        flight_days = None
        if flight_start and flight_end:
            flight_days = max(1, (flight_end - flight_start).days)

        # 1. Stack too narrow for AND logic
        if use_and and est_reach < 50000:
            st.markdown(
                '<div class="warn-box">⚠️  Stack may be too narrow to deliver at scale. '
                "AND logic produces under 50,000 addressable users — consider switching to "
                "OR logic or adding a broader segment.</div>",
                unsafe_allow_html=True,
            )

        # 2. Budget vs available impressions
        if budget and budget > 0 and flight_days:
            mid_cpm = (wt_cpm_min + wt_cpm_max) / 2
            if mid_cpm > 0:
                implied_imps  = int((budget / mid_cpm) * 1000)
                flight_months = flight_days / 30
                avail_imps    = int(est_reach * flight_months)
                if implied_imps > avail_imps:
                    st.markdown(
                        f'<div class="warn-box">⚠️  Requested volume exceeds available reach. '
                        f"A${budget:,} at A${mid_cpm:.2f} blended CPM implies {implied_imps:,} impressions "
                        f"over {flight_days} days — but the stack has only ~{avail_imps:,} available "
                        f"impressions for this flight length.</div>",
                        unsafe_allow_html=True,
                    )

        # 3. Notes-based stacking caveats (surface each unique segment note once)
        stack_kws = ["do not stack", "avoid stacking", "stacking", "frequency cap", "saturation", "over-invest"]
        seen_caveats = set()
        for seg in selected_segs:
            notes_lower = seg.get("notes", "").lower()
            if any(kw in notes_lower for kw in stack_kws):
                if seg["segment_id"] not in seen_caveats:
                    seen_caveats.add(seg["segment_id"])
                    st.markdown(
                        f'<div class="warn-box">⚠️  <strong>{seg["segment_id"]} — '
                        f'{seg["segment_name"]}:</strong> {seg["notes"]}</div>',
                        unsafe_allow_html=True,
                    )


        # ── SECTION 4: DEAL STRUCTURE ─────────────────────────────────────────
        st.markdown('<hr class="sec-rule">', unsafe_allow_html=True)
        st.markdown("### 4 — Deal Structure Recommendation")

        # Rules-based deal type: budget + flight length drive the recommendation
        if budget and budget >= 50000 and flight_days and flight_days >= 28:
            deal_type      = "Programmatic Guaranteed (PG)"
            deal_rationale = "Budget and flight length support a committed volume deal; PG locks priority access and guaranteed delivery."
        elif budget and budget >= 20000:
            deal_type      = "Private Marketplace (PMP)"
            deal_rationale = "Mid-range budget suits preferred access to first-party inventory; PMP delivers priority without a guaranteed commitment."
        elif budget and 0 < budget < 20000:
            deal_type      = "Curated Deal"
            deal_rationale = "Smaller budget fits a curated package at a fixed floor — low friction with no guaranteed delivery risk."
        else:
            deal_type      = "PMP or Curated (confirm budget)"
            deal_rationale = "No budget entered — defaulting to PMP or curated. Enter a budget to sharpen the recommendation."

        # Suggested floor CPM: blended midpoint + 10% deal access premium
        mid_cpm   = (wt_cpm_min + wt_cpm_max) / 2
        floor_cpm = round(mid_cpm * 1.10, 2)

        # Frequency cap by objective
        freq_caps = {
            "Awareness":     "3–5 impressions / user / day",
            "Consideration": "2–3 impressions / user / day",
            "Conversion":    "1–2 impressions / user / day",
        }
        freq_cap = freq_caps.get(objective, "2–3 impressions / user / day")

        # Targeting parameter summary
        categories   = sorted({s["category"]   for s in selected_segs})
        signal_types = sorted({s["signal_type"] for s in selected_segs})
        all_brands   = sorted({b for s in selected_segs for b in s["source_brands"]})
        brands_str   = ", ".join(all_brands[:4]) + ("…" if len(all_brands) > 4 else "")

        # Sense-check line
        sense_check = ""
        if budget and budget > 0 and flight_days and flight_days > 0 and mid_cpm > 0:
            implied_imps = int((budget / mid_cpm) * 1000)
            sense_check = (
                f"A${budget:,} at A${mid_cpm:.2f} blended CPM implies ~{implied_imps:,} impressions "
                f"over {flight_days} days — the selected stack has ~{est_reach:,} monthly reach."
            )

        sense_html = (
            f'<div class="sense-check"><strong>Sense-check:</strong> {sense_check}</div>'
            if sense_check else ""
        )

        st.markdown(f"""
        <div class="deal-panel">
            <div style="font-weight:700;font-size:16px;color:{SECONDARY};margin-bottom:20px;">
                Indicative Deal Structure
            </div>
            <div style="display:grid;grid-template-columns:1fr 1fr;gap:24px;">
                <div>
                    <div class="mlabel">Recommended Deal Type</div>
                    <div style="font-size:20px;font-weight:700;color:{PRIMARY};margin:5px 0 6px 0;">{deal_type}</div>
                    <div style="font-size:13px;color:{TEXT_SEC};">{deal_rationale}</div>
                </div>
                <div>
                    <div class="mlabel">Suggested Floor CPM</div>
                    <div style="font-size:20px;font-weight:700;color:{TEXT_PRI};margin:5px 0 6px 0;">A${floor_cpm:.2f}</div>
                    <div style="font-size:13px;color:{TEXT_SEC};">+10% premium on blended CPM midpoint</div>
                </div>
                <div>
                    <div class="mlabel">Suggested Frequency Cap</div>
                    <div style="font-size:16px;font-weight:600;color:{TEXT_PRI};margin:5px 0 6px 0;">{freq_cap}</div>
                    <div style="font-size:13px;color:{TEXT_SEC};">Based on <strong>{objective}</strong> objective</div>
                </div>
                <div>
                    <div class="mlabel">Targeting Parameters</div>
                    <div style="font-size:13px;color:{TEXT_PRI};margin-top:6px;line-height:1.7;">
                        <strong>Categories:</strong> {', '.join(categories)}<br>
                        <strong>Signal types:</strong> {', '.join(signal_types)}<br>
                        <strong>Source brands:</strong> {brands_str}
                    </div>
                </div>
            </div>
            {sense_html}
            <div style="margin-top:14px;font-size:11px;color:{TEXT_SEC};font-style:italic;">
                Indicative only. Validate availability in the ad server before committing.
            </div>
        </div>
        """, unsafe_allow_html=True)

        # Store full context for the PPTX builder to consume from session state
        st.session_state["_as_ctx"] = {
            "selected_segs":  selected_segs,
            "brief_text":     brief_text,
            "budget":         budget,
            "flight_start":   flight_start,
            "flight_end":     flight_end,
            "flight_days":    flight_days,
            "objective":      objective,
            "est_reach":      est_reach,
            "use_and":        use_and,
            "wt_cpm_min":     wt_cpm_min,
            "wt_cpm_max":     wt_cpm_max,
            "max_min_spend":  max_min_spend,
            "wt_index":       wt_index,
            "deal_type":      deal_type,
            "deal_rationale": deal_rationale,
            "floor_cpm":      floor_cpm,
            "freq_cap":       freq_cap,
            "categories":     categories,
            "signal_types":   signal_types,
            "sense_check":    sense_check,
        }

        # ── SECTION 5: EXPORT ─────────────────────────────────────────────────
        st.markdown('<hr class="sec-rule">', unsafe_allow_html=True)

        col_btn, col_dl, _ = st.columns([2, 3, 4])
        with col_btn:
            if st.button("📋 Generate Proposal", type="primary", key="as_open_dialog"):
                show_proposal_dialog()
        with col_dl:
            # Persistent download — visible outside the dialog once generated
            if "_as_pptx" in st.session_state:
                st.download_button(
                    "📥 Download Proposal (.pptx)",
                    data=st.session_state["_as_pptx"],
                    file_name=st.session_state.get(
                        "_as_filename", "AudienceSolutions_Proposal.pptx"
                    ),
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="as_dl_persistent",
                )
