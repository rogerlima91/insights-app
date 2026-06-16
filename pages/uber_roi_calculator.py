import io
import os
from datetime import date as _date

import anthropic
import pandas as pd
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# STYLE LOCK: Do not remove or modify this CSS block.
st.markdown("""
<style>
    /* ── Base font and body ─────────────────────────────────────── */
    html, body, [class*="css"] {
        font-family: "Inter", system-ui, -apple-system, "Segoe UI", sans-serif;
        font-size: 15px;
        color: #374151;
    }

    /* ── Page background ────────────────────────────────────────── */
    .stApp { background-color: #F3F4F6; }

    /* ── Main area headings ──────────────────────────────────────── */
    h1, h2, h3, h4, h5, h6 { font-weight: 700 !important; color: #111827 !important; }
    h2, h3 {
        margin-top: 2rem !important;
        padding-top: 0.25rem !important;
        border-bottom: none !important;
        padding-bottom: 0 !important;
        border-left: none !important;
        padding-left: 0 !important;
    }

    /* ── KPI metric cards ───────────────────────────────────────── */
    [data-testid="metric-container"] {
        background: #FFFFFF;
        border: none;
        border-radius: 16px;
        padding: 20px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        border-top: 4px solid #7C3AED;
    }
    [data-testid="metric-container"] label {
        font-size: 12px;
        font-weight: 600;
        color: #6B7280;
        text-transform: uppercase;
        letter-spacing: 0.05em;
    }
    [data-testid="metric-container"] [data-testid="stMetricValue"] {
        font-size: 26px;
        font-weight: 700;
        color: #111827;
    }

    /* ── Chart cards — wrap Plotly chart output in white card ───── */
    .element-container:has([data-testid="stPlotlyChart"]) {
        background: #FFFFFF;
        border-radius: 12px;
        padding: 16px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
    }

    /* ── Primary buttons — purple ───────────────────────────────── */
    .stButton > button[kind="primary"],
    [data-testid="baseButton-primary"] {
        background-color: #7C3AED !important;
        color: #FFFFFF !important;
        border: none !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        font-size: 15px !important;
        padding: 0.5rem 1.25rem !important;
    }
    .stButton > button[kind="primary"]:hover,
    [data-testid="baseButton-primary"]:hover {
        background-color: #6D28D9 !important;
    }

    /* ── Secondary / default buttons ────────────────────────────── */
    .stButton > button[kind="secondary"],
    [data-testid="baseButton-secondary"] {
        border-radius: 8px !important;
        font-weight: 600 !important;
    }

    /* ── Data tables ─────────────────────────────────────────────── */
    [data-testid="stDataFrame"] {
        background: #FFFFFF;
        border-radius: 12px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.06);
    }
</style>
""", unsafe_allow_html=True)
# STYLE LOCK


# ── Category benchmarks (hardcoded) ─────────────────────────────────────────────
# avg_cpo  = average cost per order (AUD)
# avg_roas = average return on ad spend (revenue / ad spend)
# avg_ctr  = average click-through rate (%)
# incr_rate = share of ad-driven orders that are truly incremental (0–1)
CATEGORY_BENCHMARKS = {
    "QSR Burgers":        {"avg_cpo": 9.50,  "avg_roas": 3.4, "avg_ctr": 2.9, "incr_rate": 0.58},
    "Pizza":              {"avg_cpo": 8.20,  "avg_roas": 3.6, "avg_ctr": 3.2, "incr_rate": 0.60},
    "Mexican":            {"avg_cpo": 10.50, "avg_roas": 2.9, "avg_ctr": 2.7, "incr_rate": 0.55},
    "Asian":              {"avg_cpo": 9.80,  "avg_roas": 3.1, "avg_ctr": 2.8, "incr_rate": 0.56},
    "Healthy/Salads":     {"avg_cpo": 11.20, "avg_roas": 2.7, "avg_ctr": 2.5, "incr_rate": 0.52},
    "Coffee & Breakfast": {"avg_cpo": 7.80,  "avg_roas": 3.8, "avg_ctr": 3.5, "incr_rate": 0.62},
    "Other":              {"avg_cpo": 10.00, "avg_roas": 3.0, "avg_ctr": 2.8, "incr_rate": 0.56},
}

# ── Daypart base weight distribution ────────────────────────────────────────────
# Reflects typical Uber Eats order volume distribution across the day
DAYPART_BASE_WEIGHTS = {
    "Breakfast":  0.08,
    "Lunch":      0.25,
    "Afternoon":  0.10,
    "Dinner":     0.42,
    "Late Night": 0.15,
}

# ── Scenario multipliers applied to the adjusted benchmark ROAS ─────────────────
SCENARIO_FACTORS = {
    "Conservative": 0.75,
    "Base Case":    1.00,
    "Optimistic":   1.35,
}

# Column header colours: (header_bg, header_text, border_css)
SCENARIO_COLOURS = {
    "Conservative": ("#6B7280", "#FFFFFF", "1px solid #E5E7EB"),
    "Base Case":    ("#2563EB", "#FFFFFF", "2px solid #2563EB"),
    "Optimistic":   ("#059669", "#FFFFFF", "1px solid #E5E7EB"),
}


# ── Reusable HTML helpers ────────────────────────────────────────────────────────

def _metric_row(label, value, value_color="#111827"):
    """Single label / value row inside a card — used in scenario cards."""
    return (
        f"<div style='display:flex;justify-content:space-between;align-items:center;"
        f"padding:6px 0;border-bottom:1px solid #F9FAFB;'>"
        f"<span style='color:#6B7280;font-size:12px;'>{label}</span>"
        f"<span style='color:{value_color};font-weight:600;font-size:13px;'>{value}</span>"
        f"</div>"
    )


def _table_row(cells, colors=None, bold_last=False):
    """
    Build a <tr> for use in a styled HTML table.
    cells  = list of cell text strings
    colors = optional list of per-cell text colours (same length as cells)
    """
    if colors is None:
        colors = ["#374151"] * len(cells)
    tds = ""
    for i, (cell, color) in enumerate(zip(cells, colors)):
        fw = "700" if bold_last and i == len(cells) - 1 else "400"
        tds += (
            f"<td style='padding:10px 14px;border-bottom:1px solid #F3F4F6;"
            f"color:{color};font-weight:{fw};font-size:13px;'>{cell}</td>"
        )
    return f"<tr>{tds}</tr>"


def _html_table(headers, rows):
    """
    Wrap headers (list of str) and rows (list of <tr> strings) into a
    full white-card HTML table matching the app style.
    """
    header_cells = "".join(
        f"<th style='padding:10px 14px;text-align:left;font-size:11px;"
        f"font-weight:700;color:#6B7280;text-transform:uppercase;"
        f"letter-spacing:0.05em;border-bottom:2px solid #E5E7EB;'>{h}</th>"
        for h in headers
    )
    body = "".join(rows)
    return (
        "<div style='background:#FFFFFF;border-radius:12px;overflow:hidden;"
        "box-shadow:0 4px 12px rgba(0,0,0,0.08);margin-top:12px;'>"
        f"<table style='width:100%;border-collapse:collapse;'>"
        f"<thead><tr>{header_cells}</tr></thead>"
        f"<tbody>{body}</tbody>"
        "</table></div>"
    )


# ── API key loader ───────────────────────────────────────────────────────────────

def get_api_key():
    """Load Anthropic API key from Streamlit secrets or environment variables."""
    if "ANTHROPIC_API_KEY" in st.secrets:
        return st.secrets.get("ANTHROPIC_API_KEY")
    return os.environ.get("ANTHROPIC_API_KEY")


# ── Status helper (for scorecard) ────────────────────────────────────────────────

def _get_status(actual, target, lower_is_better=False):
    """
    Return (emoji, hex_colour) based on actual vs target ratio.
    lower_is_better=True inverts the comparison (e.g. CPO, spend).
    """
    if target <= 0:
        return "—", "#6B7280"
    ratio = (target / actual) if lower_is_better and actual > 0 else (actual / target)
    if ratio >= 1.0:
        return "✅", "#10B981"
    elif ratio >= 0.85:
        return "⚠️", "#F59E0B"
    else:
        return "❌", "#EF4444"


# ── Campaign file parser (Tab 2) ────────────────────────────────────────────────

# Column name candidates for each metric — checked in order, first match wins
_UPLOAD_COLUMN_CANDIDATES = {
    "spend":        ["spend", "cost", "total spend", "ad spend"],
    "total_orders": ["orders", "total orders", "conversions"],
    "incr_orders":  ["incremental orders", "incremental conversions"],
    "impressions":  ["impressions"],
    "clicks":       ["clicks"],
}


def parse_campaign_file(uploaded_file):
    """
    Parse a CSV or XLSX Uber Ads export.
    Auto-detects metric columns using _UPLOAD_COLUMN_CANDIDATES,
    sums all rows, and returns campaign totals.

    Returns:
        detected   — dict {metric: summed_float} for each column found
        col_found  — dict {metric: matched_column_name} for the preview label
    """
    try:
        if uploaded_file.name.lower().endswith(".xlsx"):
            df = pd.read_excel(uploaded_file)
        else:
            df = pd.read_csv(uploaded_file)
    except Exception:
        return {}, {}

    if df.empty:
        return {}, {}

    # Build a lowercase → actual column name lookup
    col_lower = {c.lower().strip(): c for c in df.columns}

    detected  = {}
    col_found = {}
    for metric, candidates in _UPLOAD_COLUMN_CANDIDATES.items():
        for cand in candidates:
            if cand in col_lower:
                actual_col = col_lower[cand]
                numeric    = pd.to_numeric(df[actual_col], errors="coerce")
                detected[metric]  = float(numeric.sum())
                col_found[metric] = actual_col
                break

    return detected, col_found


# ── Core ROI calculation (Tab 1) ─────────────────────────────────────────────────

def compute_roi(category, avg_order_value, num_locations, rating,
                monthly_budget, campaign_types, dayparts):
    """
    Apply input-based multipliers to category benchmarks, then project
    ROI across three scenarios (Conservative / Base Case / Optimistic).

    All multipliers convert to ROAS adjustments because ROAS = AOV / CPO,
    so any change to CPO or order volume flows through ROAS consistently.

    Returns:
        scenarios — dict {name: {roas, revenue, orders, incr_orders, cpo, profit}}
        adj_roas  — fully adjusted base ROAS (before scenario factor)
        base_cpo  — raw benchmark CPO for the category (for break-even text)
    """
    bench         = CATEGORY_BENCHMARKS[category]
    adj_roas      = bench["avg_roas"]
    base_cpo      = bench["avg_cpo"]
    adj_incr_rate = bench["incr_rate"]

    # Rating: higher-rated restaurants convert ad impressions to orders more efficiently
    if rating >= 4.5:
        adj_roas *= 1.15
    elif rating < 4.0:
        adj_roas *= 0.85

    # Sponsored Listings: CPO −10% → equivalent ROAS gain of 1/0.90
    if "Sponsored Listings" in campaign_types:
        adj_roas *= (1 / 0.90)

    # Homepage Banner: incremental visibility drives ROAS uplift
    if "Homepage Banner" in campaign_types:
        adj_roas *= 1.20

    # Daypart multipliers: peak meal times produce more orders per dollar
    if "Dinner" in dayparts:
        adj_roas *= 1.18
    if "Lunch" in dayparts:
        adj_roas *= 1.12

    # Multiple locations: each location generates orders independently on the same budget
    adj_roas *= num_locations

    # Compute three scenarios
    scenarios = {}
    for name, factor in SCENARIO_FACTORS.items():
        s_roas   = adj_roas * factor
        s_rev    = monthly_budget * s_roas
        s_orders = s_rev / avg_order_value if avg_order_value > 0 else 0
        s_incr   = s_orders * adj_incr_rate
        s_cpo    = monthly_budget / s_orders if s_orders > 0 else 0
        # Profit = 30% gross margin on revenue minus ad spend
        s_profit = s_rev * 0.30 - monthly_budget
        scenarios[name] = {
            "roas":        round(s_roas, 2),
            "revenue":     round(s_rev, 2),
            "orders":      round(s_orders, 1),
            "incr_orders": round(s_incr, 1),
            "cpo":         round(s_cpo, 2),
            "profit":      round(s_profit, 2),
        }

    return scenarios, adj_roas, base_cpo


# ── Campaign outcome metrics (Tab 2) ────────────────────────────────────────────

def compute_outcome_metrics(inputs):
    """
    Derive all calculated performance metrics from campaign actuals entered in Tab 2.
    Returns a dict of rounded metric values.
    """
    spend       = inputs["spend"]
    tot_orders  = inputs["total_orders"]
    incr_orders = inputs["incr_orders"]
    impressions = inputs["impressions"]
    clicks      = inputs["clicks"]
    aov         = inputs["avg_order_value"]

    # ROAS uses only incremental orders (orders that wouldn't have happened without ads)
    actual_roas  = (incr_orders * aov) / spend if spend > 0 else 0
    actual_cpo   = spend / tot_orders          if tot_orders > 0 else 0
    incr_cpo     = spend / incr_orders         if incr_orders > 0 else 0
    ctr          = (clicks / impressions * 100) if impressions > 0 else 0
    conv_rate    = (tot_orders / clicks * 100)  if clicks > 0 else 0
    incr_rate    = (incr_orders / tot_orders * 100) if tot_orders > 0 else 0

    return {
        "actual_roas": round(actual_roas, 2),
        "actual_cpo":  round(actual_cpo, 2),
        "incr_cpo":    round(incr_cpo, 2),
        "ctr":         round(ctr, 2),
        "conv_rate":   round(conv_rate, 2),
        "incr_rate":   round(incr_rate, 1),
    }


# ── Campaign health score (Tab 2) ────────────────────────────────────────────────

def compute_health_score(metrics, inputs, bench):
    """
    Score the campaign 0–100 based on performance vs targets and benchmarks.
    Weights: ROAS 30%, Orders 25%, CPO 25%, Incremental Rate 20%.
    """
    target_roas   = inputs["target_roas"]
    target_orders = inputs["target_orders"]
    bench_incr    = bench["incr_rate"] * 100   # convert to %

    # Implied CPO target = budget / order target
    impl_cpo_target = inputs["spend"] / target_orders if target_orders > 0 else 0

    def _higher_better(actual, target, cap=1.5):
        if target <= 0:
            return 50
        return min(actual / target, cap) / cap * 100

    def _lower_better(actual, target, cap=1.5):
        if actual <= 0 or target <= 0:
            return 50
        return min(target / actual, cap) / cap * 100

    roas_score   = _higher_better(metrics["actual_roas"],  target_roas)
    orders_score = _higher_better(inputs["total_orders"],   target_orders)
    cpo_score    = _lower_better(metrics["actual_cpo"],    impl_cpo_target)
    incr_score   = _higher_better(metrics["incr_rate"],    bench_incr)

    weighted = (roas_score * 0.30 + orders_score * 0.25 +
                cpo_score  * 0.25 + incr_score   * 0.20)
    return round(weighted)


# ── PowerPoint helpers ───────────────────────────────────────────────────────────

def _pptx_base():
    """
    Create a blank Presentation with the dark navy premium template.
    Returns (prs, blank_layout, colour constants, add_text fn, add_footer fn).
    """
    BG    = RGBColor(0x0D, 0x1B, 0x2A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREY  = RGBColor(0xA8, 0xB2, 0xBC)
    BLUE  = RGBColor(0x00, 0xA8, 0xE8)
    GREEN = RGBColor(0x05, 0x96, 0x69)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

    def add_text(slide, text, left, top, width, height,
                 size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text            = text
        run.font.size       = Pt(size)
        run.font.bold       = bold
        run.font.color.rgb  = color
        run.font.name       = "Calibri"
        return tb

    def add_footer(slide, num):
        add_text(slide, "Insights App", 0.3, 7.1, 3, 0.3, size=9, color=GREY)
        add_text(slide, str(num), 12.7, 7.1, 0.4, 0.3, size=9,
                 color=GREY, align=PP_ALIGN.RIGHT)

    return prs, blank, (BG, WHITE, GREY, BLUE, GREEN), set_bg, add_text, add_footer


def build_forecast_pptx(inputs, scenarios, ai_text=""):
    """Build the ROI Forecast PowerPoint. Returns BytesIO buffer."""
    prs, blank, colours, set_bg, add_text, add_footer = _pptx_base()
    _, WHITE, GREY, BLUE, GREEN = colours

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    set_bg(s1)
    add_text(s1, "Uber Ads ROI Forecast", 0.8, 1.4, 11.5, 1.0,
             size=38, bold=True, color=WHITE)
    add_text(s1, "Model and measure return on advertising investment for restaurant partners on Uber Eats.",
             0.8, 2.6, 10, 0.55, size=14, color=GREY)

    bar = s1.shapes.add_shape(1, Inches(0.8), Inches(3.3), Inches(4.5), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    campaign_types_str = ", ".join(inputs["campaign_types"]) or "None"
    dayparts_str       = ", ".join(inputs["dayparts"]) or "All dayparts"
    summary = (
        f"Category: {inputs['category']}   |   Monthly Budget: A${inputs['monthly_budget']:,.0f}   |   "
        f"AOV: A${inputs['avg_order_value']:.2f}   |   Locations: {inputs['num_locations']}   |   "
        f"Rating: {inputs['rating']:.1f}\n"
        f"Objective: {inputs['target_objective']}   |   Duration: {inputs['campaign_weeks']} weeks   |   "
        f"Campaign types: {campaign_types_str}   |   Dayparts: {dayparts_str}"
    )
    add_text(s1, summary, 0.8, 3.55, 12, 1.2, size=12, color=GREY)
    add_footer(s1, 1)

    # ── Slide 2: Three Scenarios ─────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    set_bg(s2)
    add_text(s2, "ROI Projections — Three Scenarios", 0.5, 0.25, 12, 0.6,
             size=22, bold=True, color=WHITE)

    col_pos    = [0.5, 4.7, 8.9]
    col_colors = [GREY, BLUE, GREEN]
    sc_names   = ["Conservative", "Base Case", "Optimistic"]

    for pos, col, name in zip(col_pos, col_colors, sc_names):
        add_text(s2, name, pos, 1.0, 3.8, 0.45, size=15, bold=True, color=col)

    metric_defs = [
        ("Total Orders",        "orders",      lambda v: f"{v:,.0f}"),
        ("Incremental Orders",  "incr_orders", lambda v: f"{v:,.0f}"),
        ("Revenue Driven",      "revenue",     lambda v: f"A${v:,.0f}"),
        ("ROAS",                "roas",        lambda v: f"{v:.2f}x"),
        ("CPO",                 "cpo",         lambda v: f"A${v:.2f}"),
        ("Profit on Ad Spend",  "profit",      lambda v: f"A${v:,.0f}"),
    ]
    row_tops = [1.6, 2.15, 2.7, 3.25, 3.8, 4.5]

    for j, (label, key, fmt) in enumerate(metric_defs):
        add_text(s2, label, 0.5, row_tops[j], 3.5, 0.42, size=10, color=GREY)
        for i, (pos, sname) in enumerate(zip(col_pos, sc_names)):
            val = scenarios[sname][key]
            vcol = col_colors[i]
            if key == "profit" and val < 0:
                vcol = RGBColor(0xEF, 0x44, 0x44)
            add_text(s2, fmt(val), pos, row_tops[j], 3.8, 0.42,
                     size=12, bold=True, color=vcol)
    add_footer(s2, 2)

    # ── Slide 3: Business Case (only if AI text exists) ──────────────────────
    if ai_text:
        s3 = prs.slides.add_slide(blank)
        set_bg(s3)
        add_text(s3, "Business Case", 0.5, 0.25, 12, 0.6,
                 size=22, bold=True, color=WHITE)
        display = ai_text[:1400] + ("…" if len(ai_text) > 1400 else "")
        add_text(s3, display, 0.5, 1.1, 12.3, 5.8, size=11, color=GREY)
        add_footer(s3, 3)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def build_outcome_pptx(inputs, metrics, health_score, ai_text=""):
    """Build the Campaign Outcome PowerPoint. Returns BytesIO buffer."""
    prs, blank, colours, set_bg, add_text, add_footer = _pptx_base()
    _, WHITE, GREY, BLUE, GREEN = colours
    RED = RGBColor(0xEF, 0x44, 0x44)

    bench = CATEGORY_BENCHMARKS[inputs["category"]]

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    set_bg(s1)
    restaurant = inputs.get("restaurant_name", "") or inputs["category"]
    add_text(s1, f"Campaign Outcome Report — {restaurant}", 0.8, 1.4, 11.5, 1.0,
             size=32, bold=True, color=WHITE)
    add_text(s1, f"Category: {inputs['category']}   |   "
             f"A${inputs['spend']:,.0f} spend   |   "
             f"{inputs['total_orders']:,} orders   |   "
             f"{inputs['incr_orders']:,} incremental orders",
             0.8, 2.6, 11, 0.5, size=13, color=GREY)

    bar = s1.shapes.add_shape(1, Inches(0.8), Inches(3.2), Inches(4.5), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    # Health score callout
    score_color = (GREEN if health_score >= 80 else
                   BLUE  if health_score >= 60 else
                   RGBColor(0xF5, 0x9E, 0x0B) if health_score >= 40 else RED)
    grade = ("A — Excellent" if health_score >= 80 else
             "B — Good"      if health_score >= 60 else
             "C — Fair"      if health_score >= 40 else
             "D — Needs Improvement")
    add_text(s1, f"Health Score: {health_score}/100  ({grade})",
             0.8, 3.5, 9, 0.5, size=15, bold=True, color=score_color)
    add_footer(s1, 1)

    # ── Slide 2: Calculated Metrics ──────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    set_bg(s2)
    add_text(s2, "Campaign Metrics", 0.5, 0.25, 12, 0.6,
             size=22, bold=True, color=WHITE)

    metric_labels = [
        ("Actual ROAS",         f"{metrics['actual_roas']:.2f}x"),
        ("Actual CPO",          f"A${metrics['actual_cpo']:.2f}"),
        ("Incremental CPO",     f"A${metrics['incr_cpo']:.2f}"),
        ("CTR",                 f"{metrics['ctr']:.2f}%"),
        ("Conversion Rate",     f"{metrics['conv_rate']:.2f}%"),
        ("Incremental Rate",    f"{metrics['incr_rate']:.1f}%"),
    ]
    for i, (lbl, val) in enumerate(metric_labels):
        col = i % 3
        row = i // 3
        add_text(s2, lbl, 0.5 + col * 4.3, 1.1 + row * 1.5, 3.8, 0.4,
                 size=10, color=GREY)
        add_text(s2, val, 0.5 + col * 4.3, 1.55 + row * 1.5, 3.8, 0.6,
                 size=22, bold=True, color=BLUE)

    # Benchmark comparison summary
    add_text(s2, "vs Category Benchmarks", 0.5, 4.2, 12, 0.4,
             size=14, bold=True, color=WHITE)
    bench_lines = (
        f"ROAS: {metrics['actual_roas']:.2f}x  vs  benchmark {bench['avg_roas']:.1f}x   |   "
        f"CPO: A${metrics['actual_cpo']:.2f}  vs  benchmark A${bench['avg_cpo']:.2f}   |   "
        f"CTR: {metrics['ctr']:.2f}%  vs  benchmark {bench['avg_ctr']:.1f}%   |   "
        f"Incremental Rate: {metrics['incr_rate']:.1f}%  vs  benchmark {bench['incr_rate']*100:.0f}%"
    )
    add_text(s2, bench_lines, 0.5, 4.7, 12.3, 0.8, size=11, color=GREY)
    add_footer(s2, 2)

    # ── Slide 3: AI Analysis (only if text exists) ────────────────────────────
    if ai_text:
        s3 = prs.slides.add_slide(blank)
        set_bg(s3)
        add_text(s3, "AI Campaign Analysis", 0.5, 0.25, 12, 0.6,
                 size=22, bold=True, color=WHITE)
        display = ai_text[:1400] + ("…" if len(ai_text) > 1400 else "")
        add_text(s3, display, 0.5, 1.1, 12.3, 5.8, size=11, color=GREY)
        add_footer(s3, 3)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ════════════════════════════════════════════════════════════════════════════════
# PAGE
# ════════════════════════════════════════════════════════════════════════════════

st.title("Uber Ads ROI Calculator")
st.markdown(
    "<p style='color:#6b7280;font-size:14px;margin-top:-12px;'>"
    "Model and measure return on advertising investment for restaurant partners on Uber Eats."
    "</p>",
    unsafe_allow_html=True,
)

# ── Two tabs ─────────────────────────────────────────────────────────────────────
tab_forecast, tab_outcome = st.tabs(["📈 ROI Forecast", "📊 Campaign Outcome"])


# ════════════════════════════════════════════════════════════════════════════════
# TAB 1 — ROI FORECAST (PRE-CAMPAIGN)
# ════════════════════════════════════════════════════════════════════════════════

with tab_forecast:

    # ── Section 1: Partner Inputs ────────────────────────────────────────────
    st.subheader("Partner Inputs")
    col_left, col_right = st.columns(2)

    with col_left:
        st.markdown(
            "<p style='font-weight:700;font-size:13px;color:#374151;margin-bottom:2px;'>"
            "Restaurant Profile</p>",
            unsafe_allow_html=True,
        )
        fc_category = st.selectbox(
            "Restaurant category",
            list(CATEGORY_BENCHMARKS.keys()),
            key="fc_category",
        )
        fc_aov = st.number_input(
            "Average order value, AUD",
            min_value=1.0, max_value=500.0,
            value=28.00, step=0.50, format="%.2f",
            key="fc_aov",
        )
        fc_organic = st.number_input(
            "Current monthly organic orders",
            min_value=0, max_value=1_000_000,
            value=500, step=10,
            key="fc_organic",
        )
        fc_locations = st.number_input(
            "Number of locations on Uber Eats",
            min_value=1, max_value=10_000,
            value=1, step=1,
            key="fc_locations",
        )
        fc_rating = st.slider(
            "Current Uber Eats rating",
            min_value=1.0, max_value=5.0,
            value=4.2, step=0.1,
            key="fc_rating",
        )

    with col_right:
        st.markdown(
            "<p style='font-weight:700;font-size:13px;color:#374151;margin-bottom:2px;'>"
            "Campaign Parameters</p>",
            unsafe_allow_html=True,
        )
        fc_budget = st.number_input(
            "Monthly ad budget, AUD",
            min_value=100.0, max_value=10_000_000.0,
            value=5_000.00, step=500.0, format="%.2f",
            key="fc_budget",
        )
        fc_campaign_types = st.multiselect(
            "Campaign type",
            ["Sponsored Listings", "Display Ads", "Homepage Banner", "Carousel Ads"],
            default=["Sponsored Listings"],
            key="fc_campaign_types",
        )
        fc_objective = st.selectbox(
            "Target objective",
            ["Maximise Orders", "Maximise ROAS", "New Customer Acquisition", "Reactivation"],
            key="fc_objective",
        )
        fc_weeks = st.slider(
            "Campaign duration in weeks",
            min_value=1, max_value=12,
            value=4, step=1,
            key="fc_weeks",
        )
        fc_dayparts = st.multiselect(
            "Daypart focus",
            ["Breakfast", "Lunch", "Afternoon", "Dinner", "Late Night"],
            default=["Lunch", "Dinner"],
            key="fc_dayparts",
        )

    st.divider()

    # ── Section 3: Calculate ROI Forecast button (centred) ───────────────────
    _, btn_col, _ = st.columns([1, 2, 1])
    with btn_col:
        fc_calculate = st.button(
            "Calculate ROI Forecast",
            type="primary",
            use_container_width=True,
            key="fc_calculate_btn",
        )

    if fc_calculate:
        fc_results, fc_adj_roas, fc_base_cpo = compute_roi(
            category        = fc_category,
            avg_order_value = fc_aov,
            num_locations   = fc_locations,
            rating          = fc_rating,
            monthly_budget  = fc_budget,
            campaign_types  = fc_campaign_types,
            dayparts        = fc_dayparts,
        )
        st.session_state["fc_calculated"] = True
        st.session_state["fc_results"]    = fc_results
        st.session_state["fc_adj_roas"]   = fc_adj_roas
        st.session_state["fc_base_cpo"]   = fc_base_cpo
        st.session_state["fc_inputs"]     = {
            "category":         fc_category,
            "avg_order_value":  fc_aov,
            "organic_orders":   fc_organic,
            "num_locations":    fc_locations,
            "rating":           fc_rating,
            "monthly_budget":   fc_budget,
            "campaign_types":   fc_campaign_types,
            "target_objective": fc_objective,
            "campaign_weeks":   fc_weeks,
            "dayparts":         fc_dayparts,
        }
        # Clear stale AI / PPTX whenever inputs change
        st.session_state.pop("fc_ai_case", None)
        st.session_state.pop("fc_pptx",    None)

    # ── Results ──────────────────────────────────────────────────────────────
    if st.session_state.get("fc_calculated"):
        fc_res    = st.session_state["fc_results"]
        fc_inputs = st.session_state["fc_inputs"]
        fc_roas   = st.session_state["fc_adj_roas"]
        fc_cpo    = st.session_state["fc_base_cpo"]

        st.divider()

        # ── Section 4: Three Scenario Output ─────────────────────────────────
        st.subheader("ROI Projections")

        sc1, sc2, sc3 = st.columns(3)
        scenario_cols = {
            "Conservative": sc1,
            "Base Case":    sc2,
            "Optimistic":   sc3,
        }

        for name, col in scenario_cols.items():
            s                        = fc_res[name]
            hdr_bg, hdr_txt, border  = SCENARIO_COLOURS[name]
            profit_col = "#10B981" if s["profit"] >= 0 else "#EF4444"

            with col:
                st.markdown(
                    f"<div style='border:{border};border-radius:12px;overflow:hidden;"
                    f"box-shadow:0 4px 12px rgba(0,0,0,0.08);background:#FFFFFF;'>"
                    f"<div style='background:{hdr_bg};padding:14px 16px;'>"
                    f"<p style='color:{hdr_txt};font-weight:700;font-size:15px;margin:0;'>"
                    f"{name}</p></div>"
                    f"<div style='padding:14px 16px;'>"
                    + _metric_row("Total Orders",       f"{s['orders']:,.0f}")
                    + _metric_row("Incremental Orders", f"{s['incr_orders']:,.0f}")
                    + _metric_row("Revenue Driven",     f"A${s['revenue']:,.0f}")
                    + _metric_row("ROAS",               f"{s['roas']:.2f}x")
                    + _metric_row("CPO",                f"A${s['cpo']:.2f}")
                    + "<div style='border-top:1px solid #E5E7EB;margin:8px 0;'></div>"
                    + _metric_row("Profit on Ad Spend", f"A${s['profit']:,.0f}",
                                  value_color=profit_col)
                    + "</div></div>",
                    unsafe_allow_html=True,
                )

        st.divider()

        # ── Section 5: Break-Even Analysis ───────────────────────────────────
        st.subheader("Break-Even Analysis")

        # With a 30% gross margin, ROAS must exceed 1/0.30 = 3.33x to break even
        break_even_roas = round(1 / 0.30, 2)
        base_roas       = fc_res["Base Case"]["roas"]
        base_cpo_disp   = fc_res["Base Case"]["cpo"]
        aov             = fc_inputs["avg_order_value"]

        st.markdown(
            f"<div style='background:#FFFFFF;border-radius:12px;padding:18px 22px;"
            f"box-shadow:0 4px 12px rgba(0,0,0,0.08);font-size:14px;color:#374151;'>"
            f"At your average order value of <strong>A${aov:.2f}</strong> and a base-case CPO of "
            f"<strong>A${base_cpo_disp:.2f}</strong>, you break even when ROAS exceeds "
            f"<strong>{break_even_roas:.2f}x</strong>. "
            f"Your projected base-case ROAS is <strong>{base_roas:.2f}x</strong>."
            f"</div>",
            unsafe_allow_html=True,
        )

        # Gauge: projected ROAS vs break-even threshold
        gauge_max = max(base_roas * 1.5, break_even_roas * 1.5, 6.0)
        gauge_fig = go.Figure(go.Indicator(
            mode="gauge+number+delta",
            value=base_roas,
            number={"suffix": "x", "font": {"size": 28, "family": "Inter, sans-serif"}},
            delta={
                "reference": break_even_roas,
                "increasing": {"color": "#10B981"},
                "decreasing": {"color": "#EF4444"},
                "suffix": "x vs break-even",
            },
            title={"text": "Projected ROAS (Base Case)", "font": {"size": 14}},
            gauge={
                "axis": {
                    "range": [0, gauge_max],
                    "tickwidth": 1,
                    "tickcolor": "#6B7280",
                    "tickfont": {"size": 11},
                },
                "bar": {"color": "#7C3AED", "thickness": 0.3},
                "bgcolor": "white",
                "borderwidth": 0,
                "steps": [
                    {"range": [0,               break_even_roas], "color": "#FEF2F2"},
                    {"range": [break_even_roas,  gauge_max],       "color": "#ECFDF5"},
                ],
                "threshold": {
                    "line": {"color": "#EF4444", "width": 3},
                    "thickness": 0.75,
                    "value": break_even_roas,
                },
            },
        ))
        gauge_fig.update_layout(
            height=280,
            margin=dict(t=40, b=20, l=30, r=30),
            paper_bgcolor="rgba(0,0,0,0)",
            font=dict(family="Inter, sans-serif", color="#374151"),
        )
        st.plotly_chart(gauge_fig, use_container_width=True)

        st.divider()

        # ── Section 6: Reverse Calculator ────────────────────────────────────
        st.subheader("Reverse Calculator")
        reverse_on = st.toggle(
            "I know my target — what budget do I need?",
            key="fc_reverse_toggle",
        )

        if reverse_on:
            rev1, rev2 = st.columns(2)
            with rev1:
                rev_target_orders = st.number_input(
                    "Target incremental orders per month",
                    min_value=1, max_value=1_000_000,
                    value=200, step=10,
                    key="fc_rev_orders",
                )
            with rev2:
                rev_target_roas = st.number_input(
                    "Target ROAS",
                    min_value=0.1, max_value=50.0,
                    value=3.0, step=0.1, format="%.1f",
                    key="fc_rev_roas",
                )

            # Derive totals: incremental orders = total * incr_rate
            incr_rate     = CATEGORY_BENCHMARKS[fc_inputs["category"]]["incr_rate"]
            rev_tot_orders = rev_target_orders / incr_rate if incr_rate > 0 else 0
            rev_revenue    = rev_tot_orders * fc_inputs["avg_order_value"]
            rev_budget     = rev_revenue / rev_target_roas if rev_target_roas > 0 else 0

            st.markdown(
                f"<div style='background:#FFFFFF;border-radius:12px;padding:18px 22px;"
                f"box-shadow:0 4px 12px rgba(0,0,0,0.08);margin-top:12px;'>"
                f"<p style='font-weight:700;font-size:14px;color:#111827;margin-bottom:10px;'>"
                f"Required Budget &amp; Projections</p>"
                + _metric_row("Required monthly budget", f"A${rev_budget:,.0f}")
                + _metric_row("Expected total orders",   f"{rev_tot_orders:,.0f}")
                + _metric_row("Projected revenue",       f"A${rev_revenue:,.0f}")
                + "</div>",
                unsafe_allow_html=True,
            )

        st.divider()

        # ── Section 7: Daypart Budget Recommendation ──────────────────────────
        st.subheader("Daypart Budget Recommendation")
        selected_dayparts = fc_inputs["dayparts"]
        fc_bud            = fc_inputs["monthly_budget"]
        base_orders       = fc_res["Base Case"]["orders"]

        if not selected_dayparts:
            st.info("No dayparts selected — add daypart focus in the inputs above to see a recommended split.")
        else:
            st.markdown(
                "<p style='color:#6b7280;font-size:14px;margin-top:-8px;'>"
                "Based on your category and target objective, we recommend:</p>",
                unsafe_allow_html=True,
            )
            raw_w  = {d: DAYPART_BASE_WEIGHTS[d] for d in selected_dayparts}
            total_w = sum(raw_w.values())
            norm_w  = {d: w / total_w for d, w in raw_w.items()}

            dp_rows = [
                {
                    "Daypart":         d,
                    "Budget %":        f"{pct * 100:.1f}%",
                    "Budget (AUD)":    f"A${fc_bud * pct:,.0f}",
                    "Expected Orders": f"{base_orders * pct:,.0f}",
                }
                for d, pct in norm_w.items()
            ]
            df_dp = pd.DataFrame(dp_rows)
            df_dp.index = range(1, len(df_dp) + 1)
            st.dataframe(df_dp, use_container_width=True)

        st.divider()

        # ── Section 8: AI Business Case ───────────────────────────────────────
        st.subheader("AI Business Case")
        api_key = get_api_key()

        if not api_key:
            st.warning("Add `ANTHROPIC_API_KEY` to Streamlit secrets to enable the AI Business Case.")
        else:
            if st.button("Generate Business Case", type="primary", key="fc_ai_btn"):
                base = fc_res["Base Case"]
                cons = fc_res["Conservative"]
                opti = fc_res["Optimistic"]

                prompt = (
                    f"You are a Partner Manager at Uber Advertising ANZ building a business case "
                    f"for a {fc_inputs['category']} restaurant partner to invest in Uber Ads. "
                    f"Here are the calculated projections:\n\n"
                    f"PARTNER PROFILE\n"
                    f"- Category: {fc_inputs['category']}\n"
                    f"- Average order value: A${fc_inputs['avg_order_value']:.2f}\n"
                    f"- Current monthly organic orders: {fc_inputs['organic_orders']:,}\n"
                    f"- Number of Uber Eats locations: {fc_inputs['num_locations']}\n"
                    f"- Uber Eats rating: {fc_inputs['rating']:.1f}\n\n"
                    f"CAMPAIGN PARAMETERS\n"
                    f"- Monthly ad budget: A${fc_inputs['monthly_budget']:,.0f}\n"
                    f"- Campaign types: {', '.join(fc_inputs['campaign_types']) or 'None'}\n"
                    f"- Objective: {fc_inputs['target_objective']}\n"
                    f"- Duration: {fc_inputs['campaign_weeks']} weeks\n"
                    f"- Daypart focus: {', '.join(fc_inputs['dayparts']) or 'All dayparts'}\n\n"
                    f"PROJECTIONS\n"
                    f"Conservative — {cons['orders']:,.0f} total orders, "
                    f"{cons['incr_orders']:,.0f} incremental, "
                    f"A${cons['revenue']:,.0f} revenue, {cons['roas']:.2f}x ROAS, "
                    f"A${cons['cpo']:.2f} CPO\n"
                    f"Base Case    — {base['orders']:,.0f} total orders, "
                    f"{base['incr_orders']:,.0f} incremental, "
                    f"A${base['revenue']:,.0f} revenue, {base['roas']:.2f}x ROAS, "
                    f"A${base['cpo']:.2f} CPO\n"
                    f"Optimistic   — {opti['orders']:,.0f} total orders, "
                    f"{opti['incr_orders']:,.0f} incremental, "
                    f"A${opti['revenue']:,.0f} revenue, {opti['roas']:.2f}x ROAS, "
                    f"A${opti['cpo']:.2f} CPO\n\n"
                    f"Write a compelling 3-paragraph business case that:\n"
                    f"1. Opens with the revenue opportunity using base case numbers\n"
                    f"2. Explains why incremental orders matter more than total orders\n"
                    f"3. Closes with a clear recommended starting budget and expected return\n\n"
                    f"Write in a consultative, commercially confident tone suitable for "
                    f"presenting to a restaurant partner's marketing director. "
                    f"All currency values are in AUD (A$ symbol)."
                )

                ai_client = anthropic.Anthropic(api_key=api_key)
                with st.spinner("Generating business case…"):
                    msg = ai_client.messages.create(
                        model="claude-sonnet-4-6",
                        max_tokens=700,
                        system=(
                            "You are a Partner Manager at Uber Advertising ANZ. "
                            "Write commercially confident, data-driven business cases. "
                            "Be specific with numbers. Use Australian English."
                        ),
                        messages=[{"role": "user", "content": prompt}],
                    )
                st.session_state["fc_ai_case"] = msg.content[0].text.strip()

            if st.session_state.get("fc_ai_case"):
                st.markdown(
                    "<div style='background:#FFFFFF;border-radius:12px;padding:22px 26px;"
                    "box-shadow:0 4px 12px rgba(0,0,0,0.08);margin-top:14px;font-size:14px;"
                    "line-height:1.8;color:#374151;'>"
                    + st.session_state["fc_ai_case"].replace("\n", "<br>") +
                    "</div>",
                    unsafe_allow_html=True,
                )

        st.divider()

        # ── Section 9: Export ─────────────────────────────────────────────────
        st.subheader("Export")
        if st.button("Download Business Case (.pptx)", key="fc_pptx_btn"):
            ai_text  = st.session_state.get("fc_ai_case", "")
            pptx_buf = build_forecast_pptx(fc_inputs, fc_res, ai_text)
            st.session_state["fc_pptx"] = pptx_buf

        if st.session_state.get("fc_pptx"):
            filename = f"{_date.today().isoformat()}_uber_ads_roi_forecast.pptx"
            st.download_button(
                label="Save .pptx",
                data=st.session_state["fc_pptx"],
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="fc_pptx_download",
            )


# ════════════════════════════════════════════════════════════════════════════════
# TAB 2 — CAMPAIGN OUTCOME (POST-CAMPAIGN)
# ════════════════════════════════════════════════════════════════════════════════

with tab_outcome:

    # ── Section 1: Campaign Actuals Input ────────────────────────────────────
    st.subheader("Campaign Actuals")

    oc_left, oc_right = st.columns(2)

    with oc_left:
        st.markdown(
            "<p style='font-weight:700;font-size:13px;color:#374151;margin-bottom:2px;'>"
            "Campaign Details</p>",
            unsafe_allow_html=True,
        )
        oc_restaurant = st.text_input(
            "Restaurant name",
            placeholder="e.g. Grill'd Healthy Burgers",
            key="oc_restaurant",
        )
        oc_category = st.selectbox(
            "Restaurant category",
            list(CATEGORY_BENCHMARKS.keys()),
            key="oc_category",
        )
        oc_aov = st.number_input(
            "Average order value, AUD",
            min_value=1.0, max_value=500.0,
            value=28.00, step=0.50, format="%.2f",
            key="oc_aov",
        )
        oc_campaign_types = st.multiselect(
            "Campaign type",
            ["Sponsored Listings", "Display Ads", "Homepage Banner", "Carousel Ads"],
            default=["Sponsored Listings"],
            key="oc_campaign_types",
        )
        oc_start = st.date_input("Campaign start date", key="oc_start")
        oc_end   = st.date_input("Campaign end date",   key="oc_end")

    with oc_right:
        st.markdown(
            "<p style='font-weight:700;font-size:13px;color:#374151;margin-bottom:2px;'>"
            "Actual Results</p>",
            unsafe_allow_html=True,
        )

        # ── File uploader ─────────────────────────────────────────────────────
        # Accepts CSV or XLSX Uber Ads exports. Auto-detects and sums metric columns.
        # Falls back to manual inputs when no file is uploaded.
        oc_file = st.file_uploader(
            "Upload Uber Ads campaign export (CSV)",
            type=["csv", "xlsx"],
            key="oc_upload",
        )

        if oc_file:
            # Use filename + size as a stable ID so override widget keys change
            # whenever a different file is uploaded, forcing them to reset.
            file_id = f"{oc_file.name}_{oc_file.size}"
            detected, col_found = parse_campaign_file(oc_file)

            if not detected:
                st.error(
                    "Could not read this file. Please check it is a valid CSV or XLSX "
                    "and that column headers are on the first row."
                )
                # Fallback manual inputs when file is unreadable
                oc_spend        = st.number_input("Total ad spend, AUD",      min_value=0.0, max_value=10_000_000.0, value=0.0,   step=100.0, format="%.2f", key=f"oc_spend_{file_id}")
                oc_total_orders = st.number_input("Total orders driven",       min_value=0,   max_value=1_000_000,     value=0,     step=10,    key=f"oc_total_orders_{file_id}")
                oc_incr_orders  = st.number_input("Incremental orders",        min_value=0,   max_value=1_000_000,     value=0,     step=10,    key=f"oc_incr_orders_{file_id}", help="Orders that would not have occurred without advertising")
                oc_impressions  = st.number_input("Total impressions",         min_value=0,   max_value=1_000_000_000, value=0,     step=1_000, key=f"oc_impressions_{file_id}")
                oc_clicks       = st.number_input("Total clicks",              min_value=0,   max_value=10_000_000,    value=0,     step=100,   key=f"oc_clicks_{file_id}")
            else:
                # ── Detected values preview ───────────────────────────────────
                # Show which value was found for each metric and from which column
                def _preview_row(metric, label, fmt):
                    if metric in detected:
                        src = col_found.get(metric, "")
                        src_note = f" <span style='color:#6B7280;font-size:11px;'>from \"{src}\"</span>" if src else ""
                        return (
                            f"<div style='display:flex;justify-content:space-between;"
                            f"padding:4px 0;border-bottom:1px solid #D1FAE5;'>"
                            f"<span style='font-size:13px;color:#065F46;'>{label}{src_note}</span>"
                            f"<span style='font-weight:600;font-size:13px;color:#065F46;'>{fmt(detected[metric])}</span>"
                            f"</div>"
                        )
                    return (
                        f"<div style='display:flex;justify-content:space-between;"
                        f"padding:4px 0;border-bottom:1px solid #D1FAE5;'>"
                        f"<span style='font-size:13px;color:#9CA3AF;'>{label}</span>"
                        f"<span style='font-size:12px;color:#9CA3AF;font-style:italic;'>Not detected</span>"
                        f"</div>"
                    )

                preview_html = (
                    "<div style='background:#ECFDF5;border:1px solid #6EE7B7;border-radius:10px;"
                    "padding:14px 16px;margin-bottom:12px;'>"
                    "<p style='font-weight:700;font-size:13px;color:#065F46;margin:0 0 10px 0;'>"
                    "✅ Detected from upload:</p>"
                    + _preview_row("spend",        "Spend",               lambda v: f"A${v:,.2f}")
                    + _preview_row("total_orders",  "Total Orders",        lambda v: f"{v:,.0f}")
                    + _preview_row("incr_orders",   "Incremental Orders",  lambda v: f"{v:,.0f}")
                    + _preview_row("impressions",   "Impressions",         lambda v: f"{v:,.0f}")
                    + _preview_row("clicks",        "Clicks",              lambda v: f"{v:,.0f}")
                    + "</div>"
                )
                st.markdown(preview_html, unsafe_allow_html=True)

                # ── Override expander ─────────────────────────────────────────
                # Pre-filled with detected values. User can edit if the export
                # needs correction (e.g. partial date range, test rows, etc.).
                with st.expander("✏️ Edit detected values"):
                    oc_spend        = st.number_input("Spend (AUD)",          min_value=0.0, max_value=10_000_000.0, value=float(detected.get("spend",        0.0)), step=100.0, format="%.2f", key=f"oc_spend_{file_id}")
                    oc_total_orders = st.number_input("Total orders",          min_value=0,   max_value=1_000_000,     value=int(detected.get("total_orders",  0)),   step=10,    key=f"oc_total_orders_{file_id}")
                    oc_incr_orders  = st.number_input("Incremental orders",    min_value=0,   max_value=1_000_000,     value=int(detected.get("incr_orders",   0)),   step=10,    key=f"oc_incr_orders_{file_id}", help="Orders that would not have occurred without advertising")
                    oc_impressions  = st.number_input("Impressions",           min_value=0,   max_value=1_000_000_000, value=int(detected.get("impressions",   0)),   step=1_000, key=f"oc_impressions_{file_id}")
                    oc_clicks       = st.number_input("Clicks",                min_value=0,   max_value=10_000_000,    value=int(detected.get("clicks",        0)),   step=100,   key=f"oc_clicks_{file_id}")

        else:
            # ── No file uploaded — manual inputs ──────────────────────────────
            oc_spend = st.number_input(
                "Total ad spend, AUD",
                min_value=0.0, max_value=10_000_000.0,
                value=5_000.00, step=100.0, format="%.2f",
                key="oc_spend_manual",
            )
            oc_total_orders = st.number_input(
                "Total orders driven",
                min_value=0, max_value=1_000_000,
                value=500, step=10,
                key="oc_total_orders_manual",
            )
            oc_incr_orders = st.number_input(
                "Incremental orders",
                min_value=0, max_value=1_000_000,
                value=290, step=10,
                key="oc_incr_orders_manual",
                help="Orders that would not have occurred without advertising",
            )
            oc_impressions = st.number_input(
                "Total impressions",
                min_value=0, max_value=1_000_000_000,
                value=250_000, step=1_000,
                key="oc_impressions_manual",
            )
            oc_clicks = st.number_input(
                "Total clicks",
                min_value=0, max_value=10_000_000,
                value=7_500, step=100,
                key="oc_clicks_manual",
            )

        # ── Campaign targets — always manual (from the original brief) ────────
        st.markdown(
            "<p style='font-weight:700;font-size:13px;color:#374151;"
            "margin-top:16px;margin-bottom:2px;'>Campaign Targets</p>",
            unsafe_allow_html=True,
        )
        oc_budget_target = st.number_input(
            "Original budget target, AUD",
            min_value=0.0, max_value=10_000_000.0,
            value=5_000.00, step=100.0, format="%.2f",
            key="oc_budget_target",
        )
        oc_order_target = st.number_input(
            "Original order target",
            min_value=0, max_value=1_000_000,
            value=450, step=10,
            key="oc_order_target",
        )
        oc_roas_target = st.number_input(
            "Original ROAS target",
            min_value=0.0, max_value=50.0,
            value=3.0, step=0.1, format="%.1f",
            key="oc_roas_target",
        )

    st.divider()

    # ── Section 2: Analyse button (centred) ──────────────────────────────────
    _, oc_btn_col, _ = st.columns([1, 2, 1])
    with oc_btn_col:
        oc_analyse = st.button(
            "Analyse Campaign Outcome",
            type="primary",
            use_container_width=True,
            key="oc_analyse_btn",
        )

    if oc_analyse:
        oc_inputs = {
            "restaurant_name":  oc_restaurant,
            "category":         oc_category,
            "avg_order_value":  oc_aov,
            "campaign_types":   oc_campaign_types,
            "start_date":       oc_start,
            "end_date":         oc_end,
            "spend":            oc_spend,
            "total_orders":     oc_total_orders,
            "incr_orders":      oc_incr_orders,
            "impressions":      oc_impressions,
            "clicks":           oc_clicks,
            "budget_target":    oc_budget_target,
            "target_orders":    oc_order_target,
            "target_roas":      oc_roas_target,
        }
        oc_metrics   = compute_outcome_metrics(oc_inputs)
        oc_bench     = CATEGORY_BENCHMARKS[oc_category]
        oc_health    = compute_health_score(oc_metrics, oc_inputs, oc_bench)

        st.session_state["oc_calculated"] = True
        st.session_state["oc_inputs"]     = oc_inputs
        st.session_state["oc_metrics"]    = oc_metrics
        st.session_state["oc_health"]     = oc_health
        st.session_state.pop("oc_ai_text", None)
        st.session_state.pop("oc_pptx",    None)

    # ── Results ──────────────────────────────────────────────────────────────
    if st.session_state.get("oc_calculated"):
        oc_inp  = st.session_state["oc_inputs"]
        oc_met  = st.session_state["oc_metrics"]
        oc_h    = st.session_state["oc_health"]
        oc_b    = CATEGORY_BENCHMARKS[oc_inp["category"]]

        st.divider()

        # ── Section 3: Calculated Metrics ────────────────────────────────────
        st.subheader("Calculated Metrics")

        m1, m2, m3 = st.columns(3)
        m4, m5, m6 = st.columns(3)

        m1.metric("Actual ROAS",       f"{oc_met['actual_roas']:.2f}x")
        m2.metric("Actual CPO",        f"A${oc_met['actual_cpo']:.2f}")
        m3.metric("Incremental CPO",   f"A${oc_met['incr_cpo']:.2f}")
        m4.metric("CTR",               f"{oc_met['ctr']:.2f}%")
        m5.metric("Conversion Rate",   f"{oc_met['conv_rate']:.2f}%")
        m6.metric("Incremental Rate",  f"{oc_met['incr_rate']:.1f}%")

        st.divider()

        # ── Section 4: Target vs Actual Scorecard ─────────────────────────────
        st.subheader("Target vs Actual Scorecard")

        # Derive implied targets for metrics without a direct target input
        impl_incr_target = oc_inp["target_orders"] * oc_b["incr_rate"]
        impl_cpo_target  = (oc_inp["budget_target"] / oc_inp["target_orders"]
                            if oc_inp["target_orders"] > 0 else 0)

        scorecard_rows_def = [
            # (label, target, actual, fmt_fn, lower_is_better)
            ("Ad Spend (AUD)", oc_inp["budget_target"], oc_inp["spend"],
             lambda t, a: (f"A${t:,.0f}", f"A${a:,.0f}", f"A${a - t:+,.0f}"), True),
            ("Total Orders",   oc_inp["target_orders"], oc_inp["total_orders"],
             lambda t, a: (f"{t:,.0f}", f"{a:,.0f}", f"{a - t:+,.0f}"), False),
            ("Incremental Orders", impl_incr_target, oc_inp["incr_orders"],
             lambda t, a: (f"{t:,.0f}", f"{a:,.0f}", f"{a - t:+,.0f}"), False),
            ("ROAS",           oc_inp["target_roas"], oc_met["actual_roas"],
             lambda t, a: (f"{t:.2f}x", f"{a:.2f}x", f"{a - t:+.2f}x"), False),
            ("CPO (AUD)",      impl_cpo_target, oc_met["actual_cpo"],
             lambda t, a: (f"A${t:.2f}", f"A${a:.2f}", f"A${a - t:+.2f}"), True),
        ]

        sc_rows_html = []
        for label, target, actual, fmt_fn, lower_better in scorecard_rows_def:
            t_str, a_str, v_str = fmt_fn(target, actual)
            emoji, status_color = _get_status(actual, target, lower_is_better=lower_better)
            sc_rows_html.append(_table_row(
                [label, t_str, a_str, v_str, emoji],
                colors=["#374151", "#374151", "#374151", status_color, status_color],
            ))

        st.markdown(
            _html_table(["Metric", "Target", "Actual", "Variance", "Status"], sc_rows_html),
            unsafe_allow_html=True,
        )

        st.divider()

        # ── Section 5: Benchmark Comparison ──────────────────────────────────
        st.subheader("Benchmark Comparison")

        bench_rows_def = [
            # (label, actual, benchmark, unit, lower_is_better)
            ("ROAS",             oc_met["actual_roas"], oc_b["avg_roas"],          "x",  False),
            ("CPO (AUD)",        oc_met["actual_cpo"],  oc_b["avg_cpo"],           "$",  True),
            ("CTR (%)",          oc_met["ctr"],         oc_b["avg_ctr"],           "%",  False),
            ("Incremental Rate", oc_met["incr_rate"],   oc_b["incr_rate"] * 100,   "%",  False),
        ]

        bench_rows_html = []
        for label, actual, bench_val, unit, lower_better in bench_rows_def:
            # Format display values
            if unit == "$":
                a_disp = f"A${actual:.2f}"
                b_disp = f"A${bench_val:.2f}"
            elif unit == "x":
                a_disp = f"{actual:.2f}x"
                b_disp = f"{bench_val:.1f}x"
            else:
                a_disp = f"{actual:.2f}%"
                b_disp = f"{bench_val:.1f}%"

            # Variance vs benchmark (positive = above = usually good)
            if lower_better:
                diff_pct = (bench_val - actual) / bench_val * 100 if bench_val > 0 else 0
            else:
                diff_pct = (actual - bench_val) / bench_val * 100 if bench_val > 0 else 0

            vs_color  = "#10B981" if diff_pct >= 0 else "#EF4444"
            vs_label  = f"{diff_pct:+.1f}%"

            bench_rows_html.append(_table_row(
                [label, a_disp, b_disp, vs_label],
                colors=["#374151", "#374151", "#374151", vs_color],
            ))

        st.markdown(
            _html_table(["Metric", "Your Result", "Category Benchmark", "vs Benchmark"],
                        bench_rows_html),
            unsafe_allow_html=True,
        )

        st.divider()

        # ── Section 6: Performance Visualisations ────────────────────────────
        st.subheader("Performance Visualisations")

        chart_left, chart_right = st.columns(2)

        # LEFT: Donut — Incremental vs Organic order split
        with chart_left:
            organic_orders = max(oc_inp["total_orders"] - oc_inp["incr_orders"], 0)
            donut_fig = go.Figure(go.Pie(
                labels=["Incremental Orders", "Organic (Ad-Influenced)"],
                values=[oc_inp["incr_orders"], organic_orders],
                hole=0.52,
                marker=dict(colors=["#7C3AED", "#E5E7EB"]),
                textinfo="label+percent",
                textfont=dict(family="Inter, sans-serif", size=11),
                hovertemplate="%{label}: %{value:,} (%{percent})<extra></extra>",
            ))
            donut_fig.update_layout(
                title=dict(text="Order Split", font=dict(size=13, family="Inter, sans-serif")),
                showlegend=False,
                height=300,
                margin=dict(t=50, b=20, l=20, r=20),
                paper_bgcolor="rgba(0,0,0,0)",
                font=dict(family="Inter, sans-serif"),
            )
            st.plotly_chart(donut_fig, use_container_width=True)

        # RIGHT: Grouped bars — Target vs Actual (ROAS, CPO, Orders)
        # Each metric is plotted in its own subplot so units don't conflict
        with chart_right:
            t_roas   = oc_inp["target_roas"]
            a_roas   = oc_met["actual_roas"]
            t_cpo    = (oc_inp["budget_target"] / oc_inp["target_orders"]
                        if oc_inp["target_orders"] > 0 else 0)
            a_cpo    = oc_met["actual_cpo"]
            t_orders = oc_inp["target_orders"]
            a_orders = oc_inp["total_orders"]

            bar_fig = make_subplots(
                rows=1, cols=3,
                subplot_titles=["ROAS", "CPO (AUD)", "Total Orders"],
            )
            bar_colours = ["#E5E7EB", "#7C3AED"]   # Target = grey, Actual = purple

            for col_i, (t_val, a_val, suffix) in enumerate([
                (t_roas,   a_roas,   "x"),
                (t_cpo,    a_cpo,    "$"),
                (t_orders, a_orders, ""),
            ], start=1):
                bar_fig.add_trace(
                    go.Bar(
                        name="Target",
                        x=["Target"],
                        y=[t_val],
                        marker_color=bar_colours[0],
                        showlegend=(col_i == 1),
                        hovertemplate=f"Target: {t_val:.2f}{suffix}<extra></extra>",
                    ),
                    row=1, col=col_i,
                )
                bar_fig.add_trace(
                    go.Bar(
                        name="Actual",
                        x=["Actual"],
                        y=[a_val],
                        marker_color=bar_colours[1],
                        showlegend=(col_i == 1),
                        hovertemplate=f"Actual: {a_val:.2f}{suffix}<extra></extra>",
                    ),
                    row=1, col=col_i,
                )

            bar_fig.update_layout(
                title=dict(text="Target vs Actual",
                           font=dict(size=13, family="Inter, sans-serif")),
                height=300,
                barmode="group",
                showlegend=True,
                legend=dict(orientation="h", y=-0.15),
                margin=dict(t=50, b=40, l=20, r=20),
                paper_bgcolor="rgba(0,0,0,0)",
                plot_bgcolor="rgba(0,0,0,0)",
                font=dict(family="Inter, sans-serif"),
            )
            bar_fig.update_traces(width=0.5)
            st.plotly_chart(bar_fig, use_container_width=True)

        st.divider()

        # ── Section 7: Campaign Health Score ──────────────────────────────────
        st.subheader("Campaign Health Score")

        # Grade and colour based on score
        if oc_h >= 80:
            grade, score_bg, score_txt = "A", "#ECFDF5", "#059669"
            grade_label = "Excellent"
        elif oc_h >= 60:
            grade, score_bg, score_txt = "B", "#EFF6FF", "#2563EB"
            grade_label = "Good"
        elif oc_h >= 40:
            grade, score_bg, score_txt = "C", "#FFFBEB", "#D97706"
            grade_label = "Fair"
        else:
            grade, score_bg, score_txt = "D", "#FEF2F2", "#DC2626"
            grade_label = "Needs Improvement"

        score_col, detail_col = st.columns([1, 3])

        with score_col:
            st.markdown(
                f"<div style='background:{score_bg};border-radius:16px;padding:28px 20px;"
                f"text-align:center;box-shadow:0 4px 12px rgba(0,0,0,0.08);'>"
                f"<p style='font-size:56px;font-weight:900;color:{score_txt};"
                f"margin:0;line-height:1;'>{oc_h}</p>"
                f"<p style='font-size:13px;color:#6B7280;margin:4px 0 0 0;'>out of 100</p>"
                f"<p style='font-size:20px;font-weight:700;color:{score_txt};"
                f"margin:8px 0 0 0;'>Grade {grade}</p>"
                f"<p style='font-size:13px;color:#6B7280;margin:2px 0 0 0;'>{grade_label}</p>"
                f"</div>",
                unsafe_allow_html=True,
            )

        with detail_col:
            # Show the four weighted dimensions
            bench_incr_pct = oc_b["incr_rate"] * 100
            impl_cpo_t = (oc_inp["budget_target"] / oc_inp["target_orders"]
                          if oc_inp["target_orders"] > 0 else 0)

            def _dim_bar(label, actual, target, pct_weight, lower_better=False):
                if target <= 0:
                    ratio = 0.5
                else:
                    ratio = (target / actual if lower_better and actual > 0
                             else actual / target)
                ratio = min(ratio, 1.5) / 1.5   # cap at 1.5x for display
                bar_w = round(ratio * 100)
                color = "#10B981" if ratio >= 0.667 else "#F59E0B" if ratio >= 0.5 else "#EF4444"
                return (
                    f"<div style='margin-bottom:12px;'>"
                    f"<div style='display:flex;justify-content:space-between;"
                    f"font-size:12px;color:#374151;margin-bottom:4px;'>"
                    f"<span><strong>{label}</strong> ({pct_weight}% of score)</span>"
                    f"<span style='color:{color};font-weight:600;'>"
                    f"{'A$' if lower_better else ''}{actual:.2f}{'x' if not lower_better and label == 'ROAS' else '%' if '%' in label else ''} "
                    f"/ target {'A$' if lower_better else ''}{target:.2f}{'x' if not lower_better and label == 'ROAS' else '%' if '%' in label else ''}</span>"
                    f"</div>"
                    f"<div style='background:#F3F4F6;border-radius:4px;height:8px;'>"
                    f"<div style='background:{color};border-radius:4px;height:8px;"
                    f"width:{bar_w}%;'></div></div></div>"
                )

            st.markdown(
                "<div style='background:#FFFFFF;border-radius:12px;padding:20px 24px;"
                "box-shadow:0 4px 12px rgba(0,0,0,0.08);'>"
                "<p style='font-size:13px;font-weight:600;color:#374151;margin-bottom:14px;'>"
                "Score Breakdown</p>"
                + _dim_bar("ROAS",             oc_met["actual_roas"], oc_inp["target_roas"],       30, False)
                + _dim_bar("Total Orders",      oc_inp["total_orders"], oc_inp["target_orders"],    25, False)
                + _dim_bar("CPO (AUD)",         oc_met["actual_cpo"], impl_cpo_t,                   25, True)
                + _dim_bar("Incremental Rate %",oc_met["incr_rate"],   bench_incr_pct,              20, False)
                + "</div>",
                unsafe_allow_html=True,
            )

        st.divider()

        # ── Section 8: AI Campaign Analysis ───────────────────────────────────
        st.subheader("AI Campaign Analysis")
        api_key_oc = get_api_key()

        if not api_key_oc:
            st.warning("Add `ANTHROPIC_API_KEY` to Streamlit secrets to enable AI analysis.")
        else:
            if st.button("Generate AI Analysis", type="primary", key="oc_ai_btn"):
                bench_roas  = oc_b["avg_roas"]
                bench_cpo   = oc_b["avg_cpo"]
                bench_ctr   = oc_b["avg_ctr"]
                bench_incr  = oc_b["incr_rate"] * 100

                oc_prompt = (
                    f"You are a senior Uber Advertising ANZ Partner Manager reviewing a completed "
                    f"campaign for a {oc_inp['category']} restaurant.\n\n"
                    f"CAMPAIGN DETAILS\n"
                    f"- Restaurant: {oc_inp.get('restaurant_name') or oc_inp['category']}\n"
                    f"- Spend: A${oc_inp['spend']:,.0f}\n"
                    f"- Duration: {oc_inp['start_date']} to {oc_inp['end_date']}\n\n"
                    f"ACTUALS vs TARGETS\n"
                    f"- ROAS: {oc_met['actual_roas']:.2f}x (target {oc_inp['target_roas']:.1f}x)\n"
                    f"- Total Orders: {oc_inp['total_orders']:,} (target {oc_inp['target_orders']:,})\n"
                    f"- Incremental Orders: {oc_inp['incr_orders']:,}\n"
                    f"- CPO: A${oc_met['actual_cpo']:.2f}\n"
                    f"- CTR: {oc_met['ctr']:.2f}%  |  Conv Rate: {oc_met['conv_rate']:.2f}%\n"
                    f"- Incremental Rate: {oc_met['incr_rate']:.1f}%\n"
                    f"- Campaign Health Score: {oc_h}/100\n\n"
                    f"CATEGORY BENCHMARKS\n"
                    f"- Benchmark ROAS: {bench_roas:.1f}x  |  Benchmark CPO: A${bench_cpo:.2f}\n"
                    f"- Benchmark CTR: {bench_ctr:.1f}%  |  Benchmark Incremental Rate: {bench_incr:.0f}%\n\n"
                    f"Write a 3-paragraph campaign analysis that:\n"
                    f"1. Summarises overall campaign performance and health score in context\n"
                    f"2. Identifies the 1-2 strongest and weakest metrics vs benchmark\n"
                    f"3. Gives 3 specific, actionable recommendations for the next campaign\n\n"
                    f"Write in a direct, commercially confident tone for a restaurant marketing director. "
                    f"Use A$ for all currency. Be specific with numbers."
                )

                ai_client_oc = anthropic.Anthropic(api_key=api_key_oc)
                with st.spinner("Generating campaign analysis…"):
                    oc_msg = ai_client_oc.messages.create(
                        model="claude-sonnet-4-6",
                        max_tokens=700,
                        system=(
                            "You are a senior Uber Advertising ANZ Partner Manager. "
                            "Write direct, data-led campaign reviews with clear next steps. "
                            "Use Australian English."
                        ),
                        messages=[{"role": "user", "content": oc_prompt}],
                    )
                st.session_state["oc_ai_text"] = oc_msg.content[0].text.strip()

            if st.session_state.get("oc_ai_text"):
                st.markdown(
                    "<div style='background:#FFFFFF;border-radius:12px;padding:22px 26px;"
                    "box-shadow:0 4px 12px rgba(0,0,0,0.08);margin-top:14px;font-size:14px;"
                    "line-height:1.8;color:#374151;'>"
                    + st.session_state["oc_ai_text"].replace("\n", "<br>") +
                    "</div>",
                    unsafe_allow_html=True,
                )

        st.divider()

        # ── Section 9: Export ─────────────────────────────────────────────────
        st.subheader("Export")
        if st.button("Download Outcome Report (.pptx)", key="oc_pptx_btn"):
            oc_ai    = st.session_state.get("oc_ai_text", "")
            oc_pptx  = build_outcome_pptx(oc_inp, oc_met, oc_h, oc_ai)
            st.session_state["oc_pptx"] = oc_pptx

        if st.session_state.get("oc_pptx"):
            filename_oc = f"{_date.today().isoformat()}_uber_ads_campaign_outcome.pptx"
            st.download_button(
                label="Save .pptx",
                data=st.session_state["oc_pptx"],
                file_name=filename_oc,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="oc_pptx_download",
            )


print("Uber Ads ROI Calculator page loaded.")
