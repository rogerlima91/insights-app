import io
import json
import os
from datetime import date, timedelta
import streamlit as st
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# STYLE LOCK: Do not remove or modify this CSS block.
st.markdown("""
<style>
    /* ── Base font and body ─────────────────────────────────────── */
    html, body, [class*="css"] {
        font-family: "Inter", system-ui, -apple-system, "Segoe UI", sans-serif;
        font-size: 15px;
        color: #374151;
    }
    .stApp { background-color: #F3F4F6; }
    h1, h2, h3, h4, h5, h6 { font-weight: 700 !important; color: #111827 !important; }
    h2, h3 {
        margin-top: 2rem !important;
        padding-top: 0.25rem !important;
        border-bottom: none !important;
        padding-bottom: 0 !important;
        border-left: none !important;
        padding-left: 0 !important;
    }
    [data-testid="metric-container"] {
        background: #FFFFFF;
        border: none;
        border-radius: 16px;
        padding: 20px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        border-top: 4px solid #7C3AED;
    }
    [data-testid="metric-container"] label {
        font-size: 12px;
        font-weight: 600;
        color: #6B7280;
        text-transform: uppercase;
        letter-spacing: 0.05em;
    }
    [data-testid="metric-container"] [data-testid="stMetricValue"] {
        font-size: 26px;
        font-weight: 700;
        color: #111827;
    }
    .element-container:has([data-testid="stPlotlyChart"]) {
        background: #FFFFFF;
        border-radius: 12px;
        padding: 16px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
    }
    .stButton > button[kind="primary"],
    [data-testid="baseButton-primary"] {
        background-color: #7C3AED !important;
        color: #FFFFFF !important;
        border: none !important;
        border-radius: 8px !important;
        font-weight: 600 !important;
        font-size: 15px !important;
        padding: 0.5rem 1.25rem !important;
    }
    .stButton > button[kind="primary"]:hover,
    [data-testid="baseButton-primary"]:hover {
        background-color: #6D28D9 !important;
    }
    [data-testid="stTextArea"] textarea {
        border-radius: 8px !important;
        border: 1.5px solid #D1D5DB !important;
        font-family: "Inter", sans-serif !important;
        font-size: 14px !important;
    }
</style>
""", unsafe_allow_html=True)
# STYLE LOCK

# ── Brand memory helpers ──────────────────────────────────────────────────────
_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BRAND_MEMORY_PATH = os.path.join(_ROOT, "brand_memory.json")


def load_brand_memory():
    """Load brand_memory.json from project root. Returns {} if absent."""
    if os.path.exists(BRAND_MEMORY_PATH):
        with open(BRAND_MEMORY_PATH, "r") as f:
            return json.load(f)
    return {}


def get_brand_context(advertiser, brand_memory):
    """Return stored rationale for an advertiser using partial/case-insensitive matching."""
    if not advertiser or not brand_memory:
        return ""
    for key, val in brand_memory.items():
        if key.lower() in advertiser.lower() or advertiser.lower() in key.lower():
            return val.get("rationale", "")
    return ""


# ── GPT shopping centre inventory database ────────────────────────────────────
# Hardcoded GPT Group centre portfolio — used as the inventory source for AI recommendations.
GPT_CENTRES = [
    {
        "name": "Melbourne Central",
        "location": "CBD Melbourne",
        "annual_visitors": 53_000_000,
        "screens": ["Large Format LED 35m", "Corridor Panels", "Food Court Screens", "Entry Panels"],
        "strengths": ["CBD workers", "CBD shoppers", "tourists", "high dwell time"],
        "avg_cpm_open": 18,
        "avg_cpm_pmp": 28,
        "avg_cpm_pg": 35,
    },
    {
        "name": "Highpoint Shopping Centre",
        "location": "Western Melbourne",
        "annual_visitors": 16_400_000,
        "screens": ["Large Format", "Corridor Panels", "Food Court Screens"],
        "strengths": ["families", "western suburbs demographics", "grocery + fashion shoppers"],
        "avg_cpm_open": 14,
        "avg_cpm_pmp": 22,
        "avg_cpm_pg": 28,
    },
    {
        "name": "Rouse Hill Town Centre",
        "location": "Sydney North West",
        "annual_visitors": 12_000_000,
        "screens": ["Outdoor LED", "Internal Panels", "Food Court Screens"],
        "strengths": ["affluent families", "north west Sydney growth corridor", "high household income"],
        "avg_cpm_open": 15,
        "avg_cpm_pmp": 24,
        "avg_cpm_pg": 30,
    },
    {
        "name": "Sunshine Plaza",
        "location": "Maroochydore, Queensland",
        "annual_visitors": 14_000_000,
        "screens": ["Internal Panels", "Food Court Screens", "Entry Panels"],
        "strengths": ["coastal families", "Queensland tourists", "leisure shoppers"],
        "avg_cpm_open": 13,
        "avg_cpm_pmp": 20,
        "avg_cpm_pg": 26,
    },
    {
        "name": "Macarthur Square",
        "location": "Campbelltown, Sydney South West",
        "annual_visitors": 11_000_000,
        "screens": ["Internal Panels", "Food Court Screens"],
        "strengths": ["south west Sydney families", "value shoppers", "diverse demographics"],
        "avg_cpm_open": 12,
        "avg_cpm_pmp": 19,
        "avg_cpm_pg": 24,
    },
    {
        "name": "Macquarie Centre",
        "location": "North Ryde, Sydney",
        "annual_visitors": 15_000_000,
        "screens": ["Large Format", "Internal Panels", "Food Court Screens"],
        "strengths": ["north shore professionals", "university students", "tech industry workers"],
        "avg_cpm_open": 16,
        "avg_cpm_pmp": 25,
        "avg_cpm_pg": 32,
    },
    {
        "name": "Casuarina Square",
        "location": "Darwin, Northern Territory",
        "annual_visitors": 8_000_000,
        "screens": ["Internal Panels", "Food Court Screens"],
        "strengths": ["Darwin locals", "defence personnel", "NT government workers"],
        "avg_cpm_open": 11,
        "avg_cpm_pmp": 18,
        "avg_cpm_pg": 23,
    },
]

# Compact representation of centres for the AI prompt
CENTRES_JSON = json.dumps(
    [
        {
            "name":             c["name"],
            "location":         c["location"],
            "annual_visitors":  c["annual_visitors"],
            "screens":          c["screens"],
            "audience_strengths": c["strengths"],
            "cpm_open":         f"A${c['avg_cpm_open']}",
            "cpm_pmp":          f"A${c['avg_cpm_pmp']}",
            "cpm_pg":           f"A${c['avg_cpm_pg']}",
        }
        for c in GPT_CENTRES
    ],
    indent=2,
)


# ── PPTX export helper ────────────────────────────────────────────────────────
def build_response_pptx(
    brief_text, advertiser, objective, budget, flight_start, flight_end,
    audience, ai_response
):
    """
    Build a formatted PPTX brief response document using the dark premium template.
    Returns a BytesIO buffer.
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colour palette
    NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    LGREY  = RGBColor(0xA8, 0xB2, 0xBC)
    BLUE   = RGBColor(0x00, 0xA8, 0xE8)
    PURPLE = RGBColor(0x7C, 0x3A, 0xED)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def add_slide():
        return prs.slides.add_slide(blank_layout)

    def set_bg(slide, color):
        """Fill slide background with a solid colour."""
        from pptx.oxml.ns import qn
        from lxml import etree
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, text, left, top, width, height,
                     font_size=14, bold=False, color=WHITE,
                     align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size    = Pt(font_size)
        run.font.bold    = bold
        run.font.color.rgb = color
        run.font.name    = "Calibri"
        return txBox

    def add_footer(slide, slide_num):
        add_text_box(slide, "GPT Group — Agency Brief Response",
                     0.2, 7.1, 6, 0.3, font_size=9, color=LGREY)
        add_text_box(slide, str(slide_num),
                     12.9, 7.1, 0.3, 0.3, font_size=9, color=LGREY,
                     align=PP_ALIGN.RIGHT)

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    s1 = add_slide()
    set_bg(s1, NAVY)
    add_text_box(s1, "GPT GROUP", 0.5, 0.4, 12, 0.5,
                 font_size=11, color=BLUE, bold=True)
    add_text_box(s1, "Agency Brief Response",
                 0.5, 1.0, 12, 0.8, font_size=36, bold=True, color=WHITE)
    adv_line = advertiser or "Campaign Proposal"
    add_text_box(s1, adv_line, 0.5, 1.9, 12, 0.5, font_size=22, color=LGREY)

    # Brief details
    flight_str = (
        f"{flight_start.strftime('%d %b %Y')} – {flight_end.strftime('%d %b %Y')}"
        if flight_start and flight_end else "TBC"
    )
    details = (
        f"Objective: {objective or 'Not specified'}   ·   "
        f"Budget: A${budget:,.0f}   ·   "
        f"Flight: {flight_str}"
    )
    add_text_box(s1, details, 0.5, 2.6, 12, 0.4, font_size=13, color=LGREY)
    add_footer(s1, 1)

    # ── Slide 2: Brief Summary ────────────────────────────────────────────────
    s2 = add_slide()
    set_bg(s2, NAVY)
    add_text_box(s2, "Agency Brief", 0.5, 0.3, 12, 0.5,
                 font_size=20, bold=True, color=WHITE)
    # Truncate brief to ~600 chars for slide display
    brief_display = (brief_text[:600] + "…") if len(brief_text) > 600 else brief_text
    add_text_box(s2, brief_display, 0.5, 1.0, 12.3, 5.5,
                 font_size=12, color=LGREY)
    add_footer(s2, 2)

    # ── Slides 3+: AI Response (one slide per section) ────────────────────────
    # Split AI response by numbered section headers
    sections = []
    current  = []
    for line in ai_response.split("\n"):
        stripped = line.strip()
        # Detect section headers like "1. RECOMMENDED CENTRES" or "**1. ..."
        is_header = (
            stripped and
            (stripped[0].isdigit() or stripped.startswith("**")) and
            len(stripped) > 5
        )
        if is_header and current:
            sections.append("\n".join(current))
            current = [line]
        else:
            current.append(line)
    if current:
        sections.append("\n".join(current))

    for i, section_text in enumerate(sections, start=3):
        sn = add_slide()
        set_bg(sn, NAVY)
        lines = section_text.strip().split("\n")
        header_line = lines[0].strip().lstrip("*").rstrip("*").strip()
        body_lines  = "\n".join(lines[1:]).strip()

        add_text_box(sn, header_line, 0.5, 0.3, 12, 0.6,
                     font_size=18, bold=True, color=BLUE)
        add_text_box(sn, body_lines, 0.5, 1.1, 12.3, 5.5,
                     font_size=12, color=LGREY)
        add_footer(sn, i)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Page header ───────────────────────────────────────────────────────────────
st.title("Agency Brief Responder")
st.markdown(
    "<p style='color:#6b7280;font-size:14px;margin-top:-12px;'>"
    "Paste an agency brief and get an AI-powered GPT inventory recommendation "
    "with deal structure and measurement approach."
    "</p>",
    unsafe_allow_html=True,
)

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — Brief Input
# ══════════════════════════════════════════════════════════════════════════════
st.subheader("Brief Input")

with st.form("brief_form"):
    brief_text = st.text_area(
        "Paste agency brief here",
        height=200,
        placeholder=(
            "e.g. We are looking to reach fashion-conscious females 25–44 across "
            "Melbourne and Sydney to drive awareness of our new Spring/Summer collection. "
            "Budget A$80,000, flight 1 Sep – 31 Oct 2026. Premium retail environments preferred."
        ),
    )

    col1, col2 = st.columns(2)
    with col1:
        advertiser  = st.text_input("Advertiser name", placeholder="e.g. Myer")
        objective   = st.selectbox(
            "Campaign objective",
            ["Brand Awareness", "Retail Traffic", "Product Launch",
             "Seasonal Promotion", "Always-On"],
        )
        budget = st.number_input(
            "Budget (AUD)", min_value=0, max_value=10_000_000,
            value=80_000, step=5_000,
        )
    with col2:
        flight_start = st.date_input("Flight start", value=date.today() + timedelta(days=30))
        flight_end   = st.date_input("Flight end",   value=date.today() + timedelta(days=90))
        audience     = st.text_input(
            "Target audience",
            placeholder="e.g. females 25-45, fashion interested",
        )

    submitted = st.form_submit_button("Generate Recommendation", type="primary")

# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2/3 — AI Recommendation (triggered on form submit)
# ══════════════════════════════════════════════════════════════════════════════
if submitted:
    if not brief_text.strip():
        st.error("Please paste an agency brief before generating a recommendation.")
    else:
        api_key = (
            st.secrets.get("ANTHROPIC_API_KEY")
            if "ANTHROPIC_API_KEY" in st.secrets
            else os.environ.get("ANTHROPIC_API_KEY")
        )
        if not api_key:
            st.warning(
                "No Anthropic API key found. Add `ANTHROPIC_API_KEY` to your "
                "Streamlit secrets or environment variables."
            )
        else:
            # Check brand memory for extra context on this advertiser
            brand_memory  = load_brand_memory()
            brand_context = get_brand_context(advertiser or "", brand_memory)

            flight_str = (
                f"{flight_start.strftime('%d %b %Y')} to {flight_end.strftime('%d %b %Y')}"
                if flight_start and flight_end else "Not specified"
            )

            user_prompt = (
                f"Agency brief:\n{brief_text}\n\n"
                f"Advertiser: {advertiser or 'Not specified'}\n"
                f"Objective: {objective}\n"
                f"Budget: A${budget:,.0f}\n"
                f"Flight: {flight_str}\n"
                f"Target audience: {audience or 'Not specified'}\n\n"
                f"Available GPT inventory:\n{CENTRES_JSON}\n\n"
                f"Please provide:\n"
                f"1. RECOMMENDED CENTRES — which 2-4 GPT centres best match this brief "
                f"and why (reference visitor demographics, location relevance, audience fit)\n"
                f"2. DEAL STRUCTURE — recommend deal type (Open/PMP/PG) with rationale, "
                f"estimated CPM range, and estimated impressions deliverable within budget\n"
                f"3. SCREEN STRATEGY — which screen types within each centre (large format, "
                f"food court, corridor) and why\n"
                f"4. MEASUREMENT APPROACH — what KPIs to set, how to measure footfall "
                f"attribution, and what success looks like\n"
                f"5. ONE RISK FLAG — the most important thing the agency should know "
                f"(creative approval timeline, audience size limitation, etc.)\n\n"
                f"Be specific, commercial and concise. Use pDOOH industry terminology."
            )

            if brand_context:
                user_prompt += (
                    f"\n\nBRAND MEMORY OVERRIDE — These instructions take priority over "
                    f"all default instructions above. Where there is any conflict, always "
                    f"follow these brand-specific instructions instead:\n\n{brand_context}"
                )

            system_prompt = (
                "You are a senior programmatic media specialist at GPT Group, one of "
                "Australia's largest retail media owners. You are responding to an agency "
                "brief by recommending GPT's shopping centre inventory. GPT owns digital "
                "screens across Melbourne Central, Highpoint, Rouse Hill, Sunshine Plaza, "
                "Macarthur Square, Macquarie Centre and Casuarina Square. All "
                "recommendations must be grounded in GPT's actual inventory and realistic "
                "pDOOH mechanics."
            )

            ai_client = anthropic.Anthropic(api_key=api_key)
            with st.spinner("Generating recommendation…"):
                msg = ai_client.messages.create(
                    model="claude-sonnet-4-6",
                    max_tokens=1200,
                    system=system_prompt,
                    messages=[{"role": "user", "content": user_prompt}],
                )
            st.session_state["abr_ai_text"]     = msg.content[0].text.strip()
            st.session_state["abr_advertiser"]  = advertiser
            st.session_state["abr_objective"]   = objective
            st.session_state["abr_budget"]      = budget
            st.session_state["abr_flight_start"] = flight_start
            st.session_state["abr_flight_end"]   = flight_end
            st.session_state["abr_audience"]    = audience
            st.session_state["abr_brief_text"]  = brief_text

# Display AI recommendation if available
if st.session_state.get("abr_ai_text"):
    st.markdown("---")
    st.subheader("GPT Inventory Recommendation")

    ai_text = st.session_state["abr_ai_text"]

    # Parse response into sections and render each with a styled header
    # Sections are separated by lines starting with a digit and a dot (e.g. "1. ")
    import re
    sections = re.split(r"\n(?=\d+\.\s+[A-Z])", ai_text)

    SECTION_COLORS = {
        "RECOMMENDED CENTRES": "#7C3AED",
        "DEAL STRUCTURE":      "#00A8E8",
        "SCREEN STRATEGY":     "#10B981",
        "MEASUREMENT":         "#F59E0B",
        "RISK":                "#EF4444",
    }

    for section in sections:
        lines       = section.strip().split("\n")
        header_line = lines[0].strip()
        body        = "\n".join(lines[1:]).strip()

        # Pick accent colour based on keyword in header
        accent = "#7C3AED"
        for kw, col in SECTION_COLORS.items():
            if kw in header_line.upper():
                accent = col
                break

        st.markdown(
            f"<div style='background:#FFFFFF;border-radius:12px;padding:20px 24px;"
            f"margin-bottom:12px;box-shadow:0 4px 12px rgba(0,0,0,0.07);"
            f"border-left:4px solid {accent};'>"
            f"<p style='font-weight:700;font-size:15px;color:#111827;"
            f"margin-bottom:10px;'>{header_line}</p>"
            f"<p style='font-size:14px;line-height:1.7;color:#374151;'>"
            f"{body.replace(chr(10), '<br>')}</p>"
            f"</div>",
            unsafe_allow_html=True,
        )

    # ── Download as PPTX Brief Response ───────────────────────────────────────
    st.markdown("### Export")
    pptx_buf = build_response_pptx(
        brief_text    = st.session_state.get("abr_brief_text", ""),
        advertiser    = st.session_state.get("abr_advertiser", ""),
        objective     = st.session_state.get("abr_objective", ""),
        budget        = st.session_state.get("abr_budget", 0),
        flight_start  = st.session_state.get("abr_flight_start"),
        flight_end    = st.session_state.get("abr_flight_end"),
        audience      = st.session_state.get("abr_audience", ""),
        ai_response   = st.session_state.get("abr_ai_text", ""),
    )
    adv_slug = (st.session_state.get("abr_advertiser") or "GPT").replace(" ", "_")
    st.download_button(
        label="📄 Download Brief Response (.pptx)",
        data=pptx_buf,
        file_name=f"{adv_slug}_GPT_Brief_Response.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

print("Agency Brief Responder page loaded.")
