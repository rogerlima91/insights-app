import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from datetime import date, timedelta
import io
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
from utils.design_system import get_css, PRIMARY, SECONDARY, SUCCESS, WARNING, DANGER, WHITE, TEXT_SEC, TEXT_PRI, CHART_PALETTE

# ── Apply Pacebird design system CSS ──────────────────────────────────────────
# STYLE LOCK: Pacebird design system — primary #F5A623 orange, secondary #1B2A4A navy, font Poppins.
st.markdown(get_css(), unsafe_allow_html=True)

# ── Mock data generator ───────────────────────────────────────────────────────
# Generates 180 days of daily publisher data covering two full quarters.
# All currency values are in AUD (Australian Dollars).
@st.cache_data
def generate_qbr_data():
    rng = np.random.default_rng(seed=42)

    publishers = ["Nine.com.au", "news.com.au", "Domain", "realestate.com.au", "Stuff NZ"]
    formats    = ["inRead Video", "inFeed", "Display", "CTV"]

    # Base eCPM per publisher × format (AUD) — sell-side rates
    ecpm_table = {
        "Nine.com.au":         {"inRead Video": 18.5, "inFeed": 9.2,  "Display": 4.8, "CTV": 28.0},
        "news.com.au":         {"inRead Video": 16.8, "inFeed": 8.5,  "Display": 4.3, "CTV": 25.5},
        "Domain":              {"inRead Video": 24.0, "inFeed": 14.5, "Display": 6.5, "CTV": 0.0},
        "realestate.com.au":   {"inRead Video": 22.5, "inFeed": 13.0, "Display": 6.0, "CTV": 0.0},
        "Stuff NZ":            {"inRead Video": 12.0, "inFeed": 6.5,  "Display": 3.5, "CTV": 0.0},
    }

    # Base fill rate per publisher (%) — higher-CPM publishers typically have lower fill
    fill_table = {
        "Nine.com.au":       78.0,
        "news.com.au":       75.0,
        "Domain":            65.0,
        "realestate.com.au": 62.0,
        "Stuff NZ":          80.0,
    }

    # Base daily impressions per publisher × format (thousands)
    imp_table = {
        "Nine.com.au":         {"inRead Video": 420, "inFeed": 850,  "Display": 2200, "CTV": 95},
        "news.com.au":         {"inRead Video": 380, "inFeed": 780,  "Display": 2000, "CTV": 85},
        "Domain":              {"inRead Video": 120, "inFeed": 220,  "Display": 600,  "CTV": 0},
        "realestate.com.au":   {"inRead Video": 110, "inFeed": 200,  "Display": 550,  "CTV": 0},
        "Stuff NZ":            {"inRead Video": 90,  "inFeed": 180,  "Display": 480,  "CTV": 0},
    }

    # 180 days starting 2025-01-01 covers Q1 and Q2 2025
    start_date = date(2025, 1, 1)
    rows = []
    for day_offset in range(180):
        d = start_date + timedelta(days=day_offset)
        # Slight weekend dip (Sat=5, Sun=6)
        weekday_mult = 0.70 if d.weekday() >= 5 else 1.0

        for pub in publishers:
            for fmt in formats:
                base_imp = imp_table[pub].get(fmt, 0)
                base_ecpm = ecpm_table[pub].get(fmt, 0.0)
                if base_imp == 0 or base_ecpm == 0.0:
                    continue

                # Add daily noise (±15%) and weekday effect
                noise  = rng.normal(1.0, 0.08)
                imps   = int(base_imp * 1000 * weekday_mult * noise)
                imps   = max(imps, 0)

                ecpm_noise = rng.normal(1.0, 0.06)
                ecpm   = round(base_ecpm * ecpm_noise, 2)
                ecpm   = max(ecpm, 1.0)

                revenue = round(imps / 1000 * ecpm, 2)

                fill_noise = rng.normal(1.0, 0.05)
                fill_rate  = round(min(fill_table[pub] * fill_noise, 99.9), 1)

                # Placement = publisher + format combination (simulates ad slot name)
                placement = f"{pub} — {fmt}"

                rows.append({
                    "date":        d,
                    "publisher":   pub,
                    "placement":   placement,
                    "format":      fmt,
                    "impressions": imps,
                    "revenue_aud": revenue,
                    "ecpm":        ecpm,
                    "fill_rate":   fill_rate,
                })

    return pd.DataFrame(rows)


def quarter_label(d):
    """Return e.g. 'Q1 2025' for a given date."""
    q = (d.month - 1) // 3 + 1
    return f"Q{q} {d.year}"


def get_quarter_dates(label):
    """Return (start_date, end_date) for a quarter label like 'Q1 2025'."""
    parts = label.split()
    q, yr = int(parts[0][1]), int(parts[1])
    month_start = (q - 1) * 3 + 1
    start = date(yr, month_start, 1)
    end_month = month_start + 2
    # Last day of quarter
    if end_month == 12:
        end = date(yr, 12, 31)
    else:
        end = date(yr, end_month + 1, 1) - timedelta(days=1)
    return start, end


# ── Page header ───────────────────────────────────────────────────────────────
st.markdown("# 📋 QBR Generator")
st.markdown("Quarterly Business Review reports for ANZ publishers. Select a publisher and quarter to generate a performance summary.")

# ── Load data ─────────────────────────────────────────────────────────────────
df_all = generate_qbr_data()

# Identify available quarters in the data
df_all["quarter"] = df_all["date"].apply(quarter_label)
available_quarters = sorted(df_all["quarter"].unique(), reverse=True)  # newest first

# ── Filters ───────────────────────────────────────────────────────────────────
col_f1, col_f2, col_f3 = st.columns([2, 2, 4])
with col_f1:
    publishers   = ["All Publishers"] + sorted(df_all["publisher"].unique().tolist())
    selected_pub = st.selectbox("Publisher", publishers, key="qbr_pub")
with col_f2:
    # Current quarter is the most recent; default to Q2 for the "current" selection
    quarter_options = available_quarters
    selected_q = st.selectbox("Quarter", quarter_options, index=0, key="qbr_quarter")

# Determine the previous quarter for comparison
def prev_quarter(label):
    parts = label.split()
    q, yr = int(parts[0][1]), int(parts[1])
    if q == 1:
        return f"Q4 {yr - 1}"
    return f"Q{q - 1} {yr}"

prev_q = prev_quarter(selected_q)

# Filter data
def filter_data(df, pub, quarter):
    out = df[df["quarter"] == quarter]
    if pub != "All Publishers":
        out = out[out["publisher"] == pub]
    return out

df_curr = filter_data(df_all, selected_pub, selected_q)
df_prev = filter_data(df_all, selected_pub, prev_q)

if df_curr.empty:
    st.warning("No data available for the selected filters.")
    st.stop()

# ── Summary metrics ───────────────────────────────────────────────────────────
def summarise(df):
    return {
        "revenue":     df["revenue_aud"].sum(),
        "impressions": df["impressions"].sum(),
        "ecpm":        (df["revenue_aud"].sum() / df["impressions"].sum() * 1000) if df["impressions"].sum() > 0 else 0,
        "fill_rate":   df["fill_rate"].mean(),
    }

curr = summarise(df_curr)
prev = summarise(df_prev) if not df_prev.empty else None

def pct_chg(c, p):
    if p and p != 0:
        return (c - p) / abs(p) * 100
    return None

def delta_str(val, suffix=""):
    if val is None:
        return None
    sign = "+" if val >= 0 else ""
    return f"{sign}{val:.1f}%{suffix}"

st.markdown(f"### {selected_q} — Performance Summary")
m1, m2, m3, m4 = st.columns(4)

rev_delta   = pct_chg(curr["revenue"],     prev["revenue"])     if prev else None
imp_delta   = pct_chg(curr["impressions"], prev["impressions"]) if prev else None
ecpm_delta  = pct_chg(curr["ecpm"],        prev["ecpm"])        if prev else None
fill_delta  = pct_chg(curr["fill_rate"],   prev["fill_rate"])   if prev else None

with m1:
    st.metric("Revenue (AUD)", f"A${curr['revenue']:,.0f}", delta=delta_str(rev_delta))
with m2:
    st.metric("Impressions", f"{curr['impressions']:,.0f}", delta=delta_str(imp_delta))
with m3:
    st.metric("eCPM", f"A${curr['ecpm']:.2f}", delta=delta_str(ecpm_delta))
with m4:
    st.metric("Avg Fill Rate", f"{curr['fill_rate']:.1f}%", delta=delta_str(fill_delta))

# ── Quarter-over-quarter comparison table ─────────────────────────────────────
st.markdown("### Quarter-over-Quarter Comparison")

if prev is not None:
    # Build comparison by publisher × format for depth
    def summarise_by(df, grp):
        g = df.groupby(grp).agg(
            revenue_aud=("revenue_aud", "sum"),
            impressions=("impressions", "sum"),
            fill_rate=("fill_rate", "mean"),
        ).reset_index()
        g["ecpm"] = (g["revenue_aud"] / g["impressions"] * 1000).round(2)
        return g

    grp_col = "format" if selected_pub != "All Publishers" else "publisher"
    tbl_curr = summarise_by(df_curr, grp_col).rename(columns={
        "revenue_aud": f"Revenue {selected_q}",
        "impressions": f"Impressions {selected_q}",
        "ecpm":        f"eCPM {selected_q}",
        "fill_rate":   f"Fill Rate {selected_q}",
    })
    tbl_prev = summarise_by(df_prev, grp_col).rename(columns={
        "revenue_aud": f"Revenue {prev_q}",
        "impressions": f"Impressions {prev_q}",
        "ecpm":        f"eCPM {prev_q}",
        "fill_rate":   f"Fill Rate {prev_q}",
    })

    tbl = tbl_curr.merge(tbl_prev, on=grp_col, how="outer").fillna(0)

    # Variance columns with directional indicators
    def var_col(c, p, higher_better=True):
        if p == 0:
            return "N/A"
        pct = (c - p) / abs(p) * 100
        arrow = "▲" if pct >= 0 else "▼"
        color_ok = (pct >= 0) == higher_better
        # We'll return a formatted string; coloring is done via st.dataframe styler
        return f"{arrow} {abs(pct):.1f}%"

    tbl["Rev Var %"]  = tbl.apply(lambda r: var_col(r[f"Revenue {selected_q}"], r[f"Revenue {prev_q}"]), axis=1)
    tbl["eCPM Var %"] = tbl.apply(lambda r: var_col(r[f"eCPM {selected_q}"], r[f"eCPM {prev_q}"]), axis=1)
    tbl["Fill Var %"] = tbl.apply(lambda r: var_col(r[f"Fill Rate {selected_q}"], r[f"Fill Rate {prev_q}"]), axis=1)

    # Format currency and percentage columns for display
    display_tbl = tbl[[grp_col,
        f"Revenue {prev_q}", f"Revenue {selected_q}", "Rev Var %",
        f"eCPM {prev_q}",    f"eCPM {selected_q}",    "eCPM Var %",
        f"Fill Rate {prev_q}", f"Fill Rate {selected_q}", "Fill Var %",
    ]].copy()

    display_tbl[f"Revenue {prev_q}"]    = display_tbl[f"Revenue {prev_q}"].apply(lambda x: f"A${x:,.0f}")
    display_tbl[f"Revenue {selected_q}"] = display_tbl[f"Revenue {selected_q}"].apply(lambda x: f"A${x:,.0f}")
    display_tbl[f"eCPM {prev_q}"]       = display_tbl[f"eCPM {prev_q}"].apply(lambda x: f"A${x:.2f}")
    display_tbl[f"eCPM {selected_q}"]   = display_tbl[f"eCPM {selected_q}"].apply(lambda x: f"A${x:.2f}")
    display_tbl[f"Fill Rate {prev_q}"]  = display_tbl[f"Fill Rate {prev_q}"].apply(lambda x: f"{x:.1f}%")
    display_tbl[f"Fill Rate {selected_q}"] = display_tbl[f"Fill Rate {selected_q}"].apply(lambda x: f"{x:.1f}%")

    st.dataframe(display_tbl.rename(columns={grp_col: grp_col.title()}), use_container_width=True, hide_index=True)
else:
    st.info(f"No data for {prev_q} — cannot show quarter-over-quarter comparison.")

# ── Revenue trend line chart ───────────────────────────────────────────────────
st.markdown("### Revenue Trend")

# Aggregate to daily totals
daily_rev = df_curr.groupby("date")["revenue_aud"].sum().reset_index()
daily_rev.columns = ["Date", "Revenue (AUD)"]

fig_trend = go.Figure()
fig_trend.add_trace(go.Scatter(
    x=daily_rev["Date"],
    y=daily_rev["Revenue (AUD)"],
    mode="lines+markers",
    line=dict(color=PRIMARY, width=2),
    marker=dict(size=4, color=PRIMARY),
    fill="tozeroy",
    fillcolor=f"rgba(245,166,35,0.12)",  # orange tint
    name="Daily Revenue",
))
fig_trend.update_layout(
    title=f"Daily Revenue — {selected_q}",
    xaxis_title="Date",
    yaxis_title="Revenue (A$)",
    yaxis_tickprefix="A$",
    yaxis_tickformat=",.0f",
    plot_bgcolor="white",
    paper_bgcolor="white",
    font=dict(family="Poppins, sans-serif", size=12),
    margin=dict(t=40, b=40, l=60, r=20),
    height=320,
)
st.plotly_chart(fig_trend, use_container_width=True)

# ── Revenue by ad format bar chart ─────────────────────────────────────────────
st.markdown("### Revenue by Ad Format")

rev_fmt = df_curr.groupby("format")["revenue_aud"].sum().reset_index().sort_values("revenue_aud", ascending=False)
rev_fmt.columns = ["Format", "Revenue (AUD)"]

fig_fmt = px.bar(
    rev_fmt,
    x="Format",
    y="Revenue (AUD)",
    color="Format",
    color_discrete_sequence=CHART_PALETTE,
    text=rev_fmt["Revenue (AUD)"].apply(lambda x: f"A${x:,.0f}"),
)
fig_fmt.update_traces(textposition="outside", textfont_size=11)
fig_fmt.update_layout(
    showlegend=False,
    plot_bgcolor="white",
    paper_bgcolor="white",
    yaxis_tickprefix="A$",
    yaxis_tickformat=",.0f",
    font=dict(family="Poppins, sans-serif", size=12),
    margin=dict(t=20, b=40, l=60, r=20),
    height=320,
)
st.plotly_chart(fig_fmt, use_container_width=True)

# ── Top 5 and Bottom 5 placements by eCPM ─────────────────────────────────────
st.markdown("### Placement Performance")

# Aggregate revenue, impressions, fill rate by placement
place_agg = df_curr.groupby("placement").agg(
    revenue_aud=("revenue_aud", "sum"),
    impressions=("impressions", "sum"),
    fill_rate=("fill_rate", "mean"),
).reset_index()
place_agg["eCPM"] = (place_agg["revenue_aud"] / place_agg["impressions"] * 1000).round(2)
place_agg = place_agg[place_agg["impressions"] > 0]

top5    = place_agg.nlargest(5, "eCPM")
bottom5 = place_agg.nsmallest(5, "eCPM")

col_t, col_b = st.columns(2)

with col_t:
    st.markdown("**Top 5 Placements by eCPM**")
    top5_display = top5[["placement", "eCPM", "revenue_aud", "fill_rate"]].copy()
    top5_display["eCPM"]        = top5_display["eCPM"].apply(lambda x: f"A${x:.2f}")
    top5_display["revenue_aud"] = top5_display["revenue_aud"].apply(lambda x: f"A${x:,.0f}")
    top5_display["fill_rate"]   = top5_display["fill_rate"].apply(lambda x: f"{x:.1f}%")
    top5_display.columns        = ["Placement", "eCPM", "Revenue", "Fill Rate"]
    st.dataframe(top5_display, use_container_width=True, hide_index=True)

with col_b:
    st.markdown("**Bottom 5 Placements by eCPM**")
    bot5_display = bottom5[["placement", "eCPM", "revenue_aud", "fill_rate"]].copy()
    bot5_display["eCPM"]        = bot5_display["eCPM"].apply(lambda x: f"A${x:.2f}")
    bot5_display["revenue_aud"] = bot5_display["revenue_aud"].apply(lambda x: f"A${x:,.0f}")
    bot5_display["fill_rate"]   = bot5_display["fill_rate"].apply(lambda x: f"{x:.1f}%")
    bot5_display.columns        = ["Placement", "eCPM", "Revenue", "Fill Rate"]
    st.dataframe(bot5_display, use_container_width=True, hide_index=True)

# ── Generate QBR Deck dialog ───────────────────────────────────────────────────
st.markdown("---")

@st.dialog("Build Your QBR Deck")
def show_qbr_dialog():
    """Section picker dialog — same pattern as Performance & Insights page."""
    st.caption("Select the sections to include in your PowerPoint QBR export.")

    defaults = st.session_state.get("_qbr_section_defaults", {
        "summary_metrics":    True,
        "qoq_comparison":     True,
        "revenue_trend":      True,
        "format_breakdown":   True,
        "top_bottom_placements": True,
    })

    sc = st.columns(2)
    with sc[0]:
        if st.button("Select all", key="qbr_sel_all"):
            st.session_state["_qbr_section_defaults"] = {k: True for k in defaults}
            st.rerun()
    with sc[1]:
        if st.button("Deselect all", key="qbr_desel_all"):
            st.session_state["_qbr_section_defaults"] = {k: False for k in defaults}
            st.rerun()

    st.markdown("")

    s_metrics  = st.checkbox("Summary Metrics",         value=defaults.get("summary_metrics", True),    key="qbr_cb_metrics")
    s_qoq      = st.checkbox("Quarter-over-Quarter",    value=defaults.get("qoq_comparison", True),     key="qbr_cb_qoq",
                              disabled=prev is None, help="Requires prior quarter data." if prev is None else None)
    s_trend    = st.checkbox("Revenue Trend Chart",     value=defaults.get("revenue_trend", True),      key="qbr_cb_trend")
    s_format   = st.checkbox("Revenue by Format Chart", value=defaults.get("format_breakdown", True),   key="qbr_cb_format")
    s_placements = st.checkbox("Top/Bottom Placements", value=defaults.get("top_bottom_placements", True), key="qbr_cb_placements")

    st.markdown("")
    bc = st.columns(2)
    with bc[0]:
        if st.button("Generate PowerPoint", type="primary", key="qbr_gen_btn"):
            st.session_state["_qbr_section_defaults"] = {
                "summary_metrics":       s_metrics,
                "qoq_comparison":        s_qoq,
                "revenue_trend":         s_trend,
                "format_breakdown":      s_format,
                "top_bottom_placements": s_placements,
            }
            sections = st.session_state["_qbr_section_defaults"]
            if not any(sections.values()):
                st.warning("Select at least one section.")
                return
            with st.spinner("Building QBR deck..."):
                try:
                    pptx_buf = _build_qbr_pptx(
                        sections, selected_q, prev_q, curr, prev,
                        daily_rev, rev_fmt, top5_display, bot5_display,
                    )
                    st.session_state["_qbr_pptx"] = pptx_buf
                    st.session_state["_qbr_filename"] = (
                        f"QBR_{selected_pub.replace(' ', '_')}_{selected_q.replace(' ', '_')}.pptx"
                    )
                except Exception as e:
                    st.error(f"Export failed: {e}")
                    return
    with bc[1]:
        if st.button("Cancel", key="qbr_cancel_btn"):
            st.rerun()

    # Download button shown immediately once PPTX is ready (no rerun — stays in dialog)
    if "_qbr_pptx" in st.session_state:
        st.success("Deck ready!")
        st.download_button(
            "📥 Download QBR (.pptx)",
            data=st.session_state["_qbr_pptx"],
            file_name=st.session_state.get("_qbr_filename", "QBR_report.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="qbr_dl_btn",
        )


def _build_qbr_pptx(sections, curr_q, prev_q_lbl, curr_metrics, prev_metrics, daily_rev_df, fmt_df, top5_df, bot5_df):
    """
    Build a QBR PowerPoint deck using python-pptx.
    Follows the same dark premium template used in Performance & Insights.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def hex_rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    # Dark premium template colors
    BG     = hex_rgb("#0D1B2A")
    WHITE_ = hex_rgb("#FFFFFF")
    GREY_  = hex_rgb("#A8B2BC")
    ACCENT = hex_rgb("#F5A623")  # orange accent matching brand

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    def add_slide():
        sl = prs.slides.add_slide(blank_layout)
        # Dark background
        bg = sl.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG
        return sl

    def add_text(sl, text, left, top, width, height, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
        txb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.color.rgb = color or WHITE_

    def add_footer(sl, slide_num):
        add_text(sl, "Pacebird — QBR Report", 0.3, 7.1, 4, 0.3, size=9, color=GREY_)
        add_text(sl, str(slide_num), 12.8, 7.1, 0.5, 0.3, size=9, color=GREY_, align=PP_ALIGN.RIGHT)

    slide_num = 1

    # ── Slide 1: Title ─────────────────────────────────────────────────────────
    sl = add_slide()
    add_text(sl, f"Quarterly Business Review", 1.0, 2.0, 11.0, 1.0, size=32, bold=True, color=WHITE_)
    add_text(sl, f"{selected_pub}  ·  {curr_q}", 1.0, 3.2, 11.0, 0.6, size=18, color=ACCENT)
    add_text(sl, "Prepared by Pacebird", 1.0, 4.2, 6.0, 0.4, size=12, color=GREY_)
    add_footer(sl, slide_num); slide_num += 1

    # ── Slide 2: Summary Metrics ───────────────────────────────────────────────
    if sections.get("summary_metrics"):
        sl = add_slide()
        add_text(sl, "Summary Metrics", 0.5, 0.3, 12.0, 0.6, size=22, bold=True)

        metrics = [
            ("Revenue (AUD)", f"A${curr_metrics['revenue']:,.0f}"),
            ("Impressions",   f"{curr_metrics['impressions']:,.0f}"),
            ("eCPM",          f"A${curr_metrics['ecpm']:.2f}"),
            ("Avg Fill Rate", f"{curr_metrics['fill_rate']:.1f}%"),
        ]
        x_start = 0.5
        for i, (label, value) in enumerate(metrics):
            x = x_start + i * 3.1
            # Card background
            card = sl.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(2.8), Inches(1.6))
            card.fill.solid()
            card.fill.fore_color.rgb = hex_rgb("#1B2A4A")
            card.line.color.rgb = ACCENT
            add_text(sl, label, x + 0.15, 1.35, 2.5, 0.35, size=10, color=GREY_)
            add_text(sl, value, x + 0.15, 1.75, 2.5, 0.6,  size=20, bold=True, color=WHITE_)

        if prev_metrics:
            add_text(sl, f"vs {prev_q_lbl}:", 0.5, 3.1, 3.0, 0.35, size=11, color=GREY_)
            changes = [
                pct_chg(curr_metrics['revenue'],     prev_metrics['revenue']),
                pct_chg(curr_metrics['impressions'], prev_metrics['impressions']),
                pct_chg(curr_metrics['ecpm'],        prev_metrics['ecpm']),
                pct_chg(curr_metrics['fill_rate'],   prev_metrics['fill_rate']),
            ]
            for i, chg in enumerate(changes):
                if chg is not None:
                    x = x_start + i * 3.1
                    clr = hex_rgb("#10B981") if chg >= 0 else hex_rgb("#EF4444")
                    add_text(sl, delta_str(chg), x + 0.15, 3.05, 2.5, 0.4, size=13, bold=True, color=clr)

        add_footer(sl, slide_num); slide_num += 1

    # ── Slide 3: QoQ Comparison ────────────────────────────────────────────────
    if sections.get("qoq_comparison") and prev_metrics:
        sl = add_slide()
        add_text(sl, f"Quarter-over-Quarter: {prev_q_lbl} → {curr_q}", 0.5, 0.3, 12.0, 0.6, size=20, bold=True)
        add_text(sl, "Revenue, eCPM and Fill Rate by format or publisher vs prior quarter.", 0.5, 0.9, 12.0, 0.4, size=11, color=GREY_)
        # Simple text summary
        lines = [
            f"• Revenue: A${curr_metrics['revenue']:,.0f}  ({delta_str(pct_chg(curr_metrics['revenue'], prev_metrics['revenue']))} vs {prev_q_lbl})",
            f"• eCPM: A${curr_metrics['ecpm']:.2f}  ({delta_str(pct_chg(curr_metrics['ecpm'], prev_metrics['ecpm']))} vs {prev_q_lbl})",
            f"• Fill Rate: {curr_metrics['fill_rate']:.1f}%  ({delta_str(pct_chg(curr_metrics['fill_rate'], prev_metrics['fill_rate']))} vs {prev_q_lbl})",
        ]
        for j, line in enumerate(lines):
            add_text(sl, line, 1.0, 1.6 + j * 0.7, 11.0, 0.5, size=14, color=WHITE_)
        add_footer(sl, slide_num); slide_num += 1

    # ── Slide 4: Top/Bottom Placements ────────────────────────────────────────
    if sections.get("top_bottom_placements"):
        sl = add_slide()
        add_text(sl, "Placement Performance", 0.5, 0.3, 12.0, 0.6, size=20, bold=True)
        add_text(sl, "Top 5 Placements by eCPM", 0.5, 1.0, 6.0, 0.4, size=13, color=ACCENT)
        for j, (_, row) in enumerate(top5_df.iterrows()):
            vals = row.values
            add_text(sl, f"• {vals[0]}  |  eCPM: {vals[1]}  |  Rev: {vals[2]}  |  Fill: {vals[3]}", 0.5, 1.5 + j * 0.55, 12.0, 0.45, size=11, color=WHITE_)
        add_text(sl, "Bottom 5 Placements by eCPM", 0.5, 4.4, 6.0, 0.4, size=13, color=GREY_)
        for j, (_, row) in enumerate(bot5_df.iterrows()):
            vals = row.values
            add_text(sl, f"• {vals[0]}  |  eCPM: {vals[1]}  |  Rev: {vals[2]}  |  Fill: {vals[3]}", 0.5, 4.9 + j * 0.45, 12.0, 0.4, size=11, color=GREY_)
        add_footer(sl, slide_num); slide_num += 1

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# Button to open dialog + persistent download once a deck has been generated
col_btn, col_dl, _ = st.columns([2, 3, 3])
with col_btn:
    if st.button("📊 Generate QBR Deck", type="primary", key="qbr_open_dialog"):
        show_qbr_dialog()

# Persistent download — visible even after the dialog closes
with col_dl:
    if "_qbr_pptx" in st.session_state:
        st.download_button(
            "📥 Download QBR (.pptx)",
            data=st.session_state["_qbr_pptx"],
            file_name=st.session_state.get("_qbr_filename", "QBR_report.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="qbr_dl_persistent",
        )

print("Done. QBR Generator page loaded.")
